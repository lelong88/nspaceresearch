from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

W = Inches(13.333)
H = Inches(7.5)
GREEN = RGBColor(0, 167, 88)
DARK = RGBColor(26, 26, 46)
BLUE = RGBColor(0, 115, 207)
RED = RGBColor(230, 57, 70)
ORANGE = RGBColor(244, 162, 97)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(102, 102, 102)
LGRAY = RGBColor(244, 245, 247)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, w, h, fill=None, text="", font_size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return shape

def add_text(slide, left, top, w, h, text, size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, w, h, lines, size=11, color=DARK, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bold
        p.space_after = Pt(4)
    return txBox

def add_kpi(slide, left, top, val, label, accent=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(1.1))
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = LGRAY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(32)
    p.font.color.rgb = DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False

def add_table(slide, left, top, w, rows_data, col_widths=None):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(0.35 * rows))
    table = table_shape.table
    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LGRAY
    return table_shape

def slide_header(slide, tag, title, tag_color=GREEN):
    add_box(slide, Inches(0.5), Inches(0.3), Inches(1.8), Inches(0.3), fill=RGBColor(230,247,239) if tag_color==GREEN else RGBColor(224,240,255) if tag_color==BLUE else RGBColor(253,232,234) if tag_color==RED else RGBColor(255,243,224), text=tag, font_size=8, color=tag_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.5), title, size=22, bold=True, color=DARK)

def slide_footer(slide, num):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(2), Inches(0.3), "III Manulife", size=8, color=GRAY)
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3), str(num), size=8, color=GRAY, align=PP_ALIGN.RIGHT)

# ============ SLIDE 1: TITLE ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
add_text(s, Inches(1), Inches(0.8), Inches(11), Inches(0.3), "CONFIDENTIAL — INTERNAL USE ONLY", size=10, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2), "Agency VAR Threshold\nAnalysis & Proposal", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(3.5), Inches(9), Inches(1.0), "Comprehensive review of Value-at-Risk exposure for the agency salesforce,\nhistorical trend analysis, and proposed threshold recalibration for 2026", size=14, color=RGBColor(170,170,187), align=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(5.5), Inches(9), Inches(0.3), "Agency Risk & Compensation Review  |  February 2026", size=10, color=GRAY, align=PP_ALIGN.CENTER)
slide_footer(s, 1)

# ============ SLIDE 2: EXECUTIVE SUMMARY ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "EXECUTIVE SUMMARY", "Key Findings at a Glance")
add_kpi(s, Inches(0.5), Inches(1.3), "$0.3M", "CURRENT VAR (2025 ACTUAL)", GREEN)
add_kpi(s, Inches(3.5), Inches(1.3), "$3.0M", "APPROVED TRIGGER (LEGACY)", BLUE)
add_kpi(s, Inches(6.5), Inches(1.3), "$1.3M", "FORECAST VAR 2026", ORANGE)
add_kpi(s, Inches(9.5), Inches(1.3), "$2.0M", "PROPOSED TRIGGER POINT", RED)

add_text(s, Inches(0.5), Inches(2.7), Inches(6), Inches(0.35), "Why VAR Was Low in 2025", size=14, bold=True, color=BLUE)
add_multiline(s, Inches(0.5), Inches(3.1), Inches(6), Inches(2.5), [
    ("• Conduct reviews cleaned out high-risk teams (policy churning, mass cancellations), leaving a healthier salesforce", False, DARK),
    ("• Business downturn: recruitment declined ~60% in H1 2024 vs 2023; high persistency kept FYP strong", False, DARK),
    ("• After excluding timing-gap cases, clean VAR ≈ $0.24M", False, DARK),
], size=10)

add_text(s, Inches(6.8), Inches(2.7), Inches(6), Inches(0.35), "Why the Threshold Matters", size=14, bold=True, color=BLUE)
add_multiline(s, Inches(6.8), Inches(3.1), Inches(6), Inches(2.5), [
    ("• Breaching the threshold triggers a mandatory risk protocol — deep-dive analysis and escalation", False, DARK),
    ("• New compensation scheme (2026) increases UM+ income significantly, pushing VAR from $0.3M → ~$1.3M", False, DARK),
], size=10)
slide_footer(s, 2)

# ============ SLIDE 3: METHODOLOGY ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "METHODOLOGY", "How VAR Is Calculated", BLUE)
add_box(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.6), fill=DARK, text="VAR = Σ Total Payout (L12M)  −  Σ Total FYP (L12M)", font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, (title, body) in enumerate([
    ("Rolling 12 Months", "Sum of payout per month over last 12 months. Reflects full sales cycle and mitigates impact of late payments post free-look period (21 days)."),
    ("Direct Team Scope", "VAR Direct Team of UM = Σ payout L12M − Σ FYP of all FA direct reports & UM. Only counted at months where agent holds UM+ rank."),
    ("Agency Aggregation", "Only UM+ with VAR > 0 are counted (negative-VAR teams are zeroed). Agency VAR = Σ all positive direct-team VARs."),
]):
    x = Inches(0.5 + i * 4.2)
    add_box(s, x, Inches(2.1), Inches(3.9), Inches(1.5), fill=LGRAY, text="", font_size=10)
    add_text(s, x + Inches(0.15), Inches(2.2), Inches(3.6), Inches(0.3), title, size=11, bold=True, color=DARK)
    add_text(s, x + Inches(0.15), Inches(2.55), Inches(3.6), Inches(1.0), body, size=9, color=GRAY)

add_text(s, Inches(0.5), Inches(3.9), Inches(6), Inches(0.3), "Payout Includes", size=12, bold=True, color=BLUE)
add_multiline(s, Inches(0.5), Inches(4.25), Inches(6), Inches(2.5), [
    ("• First Year Commission (FYC) + rider bonus", False, DARK),
    ("• Agent compensation (monthly, quarterly, yearly bonuses)", False, DARK),
    ("• UM+ compensation (overrides, activation, development)", False, DARK),
    ("• Contests, campaigns, financing schemes, benefits", False, DARK),
    ("• GTF subsidy payments", False, DARK),
], size=9)

add_text(s, Inches(6.8), Inches(3.9), Inches(6), Inches(0.3), "Exclusions", size=12, bold=True, color=BLUE)
add_multiline(s, Inches(6.8), Inches(4.25), Inches(6), Inches(2.5), [
    ("• CSV first year (only ~0.4% of FYP due to high allocation charge)", False, DARK),
    ("• PSM+ income — monitored separately on monthly basis", False, DARK),
    ("• Tax support scheme, cash advance (loan), retirement funds", False, DARK),
], size=9)
slide_footer(s, 3)

print("Slides 1-3 done")

# ============ SLIDE 4: HISTORICAL VAR TREND ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "HISTORICAL ANALYSIS", "VAR Trend: 5-Year Rolling L12M ($ Million USD)")

var_data = [
    ["Month","2020","2021","2022","2023","2024","2025"],
    ["Jan","2.2","1.7","2.0","3.5","2.4","0.3"],
    ["Feb","2.1","1.8","2.1","3.5","2.0","0.3"],
    ["Mar","2.0","1.8","2.1","3.5","1.7","0.3"],
    ["Apr","1.9","1.8","2.1","3.4","1.5","0.2"],
    ["May","1.9","1.7","2.3","3.8","1.3","0.2"],
    ["Jun","1.9","1.7","2.2","3.6","0.8","0.2"],
    ["Jul","1.9","1.8","2.0","3.7","0.7","0.2"],
    ["Aug","1.8","1.8","2.4","3.6","0.7","0.2"],
    ["Sep","1.8","1.8","2.9","3.3","0.4","0.3"],
    ["Oct","1.8","1.8","3.0","3.3","0.4","0.3"],
    ["Nov","1.8","1.9","3.2","3.4","0.4","0.3"],
    ["Dec","1.7","2.0","3.8","2.6","0.4","0.3"],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), var_data)

add_box(s, Inches(0.5), Inches(5.6), Inches(5.9), Inches(1.2), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(5.65), Inches(5.5), Inches(0.25), "2020–2022: Stable Period (~$1.8–2.0M)", size=10, bold=True, color=GREEN)
add_text(s, Inches(0.65), Inches(5.95), Inches(5.5), Inches(0.8), "Mass recruitment era. Agency FYP ~$200M+/year. New agents contributed ~50% of business. VAR remained stable around $2.0M.", size=9, color=DARK)

add_box(s, Inches(6.7), Inches(5.6), Inches(6.1), Inches(1.2), fill=RGBColor(255,243,224), text="")
add_text(s, Inches(6.85), Inches(5.65), Inches(5.8), Inches(0.25), "⚠ 2022–2023: Spike to ~$3.5M", size=10, bold=True, color=ORANGE)
add_text(s, Inches(6.85), Inches(5.95), Inches(5.8), Inches(0.8), "Peak driven by: (1) deferred payments from 2022 peak production flowing into 2023, and (2) special retention schemes for PSM/UM/RH paying income without production. Peak $3.8M in Dec 2022.", size=9, color=DARK)

add_text(s, Inches(0.5), Inches(1.15), Inches(12), Inches(0.3), "Suggest: replace table with line chart in PowerPoint (6 series, Jan–Dec)", size=8, color=GRAY, align=PP_ALIGN.RIGHT)
slide_footer(s, 4)

# ============ SLIDE 5: VAR/FYP RATIO ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "RATIO ANALYSIS", "VAR as % of Agency FYP", BLUE)

ratio_data = [
    ["Year","Agency FYP (USD M)","Avg VAR (USD M)","Avg %VAR/FYP","Peak %VAR/FYP"],
    ["2020","222","1.9","1.0%","1.3% (Jan)"],
    ["2021","241","1.8","0.7%","0.8% (Dec)"],
    ["2022","204","2.5","1.2%","1.9% (Dec)"],
    ["2023","72","3.5","2.7%","4.0% (Nov)"],
    ["2024","45","1.1","1.8%","3.6% (Jan)"],
    ["2025","~55*","0.3","0.5%","0.6% (Jan/Dec)"],
]
add_table(s, Inches(0.5), Inches(1.3), Inches(12.3), ratio_data)

add_text(s, Inches(0.5), Inches(4.2), Inches(12), Inches(0.3), "*2025 FYP estimated. 2026 FYP target: 1,787 bn VND (~$71M USD).", size=8, color=GRAY)

add_box(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.0), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(4.75), Inches(12), Inches(0.25), "Key Insight", size=11, bold=True, color=GREEN)
add_text(s, Inches(0.65), Inches(5.05), Inches(11.8), Inches(0.6), "2023 was the outlier year — VAR/FYP peaked at 4.0% when business volume collapsed to $72M (from $204M in 2022) while compensation structures remained unchanged. The 2025 ratio of 0.5% is the lowest in 6 years, driven by high persistency and reduced manpower.", size=10, color=DARK)

add_text(s, Inches(0.5), Inches(1.15), Inches(12), Inches(0.3), "Suggest: add bar chart above table (6 bars, y-axis = %VAR/FYP)", size=8, color=GRAY, align=PP_ALIGN.RIGHT)
slide_footer(s, 5)

# ============ SLIDE 6: $0.3M DEEP DIVE ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "DEEP DIVE", "Deconstructing the $0.3M: Why 2025 VAR Was Historically Low")
add_kpi(s, Inches(0.5), Inches(1.3), "$0.3M", "TOTAL VAR (GROSS)", GREEN)
add_kpi(s, Inches(3.5), Inches(1.3), "$0.05M", "TIMING-GAP CASES", BLUE)
add_kpi(s, Inches(6.5), Inches(1.3), "$0.24M", "CLEAN VAR (NET)", GREEN)
add_kpi(s, Inches(9.5), Inches(1.3), "260", "UNITS WITH VAR > 0", ORANGE)

for i, (title, body) in enumerate([
    ("Deferred Payments ($0.05M)", "NAB and MLY agent payments from late-2024 production paid in early 2025. Contest payouts also lag past the 21-day free-look period. These create artificial VAR without fraud signal."),
    ("Reduced Manpower", "H1 2024 recruitment declined ~60% vs 2023. Fewer bonus qualifiers = fewer agents generating positive VAR. Total units with VAR > 0 dropped to 260."),
    ("High Persistency", "Unusually high policy renewal rates in 2025 kept FYP revenue strong. Fewer cancellations meant the denominator of the risk equation held up, suppressing VAR."),
]):
    x = Inches(0.5 + i * 4.2)
    add_box(s, x, Inches(2.7), Inches(3.9), Inches(1.6), fill=LGRAY, text="")
    add_text(s, x + Inches(0.15), Inches(2.8), Inches(3.6), Inches(0.3), title, size=10, bold=True, color=DARK)
    add_text(s, x + Inches(0.15), Inches(3.15), Inches(3.6), Inches(1.1), body, size=9, color=GRAY)

add_box(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.5), fill=RGBColor(255,243,224), text="")
add_text(s, Inches(0.65), Inches(4.65), Inches(12), Inches(0.25), "⚠ Critical Context: This Is a Good-Year Artifact", size=11, bold=True, color=ORANGE)
add_text(s, Inches(0.65), Inches(4.95), Inches(11.8), Inches(1.0), "The $0.3M figure reflects favorable market conditions, not a structural reduction in risk exposure. With 2026 targeting 30%+ FYP growth and a new compensation scheme that significantly increases UM+ income, VAR is expected to revert upward. The back-test of 2025 data under the new scheme already shows VAR tripling to ~$1.0M.", size=10, color=DARK)
slide_footer(s, 6)

print("Slides 4-6 done")

# ============ SLIDE 7: BONUS CONTRIBUTION ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "COMPOSITION ANALYSIS", "Bonus Contribution Breakdown — Units with VAR > 0", BLUE)

bonus_data = [
    ["Component","Amount (USD K)","% of Total"],
    ["Total Payout","1,214","100%"],
    ["1) FYC","305","25%"],
    ["2) Agent Income","392","32%"],
    ["   Agent Compensation","192","16%"],
    ["   Agent Contest","107+","~15%"],
    ["3) UM+ Income","517","43%"],
    ["   UM+ Compensation","372","31%"],
    ["   GTF Subsidy","~45","~4%"],
    ["   UM+ Contest","85","7%"],
]
add_table(s, Inches(0.5), Inches(1.3), Inches(7), bonus_data)

add_text(s, Inches(0.5), Inches(1.15), Inches(7), Inches(0.3), "Suggest: add donut chart to the right of this table", size=8, color=GRAY, align=PP_ALIGN.RIGHT)

add_box(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.2), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(5.35), Inches(12), Inches(0.25), "Key Finding: UM+ Income Dominates VAR Exposure (43%)", size=11, bold=True, color=GREEN)
add_text(s, Inches(0.65), Inches(5.65), Inches(11.8), Inches(0.7), "Manager-level compensation — particularly leader overrides ($225K, 19%) and leader activation ($125K, 10%) — is the largest contributor to positive-VAR units. Arbitrage risk is most likely at the UM+ layer where managers can influence their sub-tree's production patterns.", size=10, color=DARK)
slide_footer(s, 7)

# ============ SLIDE 8: RISK SEGMENTS ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "RISK SEGMENTS", "High-Risk Segments: Zero-FYP & Negative-FYP Units", RED)

add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3), "FYP = 0 & Payout > 0  (8 units, $36.6K VAR)", size=12, bold=True, color=BLUE)
z_data = [
    ["Component","Amount","% "],
    ["Total VAR / Payout","$36.6K","100%"],
    ["FYC","$0.1K","0%"],
    ["Agent Income","$19.3K","53%"],
    ["  Agent Contest","$15.0K","41%"],
    ["UM+ Income","$17.1K","47%"],
    ["  GTF Subsidy","$8.1K","22%"],
]
add_table(s, Inches(0.5), Inches(1.6), Inches(5.8), z_data)

add_text(s, Inches(6.8), Inches(1.2), Inches(6), Inches(0.3), "FYP < 0 & Payout > 0  (7 units, $10.9K VAR)", size=12, bold=True, color=BLUE)
n_data = [
    ["Component","Amount","% "],
    ["Total VAR","$10.9K","—"],
    ["Total Payout","$5.4K","100%"],
    ["FYC (negative)","($2.0K)","-36%"],
    ["Agent Income","$4.2K","77%"],
    ["  Agent Contest","$3.2K","59%"],
    ["UM+ Income","$3.2K","60%"],
]
add_table(s, Inches(6.8), Inches(1.6), Inches(5.8), n_data)

add_box(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.8), fill=RGBColor(255,243,224), text="")
add_text(s, Inches(0.65), Inches(4.85), Inches(12), Inches(0.25), "GTF Subsidy: Risk Component — With Claw-Back", size=11, bold=True, color=ORANGE)
add_multiline(s, Inches(0.65), Inches(5.15), Inches(11.8), Inches(1.3), [
    ("• GTF subsidy = 22% of VAR in zero-FYP units ($8.1K of $36.6K)", False, DARK),
    ("• Calculated on net production — does not require passing the free-look period", False, DARK),
    ("• Claw-back mechanism exists: if GTF terminates within 12–18 months, must repay full subsidy", False, DARK),
    ("• Residual risk: submit-cancel-resubmit loops can still generate bonus if GTF remains active", False, DARK),
], size=10)
slide_footer(s, 8)

# ============ SLIDE 9: BACK-TEST & FORECAST ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "FORWARD-LOOKING", "2026 Compensation Impact: Back-Test & Forecast", ORANGE)
add_kpi(s, Inches(0.5), Inches(1.3), "$0.3M", "2025 ACTUAL (OLD SCHEME)", GREEN)
add_kpi(s, Inches(3.5), Inches(1.3), "$1.05M", "2025 BACK-TEST (NEW SCHEME)", ORANGE)
add_kpi(s, Inches(6.5), Inches(1.3), "$1.29M", "2026 FORECAST (NEW SCHEME)", RED)
add_kpi(s, Inches(9.5), Inches(1.3), "+35%", "FYP GROWTH TARGET", BLUE)

add_text(s, Inches(0.5), Inches(2.7), Inches(6), Inches(0.3), "Back-Test 2025 (New Scheme on Actual Data)", size=11, bold=True, color=BLUE)
bt_data = [
    ["Segment","Non-MBA","Silver","Gold","Platinum","Total"],
    ["VAR (USD K)","860","138","31","21","1,049"],
    ["# Units VAR>0","374","11","3","3","391"],
    ["  GTF","102","2","1","—","105"],
    ["  Non-GTF","272","9","2","3","286"],
]
add_table(s, Inches(0.5), Inches(3.05), Inches(6), bt_data)

add_text(s, Inches(6.8), Inches(2.7), Inches(6), Inches(0.3), "Forecast 2026 (New Scheme + Growth)", size=11, bold=True, color=BLUE)
fc_data = [
    ["Segment","Non-MBA","Silver","Gold","Platinum","Total"],
    ["VAR (USD K)","694","284","142","172","1,292"],
    ["# Units VAR>0","260","18","15","13","306"],
    ["  GTF","119","—","4","1","124"],
    ["  Non-GTF","141","14","15","12","182"],
]
add_table(s, Inches(6.8), Inches(3.05), Inches(6), fc_data)

add_box(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.4), fill=RGBColor(255,243,224), text="")
add_text(s, Inches(0.65), Inches(5.25), Inches(12), Inches(0.25), "⚠ Divergent Trends by Segment", size=11, bold=True, color=ORANGE)
add_multiline(s, Inches(0.65), Inches(5.55), Inches(11.8), Inches(1.0), [
    ("• Non-MBA Pro VAR drops 19% ($860K → $694K): income only +3% while FYP grows 35%", False, DARK),
    ("• MBA Pro VAR surges 3.2× ($190K → $598K): new scheme boosts top-performer compensation", False, DARK),
    ("• GTF units with VAR>0 increase 18% (105 → 124) despite total units declining 22% — anomaly under investigation", False, RED),
], size=10)
slide_footer(s, 9)

print("Slides 7-9 done")

# ============ SLIDE 10: NOISE FACTORS ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "CALIBRATION CHALLENGES", "Why the Threshold Number Is Noisy", ORANGE)

for i, (title, body) in enumerate([
    ("Deferred Payouts", "Commissions paid on a lag (sometimes >1 year) mean this year's VAR partially reflects last year's sales. The payout schedule artificially lowers VAR in the current year but raises it the next."),
    ("One-Off Retention Packages", "A single $500K–$1M talent retention payment to a key agent can spike individual VAR without any fraud signal. These are known, approved expenses — but they inflate aggregate VAR."),
    ("Revenue Growth Paradox", "Strong revenue growth pushes VAR up because top-performers perform even better and new recruits come in large quantity. The 2026 growth target of +35% FYP will mechanically increase VAR."),
]):
    x = Inches(0.5 + i * 4.2)
    add_box(s, x, Inches(1.2), Inches(3.9), Inches(2.0), fill=RGBColor(255,243,224), text="")
    add_text(s, x + Inches(0.15), Inches(1.3), Inches(3.6), Inches(0.3), title, size=11, bold=True, color=DARK)
    add_text(s, x + Inches(0.15), Inches(1.65), Inches(3.6), Inches(1.4), body, size=9, color=GRAY)

add_text(s, Inches(0.5), Inches(3.5), Inches(12), Inches(0.3), "Formula Mismatch: Income vs. FYP Scope", size=14, bold=True, color=BLUE)
add_box(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(1.2), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(3.95), Inches(11.8), Inches(1.1), "A structural asymmetry exists in the VAR formula: Income includes earnings from the entire team hierarchy (overrides on indirect downlines), but FYP only captures direct-team production. This means a UM+ with a large indirect team can show positive VAR even when the broader team is profitable — because their income reflects a wider scope than the FYP offset.", size=10, color=DARK)
slide_footer(s, 10)

# ============ SLIDE 11: THRESHOLD PROPOSAL ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "PROPOSAL", "Proposed Threshold: $2.0M Trigger Point")

add_box(s, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.9), fill=DARK, text="$2.0M", font_size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(3), Inches(2.15), Inches(7.3), Inches(0.3), "Recommended Trigger Point  |  Reduced from legacy $3.0M", size=10, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(2.6), Inches(6.5), Inches(0.3), "Rationale for $2.0M", size=13, bold=True, color=BLUE)
add_multiline(s, Inches(0.5), Inches(2.95), Inches(6.3), Inches(3.5), [
    ("• Tighter than legacy: 33% reduction from $3.0M demonstrating proactive risk tightening", False, DARK),
    ("• Sufficient buffer: forecast $1.29M → 1.55× buffer absorbs normal fluctuations", False, DARK),
    ("• Compensation impact absorbed: new scheme triples VAR to ~$1.0M on same production, still well below $2.0M", False, DARK),
    ("• Cleaner salesforce: conduct reviews in 2024–2025 structurally reduced risk exposure", False, DARK),
    ("• 2023 peak was anomalous: driven by deferred payments + one-off retention schemes — no longer applies", False, DARK),
], size=10)

add_text(s, Inches(7.2), Inches(2.6), Inches(5.5), Inches(0.3), "Scenario Analysis", size=13, bold=True, color=BLUE)
sc_data = [
    ["Scenario","Est. VAR","Triggers?"],
    ["Base case (2026 forecast)","$1.3M","✗ No"],
    ["+50% stress on forecast","$1.9M","✗ No"],
    ["Return to 2023 levels","$3.5M","✓ Yes"],
    ["One-off retention ($0.5M)","$1.8M","✗ No"],
    ["Recruitment surge + abuse","$2.0M+","✓ Yes"],
]
add_table(s, Inches(7.2), Inches(2.95), Inches(5.5), sc_data)

add_box(s, Inches(7.2), Inches(5.4), Inches(5.5), Inches(1.0), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(7.35), Inches(5.45), Inches(5.2), Inches(0.25), "Threshold Positioning", size=10, bold=True, color=GREEN)
add_text(s, Inches(7.35), Inches(5.75), Inches(5.2), Inches(0.6), "The $2.0M trigger is a 33% reduction from the legacy $3.0M — demonstrating tighter controls while providing sufficient headroom above the $1.3M forecast.", size=9, color=DARK)
slide_footer(s, 11)

# ============ SLIDE 12: MONITORING FRAMEWORK ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "GOVERNANCE", "VAR Monitoring & Control Framework", BLUE)

mon_data = [
    ["Action Item","Owner","Frequency","Timeline"],
    ["Dashboard Channel VaR report development","ASBD","—","2025 Q1"],
    ["Monitor total payout incl. PSM+ for significant fluctuations","ASBD","Monthly","Ongoing"],
    ["Review VaR monitoring tool for enhancement","ASBD","Quarterly","Ongoing"],
    ["Deep-dive payout data: compensation components, agency segments","ASBD","Monthly","From Mar '25"],
    ["Red-flag investigation: multi-indicator screening (persistency, sale patterns, phone/bank sharing, etc.)","Line1B","Monthly","Already in place"],
    ["Update tracking in Local ARC; escalate to ARC/RCSC when breached","Line1B","Monthly","From Mar '25"],
    ["Embed VAR impact review in annual compensation review","ASBD","Annual","2025 Q4"],
]
add_table(s, Inches(0.5), Inches(1.3), Inches(12.3), mon_data)

add_box(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.2), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(5.25), Inches(12), Inches(0.25), "Escalation Protocol When Threshold Is Breached", size=11, bold=True, color=GREEN)
add_text(s, Inches(0.65), Inches(5.55), Inches(11.8), Inches(0.7), "When monthly VAR exceeds $2.0M, the local Compensation team and Line1B (front risk sales) must provide a deep-dive analysis. Line1B already manages monthly protocols for other risk controls — this is integrated into existing workflows, not a new standalone process.", size=10, color=DARK)
slide_footer(s, 12)

print("Slides 10-12 done")

# ============ SLIDE 13: OPEN ITEMS ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "OPEN ITEMS", "Outstanding Questions & Risk Considerations", RED)

add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3), "Unresolved Questions", size=13, bold=True, color=BLUE)

for i, (title, body, clr) in enumerate([
    ("1. GTF Unit Anomaly", "GTF units with VAR > 0 increase 18% (105 → 124) in 2026 forecast despite total units declining 22%. Root cause under investigation: possible tier-jump effects or team-size changes.", RED),
    ("2. GTF Subsidy Abuse Risk", "Submit-cancel-resubmit loops can generate recurring bonus without genuine sales. Should this risk be explicitly addressed in the proposal?", RED),
    ("3. MBA Pro Concentration", "MBA Pro VAR surges 3.2× under new scheme. Is the increased exposure acceptable given business growth objectives?", ORANGE),
]):
    y = Inches(1.6 + i * 1.5)
    add_box(s, Inches(0.5), y, Inches(6.2), Inches(1.3), fill=LGRAY, text="")
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(6.2), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = clr; bar.line.fill.background(); bar.shadow.inherit = False
    add_text(s, Inches(0.65), y + Inches(0.1), Inches(5.9), Inches(0.25), title, size=10, bold=True, color=DARK)
    add_text(s, Inches(0.65), y + Inches(0.4), Inches(5.9), Inches(0.8), body, size=9, color=GRAY)

add_text(s, Inches(7), Inches(1.2), Inches(6), Inches(0.3), "Risk Mitigation Already in Place", size=13, bold=True, color=BLUE)
add_multiline(s, Inches(7), Inches(1.6), Inches(5.8), Inches(2.5), [
    ("• Monthly bi-conduct screening: payout ratios cross-referenced with risk indicators", False, DARK),
    ("• Payment suspension authority for high-risk cases pending investigation", False, DARK),
    ("• Claw-back clauses in high-reward scheme agreements", False, DARK),
    ("• Campaign enrollment suspension for violating agents/UMs", False, DARK),
    ("• Contract termination for confirmed fraud cases", False, DARK),
], size=10)

add_text(s, Inches(7), Inches(4.3), Inches(6), Inches(0.3), "Recommended Next Steps", size=13, bold=True, color=BLUE)
add_multiline(s, Inches(7), Inches(4.7), Inches(5.8), Inches(2.0), [
    ("• Complete deep-dive on 119 non-MBA GTF units", False, DARK),
    ("• Quantify GTF subsidy contribution to total VAR exposure", False, DARK),
    ("• Validate 30%+ FYP growth assumption", False, DARK),
    ("• Assess whether GTF subsidy claw-back mechanism should be strengthened", False, DARK),
], size=10)
slide_footer(s, 13)

# ============ SLIDE 14: CONCLUSION ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
add_box(s, Inches(0.5), Inches(0.3), Inches(1.5), Inches(0.3), fill=RGBColor(50,50,70), text="SUMMARY", font_size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.5), "Conclusion & Recommendation", size=22, bold=True, color=WHITE)

panels = [
    ("Current State", "$0.3M", "2025 actual VAR — historically low due to high persistency, reduced manpower, and business downtrend. This is a good-year artifact, not structural improvement.", GREEN),
    ("Expected State", "$1.3M", "2026 forecast under new compensation scheme with 35% FYP growth. New scheme triples VAR on same production; growth adds further upward pressure.", ORANGE),
    ("Proposed Trigger", "$2.0M", "Reduced from legacy $3.0M — a 33% tightening. Provides 1.55× buffer above forecast. Reflects cleaner salesforce post-conduct reviews.", BLUE),
]
for i, (title, val, desc, accent) in enumerate(panels):
    x = Inches(0.5 + i * 4.2)
    add_box(s, x, Inches(1.5), Inches(3.9), Inches(2.5), fill=RGBColor(40,40,60), text="")
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(0.06), Inches(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit = False
    add_text(s, x + Inches(0.2), Inches(1.6), Inches(3.5), Inches(0.3), title, size=13, color=accent, bold=True)
    add_text(s, x + Inches(0.2), Inches(1.95), Inches(3.5), Inches(0.5), val, size=28, color=WHITE, bold=True)
    add_text(s, x + Inches(0.2), Inches(2.55), Inches(3.5), Inches(1.3), desc, size=9, color=RGBColor(170,170,187))

add_box(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(1.5), fill=RGBColor(0,40,25), text="")
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.3), Inches(0.06), Inches(1.5))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background(); bar.shadow.inherit = False
add_text(s, Inches(0.7), Inches(4.4), Inches(11.8), Inches(1.3), "The $2.0M trigger point represents a proactive 33% tightening from the legacy $3.0M threshold, demonstrating strengthened risk controls. It is calibrated against the 2026 forecast ($1.29M), stress-tested under the new compensation scheme, and supported by a structurally cleaner salesforce following 2024–2025 conduct reviews. The comprehensive monthly monitoring framework with multi-indicator screening remains in place regardless of threshold level.", size=11, color=RGBColor(221,221,238))

add_text(s, Inches(2), Inches(6.2), Inches(9), Inches(0.3), "All data as of January 2026", size=9, color=GRAY, align=PP_ALIGN.CENTER)
slide_footer(s, 14)

# ============ SAVE ============
out = "/home/ubuntu/nspaceresearch/na/var-report.pptx"
prs.save(out)
print(f"Saved to {out}")
