from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

W = Inches(13.333)
H = Inches(7.5)
GREEN = RGBColor(0, 167, 88)
DARK = RGBColor(26, 26, 46)
BLUE = RGBColor(0, 115, 207)
RED = RGBColor(230, 57, 70)
ORANGE = RGBColor(244, 162, 97)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(102, 102, 102)
LGRAY = RGBColor(244, 245, 247)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, w, h, fill=None, text="", font_size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return shape

def add_text(slide, left, top, w, h, text, size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, w, h, lines, size=11, color=DARK, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bold
        p.space_after = Pt(4)
    return txBox

def add_kpi(slide, left, top, val, label, accent=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(1.1))
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = LGRAY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(32)
    p.font.color.rgb = DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False

def add_table(slide, left, top, w, rows_data, col_widths=None):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(0.35 * rows))
    table = table_shape.table
    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LGRAY
    return table_shape

def slide_header(slide, tag, title, tag_color=GREEN):
    add_box(slide, Inches(0.5), Inches(0.3), Inches(1.8), Inches(0.3), fill=RGBColor(230,247,239) if tag_color==GREEN else RGBColor(224,240,255) if tag_color==BLUE else RGBColor(253,232,234) if tag_color==RED else RGBColor(255,243,224), text=tag, font_size=8, color=tag_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.5), title, size=22, bold=True, color=DARK)

def slide_footer(slide, num):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(2), Inches(0.3), "III Manulife", size=8, color=GRAY)
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3), str(num), size=8, color=GRAY, align=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 1: MAC Pilot Performance Overview (12 MACs)
# Per feedback: 1 comprehensive slide with map reference, performance
# overview, warnings (MAC 01 fail MOC). No individual MAC names needed.
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "MAC PILOT", "12 Pilots Performance Update — as of Mar'26")

# KPI banners
add_kpi(s, Inches(0.3), Inches(1.2), "12", "TOTAL MAC OPERATIONAL", GREEN)
add_kpi(s, Inches(3.3), Inches(1.2), "6,561", "TOTAL FYP (MIL VND)", BLUE)
add_kpi(s, Inches(6.3), Inches(1.2), "80%", "AVG TARGET ACHIEVED", GREEN)
add_kpi(s, Inches(9.3), Inches(1.2), "139", "NEW RECRUITS", BLUE)

# Performance summary callout — positive
add_box(s, Inches(0.3), Inches(2.5), Inches(12.7), Inches(0.7), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.45), Inches(2.52), Inches(12.3), Inches(0.65),
    "Overall performance: Positive — 80% aggregate FYP target, 5/12 above 90% target, 100% quality compliance (12/12 office rated green, DC check: 12/12 pass). MAC cost/FYP ~15% (excl. one-off setup) vs. Branch office CE cost ~16–20%.",
    size=9, color=DARK)

# Performance table — no individual MAC names per feedback, use MAC codes only
perf_data = [
    ["MAC", "Contracted", "FYP (mil)", "Target (mil)", "% Target", "#NR", "# 1mAA", "Allowance (mil)", "Office", "DC", "Comment"],
    ["MAC 01", "2025-05", "1,254", "1,800", "70%", "13", "19", "223", "Green", "Pass", "Below target"],
    ["MAC 02", "2025-08", "475", "1,110", "43%", "8", "21", "240", "Green", "Pass", "Underperforming"],
    ["MAC 03", "2025-09", "445", "810", "55%", "27", "28", "160", "Green", "Pass", "FYP lagging"],
    ["MAC 04", "2025-09", "1,087", "810", "134%", "22", "36", "180", "Green", "Pass", "Star performer"],
    ["MAC 05", "2025-09", "750", "810", "93%", "28", "36", "200", "Green", "Pass", "Strong"],
    ["MAC 06", "2025-10", "889", "1,334", "67%", "12", "41", "128", "Green", "Pass", "Low conversion"],
    ["MAC 07", "2025-11", "1,535", "1,364", "113%", "25", "48", "132", "Green", "Pass", "Strong — urban"],
    ["MAC 08", "2025-12", "127", "210", "60%", "4", "4", "80", "Green", "Pass", "Below target"],
    ["MAC 09", "2026-02", "88", "60", "147%", "1", "5", "0", "Green", "Pass", "Too early"],
    ["MAC 10", "2026-02", "", "", "95%", "", "4", "0", "Green", "Pass", "Too early"],
    ["MAC 11", "2026-02", "", "", "—", "", "", "0", "Green", "Pass", "Too early"],
    ["MAC 12", "2026-02", "", "", "—", "", "", "0", "Green", "Pass", "Too early"],
    ["Total", "", "6,561", "8,248", "80%", "139", "233", "1,343", "12/12", "12/12", ""],
]
add_table(s, Inches(0.3), Inches(3.3), Inches(12.7), perf_data)

# Alert callout — MAC 01 fail MOC
add_box(s, Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.6), fill=RGBColor(253,232,234), text="")
add_text(s, Inches(0.45), Inches(6.22), Inches(12.3), Inches(0.55),
    "⚠ Alert: MAC 01 failed SM MOC (M9–M12) → Guaranteed subsidy stopped per scheme rules → Next review at M12; if demoted below SM → terminate contract. MAC as % of Agency FYP: 0.6% (early stage).",
    size=9, bold=False, color=RED)

add_text(s, Inches(0.3), Inches(3.15), Inches(12.5), Inches(0.2),
    "Suggest: add Vietnam map on the left showing MAC 09–12 locations (new MACs onboarded Feb'26)",
    size=8, color=GRAY, align=PP_ALIGN.RIGHT)

slide_footer(s, 1)
print("Slide 1 (Pilot Performance) done")


# =====================================================================
# SLIDE 2: Pipeline & Expansion Funnel
# Per feedback: visualize 70 → 22 → 12 funnel, no individual MAC names.
# Include expansion target: +10 MACs by end H1'26.
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "MAC PILOT", "Pipeline & Expansion Funnel")

# Funnel KPI banners — the 4 major stages
add_kpi(s, Inches(0.3), Inches(1.2), "70", "PIPELINE", BLUE)
add_kpi(s, Inches(3.3), Inches(1.2), "22", "OFFER LETTER", ORANGE)
add_kpi(s, Inches(6.3), Inches(1.2), "12", "MAC ONBOARDED", GREEN)
add_kpi(s, Inches(9.3), Inches(1.2), "8", "OPERATIONAL", GREEN)

# Funnel detail — Pipeline stage
add_text(s, Inches(0.3), Inches(2.5), Inches(3.0), Inches(0.3), "1. Pipeline — 70", size=13, bold=True, color=BLUE)
pipeline_data = [
    ["Stage", "Count"],
    ["Nominee", "75"],
    ["Application", "70"],
]
add_table(s, Inches(0.3), Inches(2.9), Inches(3.0), pipeline_data)

# Funnel detail — Offer Letter stage
add_text(s, Inches(3.5), Inches(2.5), Inches(3.0), Inches(0.3), "2. Offer Letter — 22", size=13, bold=True, color=ORANGE)
offer_data = [
    ["Stage", "Count"],
    ["Committee Interview", "25"],
    ["Passed & Released Offer", "23"],
    ["Offer Letter", "22"],
]
add_table(s, Inches(3.5), Inches(2.9), Inches(3.0), offer_data)

# Funnel detail — Onboarding stage
add_text(s, Inches(6.7), Inches(2.5), Inches(3.0), Inches(0.3), "3. MAC Onboarding — 12", size=13, bold=True, color=GREEN)
onboard_data = [
    ["Stage", "Count"],
    ["Location Approval", "14"],
    ["MIT", "14"],
    ["Register LTD", "14"],
    ["Appointed & Contract", "12"],
]
add_table(s, Inches(6.7), Inches(2.9), Inches(3.0), onboard_data)

# Funnel detail — Operation stage
add_text(s, Inches(9.9), Inches(2.5), Inches(3.0), Inches(0.3), "4. Operation — 8", size=13, bold=True, color=GREEN)
ops_data = [
    ["Stage", "Count"],
    ["Construction", "8"],
    ["Completed", "8"],
    ["Opening", "7"],
]
add_table(s, Inches(9.9), Inches(2.9), Inches(3.0), ops_data)

# Status updates section
add_text(s, Inches(0.3), Inches(4.9), Inches(12), Inches(0.3), "Current Status Highlights", size=13, bold=True, color=BLUE)
add_multiline(s, Inches(0.3), Inches(5.25), Inches(6.0), Inches(1.5), [
    ("Pipeline & Offer Letter:", True, DARK),
    ("• 1 MACD completed interview, confirmed offer → moving to onboarding", False, DARK),
    ("• 1 MACD completed interview, awaiting confirmation", False, DARK),
    ("Onboarding:", True, DARK),
    ("• 4 MACD signed MAC contract in Feb'26 (MAC 09–12)", False, DARK),
    ("• 3 MACD completed locations, moving to final step", False, DARK),
    ("• 2 MACD looking for office location", False, DARK),
    ("• 1 MACD failed onboarding process", False, RED),
], size=9)

add_multiline(s, Inches(6.7), Inches(5.25), Inches(6.0), Inches(1.5), [
    ("Operation:", True, DARK),
    ("• 7 MACs completed construction & opened", False, DARK),
    ("• 1 MAC completed construction, grand opening 5/3/2026", False, DARK),
], size=9)

# Expansion target callout
add_box(s, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.4), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.45), Inches(6.52), Inches(12.3), Inches(0.35),
    "Target: +10 MACs onboarding & operational by end of H1'26  |  Current pipeline: 74 nominees → 20 offers → 12 operational (7 opened)",
    size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

slide_footer(s, 2)
print("Slide 2 (Pipeline & Expansion Funnel) done")


# =====================================================================
# SLIDE 3: Selection & Onboarding Process
# Per feedback: visualize full process step-by-step, include full list
# of selection criteria, show comprehensiveness of the process.
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "MAC PILOT", "MAC Recruitment & Onboarding Process", BLUE)

# Phase 1: Interview (similar GTF SM+)
add_text(s, Inches(0.3), Inches(1.15), Inches(6), Inches(0.3), "Phase 1: Interview (similar GTF SM+)", size=14, bold=True, color=BLUE)

interview_data = [
    ["Step", "Responsible", "Process"],
    ["1", "Recruiter / RH", "1st Layer Interview & Review Application"],
    ["2", "Agency Operation (AO)", "Review application, Due diligence"],
    ["3", "AO, MBA Training", "Official Interview with Manulife — Stop if Fail"],
    ["4", "CAO, Head ASBD,\nHead MBA, Head Line1B", "MAC Committee Interview — Stop if Fail"],
]
add_table(s, Inches(0.3), Inches(1.5), Inches(7.5), interview_data)

# Selection criteria box — full list per feedback
add_text(s, Inches(8.1), Inches(1.15), Inches(5), Inches(0.3), "Selection Criteria (Full List)", size=14, bold=True, color=BLUE)
add_box(s, Inches(8.1), Inches(1.5), Inches(4.9), Inches(2.1), fill=LGRAY, text="")
add_multiline(s, Inches(8.25), Inches(1.55), Inches(4.6), Inches(2.0), [
    ("• FYP L12M (last 12 months production)", False, DARK),
    ("• #CA L6M (case count last 6 months)", False, DARK),
    ("• $Income L12M (income last 12 months)", False, DARK),
    ("• P13 (persistency month 13)", False, DARK),
    ("• Working experience in Life Insurance", False, DARK),
    ("• Education & Age requirements", False, DARK),
    ("• CIC credit check", False, RED),
    ("• Reference check", False, RED),
], size=9)

# Phase 2: Onboarding (SM GTF MAC)
add_text(s, Inches(0.3), Inches(3.9), Inches(6), Inches(0.3), "Phase 2: Onboarding (SM GTF MAC)", size=14, bold=True, color=GREEN)

onboard_process_data = [
    ["Step", "Responsible", "Process"],
    ["5", "AO", "Issuing Offer Letter, Support MIT, Contract sign"],
    ["6", "CRE & AMkt", "Review → Approval/Reject of office location — Stop if rejected"],
    ["7", "SM GTF MAC", "Ensure construction complies with legal & Manulife requirements"],
    ["8", "CRE & AMkt", "Sign Project Acceptance (leasing contract, construction contract,\ninvoice, look & feel alignment with guideline)"],
]
add_table(s, Inches(0.3), Inches(4.25), Inches(12.7), onboard_process_data)

# Final step callout
add_box(s, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.4), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.45), Inches(6.02), Inches(12.3), Inches(0.35),
    "Final: Scan back documents to CPA for future payment → MAC operational",
    size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Process flow summary
add_box(s, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.4), fill=RGBColor(224,240,255), text="")
add_text(s, Inches(0.45), Inches(6.52), Inches(12.3), Inches(0.35),
    "End-to-end: Scan hồ sơ → Interview GTF → MAC Committee → Business Plan → Offer → LTD Registration → Location Approval → Sign Contract → Operational",
    size=9, bold=False, color=BLUE, align=PP_ALIGN.CENTER)

slide_footer(s, 3)
print("Slide 3 (Selection & Onboarding Process) done")

# ============ SAVE ============
import os
out_dir = os.path.dirname(os.path.abspath(__file__))
out = os.path.join(out_dir, "gen_pptx_mac_2nd_draft.pptx")
prs.save(out)
print(f"Saved to {out}")
