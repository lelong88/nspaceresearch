"""
MAC Pilot Review — 2 Slides Only
Slide 1: Updated "Manulife Agency Center | Performance update as of Jan'26" (replaces existing deck slide)
Slide 2: One additional slide — Pilot Review Summary with key insights & recommendations
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

W = Inches(13.333)
H = Inches(7.5)
GREEN = RGBColor(0, 167, 88)
DARK = RGBColor(26, 26, 46)
BLUE = RGBColor(0, 115, 207)
RED = RGBColor(230, 57, 70)
ORANGE = RGBColor(244, 162, 97)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(102, 102, 102)
LGRAY = RGBColor(244, 245, 247)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── helpers ──────────────────────────────────────────────
def add_bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def add_box(sl, l, t, w, h, fill=None, text="", fs=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.shadow.inherit = False
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(fs)
    p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return sh

def add_text(sl, l, t, w, h, text, size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tx = sl.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tx

def add_ml(sl, l, t, w, h, lines, size=11, color=DARK):
    tx = sl.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    for i, (txt, bold, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(size)
        p.font.color.rgb = clr if clr else color; p.font.bold = bold
        p.space_after = Pt(4)
    return tx

def add_kpi(sl, l, t, val, label, accent=GREEN):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(2.8), Inches(1.1))
    sh.shadow.inherit = False; sh.fill.solid(); sh.fill.fore_color.rgb = LGRAY
    sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = val; p.font.size = Pt(32)
    p.font.color.rgb = DARK; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.06), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    bar.shadow.inherit = False

def add_table(sl, l, t, w, data, hdr_color=DARK):
    rows, cols = len(data), len(data[0])
    ts = sl.shapes.add_table(rows, cols, l, t, w, Inches(0.32 * rows))
    tbl = ts.table
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            c = tbl.cell(i, j); c.text = str(val)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(8); p.alignment = PP_ALIGN.CENTER
                if i == 0: p.font.bold = True; p.font.color.rgb = WHITE
            if i == 0: c.fill.solid(); c.fill.fore_color.rgb = hdr_color
            else: c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 == 1 else LGRAY
    return ts

def hdr(sl, tag, title, tc=GREEN):
    fm = {GREEN: RGBColor(230,247,239), BLUE: RGBColor(224,240,255),
          RED: RGBColor(253,232,234), ORANGE: RGBColor(255,243,224)}
    add_box(sl, Inches(0.5), Inches(0.3), Inches(2.2), Inches(0.3),
            fill=fm.get(tc, LGRAY), text=tag, fs=8, color=tc, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(0.5), Inches(0.65), Inches(12), Inches(0.5), title, size=22, bold=True, color=DARK)

def ftr(sl, n):
    add_text(sl, Inches(0.5), Inches(7.0), Inches(2), Inches(0.3), "III Manulife", size=8, color=GRAY)
    add_text(sl, Inches(5), Inches(7.0), Inches(3.3), Inches(0.3), "Highly Confidential", size=8, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3), str(n), size=8, color=GRAY, align=PP_ALIGN.RIGHT)

sn = 0

# ════════════════════════════════════════════════════════════
# SLIDE 1: UPDATED "Manulife Agency Center | Performance update as of Jan'26"
# This replaces the existing slide in the deck with enriched data
# from the feedback input (monthly breakdown, CC, allowance, quality)
# ════════════════════════════════════════════════════════════
sn += 1
s = prs.slides.add_slide(prs.slide_layouts[6])
hdr(s, "PILOT REVIEW", "Manulife Agency Center | Performance Update as of Jan'26")

# KPI row — same layout as original but with added CC metric
add_kpi(s, Inches(0.3), Inches(1.15), "8", "TOTAL MACs", GREEN)
add_kpi(s, Inches(3.3), Inches(1.15), "6,195m", "FYP (MIL VND)", BLUE)
add_kpi(s, Inches(6.3), Inches(1.15), "139", "NEW RECRUITS", ORANGE)
add_kpi(s, Inches(9.3), Inches(1.15), "93%", "AVG TARGET ACHIEVED", GREEN)

# Main performance table — enriched with CC, Allowance, Quality from feedback data
perf = [
    ["MAC","Name","Province","Contract\nDate","FYP\n(mil)","Target\n(mil)","%Ach","NR","1mAA","CC","Allowance\nPaid (mil)","Office\nQuality","DC"],
    ["01","Ecopark","Hưng Yên","May'25","1,254","1,800","70%","13","19","45","224","Green","Pass"],
    ["02","Quảng Yên","Quảng Ninh","Aug'25","458","810","57%","8","20","25","100","Green","Pass"],
    ["03","Đạ Tẻh","Lâm Đồng","Sep'25","429","570","75%","27","27","27","160","Green","Pass"],
    ["04","Diễn Châu","Nghệ An","Sep'25","986","570","173%","22","33","53","140","Green","Pass"],
    ["05","Vị Thanh","Cần Thơ","Sep'25","698","570","122%","28","33","41","160","Green","Pass"],
    ["06","Yên Bái","Lào Cai","Oct'25","767","1,050","73%","12","37","44","84","Green","Pass"],
    ["07","Thủ Đức","Hồ Chí Minh","Nov'25","1,477","1,148","129%","25","45","70","44","Green","Pass"],
    ["08","Vĩnh Phúc","Phú Thọ","Dec'25","127","120","106%","4","4","6","0","Green","Pass"],
    ["","Total","","","6,195","6,638","93%","139","218","311","912","8/8 Green","8/8 Pass"],
]
add_table(s, Inches(0.2), Inches(2.4), Inches(12.9), perf)

# Monthly FYP trend — compact row from feedback data
add_text(s, Inches(0.3), Inches(5.65), Inches(5), Inches(0.25),
         "Monthly FYP Trend (mil VND) — All MACs Combined", size=10, bold=True, color=BLUE)
mtrend = [
    ["M01","M02","M03","M04","M05","M06","M07","M08","M09"],
    ["1,043","1,679","1,391","1,233","448","346","38","17","0.3"],
]
add_table(s, Inches(0.2), Inches(5.95), Inches(5.2), mtrend)

# Right side: key context boxes
add_box(s, Inches(5.7), Inches(5.55), Inches(3.5), Inches(1.2), fill=RGBColor(230,247,239))
add_text(s, Inches(5.85), Inches(5.6), Inches(3.2), Inches(0.2),
         "5 of 8 MACs at or above target", size=9, bold=True, color=GREEN)
add_ml(s, Inches(5.85), Inches(5.82), Inches(3.2), Inches(0.85), [
    ("• MAC 04 Diễn Châu: 173% — star performer", False, GREEN),
    ("• MAC 07 Thủ Đức: 129% — highest absolute FYP", False, GREEN),
    ("• 100% Green office quality & DC pass", False, DARK),
    ("• MAC as % of Agency FYP: 0.6% (early stage)", False, GRAY),
], size=8)

add_box(s, Inches(9.4), Inches(5.55), Inches(3.5), Inches(1.2), fill=RGBColor(255,243,224))
add_text(s, Inches(9.55), Inches(5.6), Inches(3.2), Inches(0.2),
         "Action needed: MAC 02 & 01", size=9, bold=True, color=ORANGE)
add_ml(s, Inches(9.55), Inches(5.82), Inches(3.2), Inches(0.85), [
    ("• MAC 02 (57%) — lowest target achievement", False, RED),
    ("• MAC 01 — SM MOC fail at M9-M12", False, RED),
    ("• FYP trend declining after initial burst (M04+)", False, ORANGE),
    ("• Pipeline: 11 more MACs in onboarding/construction", False, DARK),
], size=8)

ftr(s, sn)

# ════════════════════════════════════════════════════════════
# SLIDE 2: ADDITIONAL — Pilot Review Summary & Recommendations
# Consolidates: cost analysis, quality tracking, pipeline, and next steps
# ════════════════════════════════════════════════════════════
sn += 1
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
add_box(s, Inches(0.5), Inches(0.3), Inches(2.2), Inches(0.3),
        fill=RGBColor(50,50,70), text="PILOT REVIEW", fs=8, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(0.65), Inches(12), Inches(0.45),
         "MAC Pilot Summary — Cost, Quality & Recommendations", size=22, bold=True, color=WHITE)

# ── Row 1: 4 KPI cards on dark background ──
kpis = [
    ("~13%", "STEADY-STATE\nCOST/FYP", GREEN),
    ("100%", "OFFICE QUALITY\nGREEN", BLUE),
    ("7 of 8", "FYP ON TRACK\n(SM COMP)", ORANGE),
    ("11", "MACs IN\nPIPELINE", RGBColor(180,130,255)),
]
for i, (val, label, accent) in enumerate(kpis):
    x = Inches(0.4 + i * 3.2)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.2), Inches(2.9), Inches(0.95))
    sh.shadow.inherit = False; sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(40,40,60)
    sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = val; p.font.size = Pt(26)
    p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(8)
    p2.font.color.rgb = RGBColor(170,170,187); p2.alignment = PP_ALIGN.CENTER
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.2), Inches(0.05), Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit = False

# ── Row 2: Three content columns ──

# Column 1: Cost & Compensation
add_box(s, Inches(0.4), Inches(2.35), Inches(4.0), Inches(2.6), fill=RGBColor(0,40,25))
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(2.35), Inches(0.05), Inches(2.6))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background(); bar.shadow.inherit = False
add_text(s, Inches(0.6), Inches(2.4), Inches(3.6), Inches(0.25),
         "Cost & Compensation", size=11, bold=True, color=GREEN)
add_ml(s, Inches(0.6), Inches(2.7), Inches(3.6), Inches(2.1), [
    ("Setup cost: ~410m/MAC (one-off, 100% sponsored)", False, WHITE),
    ("Allowance paid to date: ~912m (8% FYP-based)", False, WHITE),
    ("Total MAC scheme income: 54.5m across 8 MACs", False, WHITE),
    ("Total SM+ compensation: 6.4m", False, WHITE),
    ("", False, None),
    ("Current total cost/FYP: ~68% (incl. setup)", False, RGBColor(170,170,187)),
    ("Excl. one-off setup: ~15% (allowance only)", False, GREEN),
    ("Target steady-state: ~13% vs. branch ~16-20%", False, GREEN),
], size=9)

# Column 2: Quality & MOC Tracking
add_box(s, Inches(4.6), Inches(2.35), Inches(4.0), Inches(2.6), fill=RGBColor(15,25,50))
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.6), Inches(2.35), Inches(0.05), Inches(2.6))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background(); bar.shadow.inherit = False
add_text(s, Inches(4.8), Inches(2.4), Inches(3.6), Inches(0.25),
         "Quality & MOC Status", size=11, bold=True, color=BLUE)
add_ml(s, Inches(4.8), Inches(2.7), Inches(3.6), Inches(2.1), [
    ("Office quality: 8/8 Green", False, GREEN),
    ("DC check: 8/8 Pass", False, GREEN),
    ("P13: Pass (MAC 06 & 07); n/a for rest (too early)", False, WHITE),
    ("FYP on track: 7/8 (all except MAC 01)", False, WHITE),
    ("", False, None),
    ("Alert: MAC 01 failed SM MOC (M9-M12)", True, RGBColor(230,100,100)),
    ("→ Guaranteed subsidy stopped per scheme rules", False, RGBColor(170,170,187)),
    ("→ Next review at M12; if demoted below SM → terminate", False, RGBColor(170,170,187)),
], size=9)

# Column 3: Pipeline & Expansion
add_box(s, Inches(8.8), Inches(2.35), Inches(4.2), Inches(2.6), fill=RGBColor(50,30,10))
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(2.35), Inches(0.05), Inches(2.6))
bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background(); bar.shadow.inherit = False
add_text(s, Inches(9.0), Inches(2.4), Inches(3.8), Inches(0.25),
         "Pipeline & Expansion", size=11, bold=True, color=ORANGE)
add_ml(s, Inches(9.0), Inches(2.7), Inches(3.8), Inches(2.1), [
    ("74 nominees → 20 offers → 8 operational (7 opened)", False, WHITE),
    ("2 awaiting confirmation (Thanh Sơn, Nam Định)", False, WHITE),
    ("2 location approved (Bảy Hiền, Sơn La)", False, WHITE),
    ("6 onboarding (Trà Vinh, Phố Nối, Quy Nhơn,", False, WHITE),
    ("   Cẩm Phả, Pleiku, Đan Phượng)", False, WHITE),
    ("1 under construction (Vĩnh Phúc — opening 05/03)", False, WHITE),
    ("", False, None),
    ("Target: ~12-15 MACs operational by end H1'26", False, ORANGE),
], size=9)

# ── Row 3: Bottom — Recommendations bar ──
add_box(s, Inches(0.4), Inches(5.15), Inches(12.5), Inches(1.65), fill=RGBColor(40,40,60))
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(5.15), Inches(0.05), Inches(1.65))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background(); bar.shadow.inherit = False
add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.25),
         "Pilot Assessment & Recommendations", size=12, bold=True, color=GREEN)
add_ml(s, Inches(0.6), Inches(5.5), Inches(5.5), Inches(1.2), [
    ("Overall: POSITIVE — 93% aggregate target, 5/8 above target,", False, GREEN),
    ("100% quality compliance, team building faster than normal GTF.", False, GREEN),
    ("", False, None),
    ("Key ask:", True, WHITE),
    ("1. Approve MAC for ongoing implementation", False, WHITE),
    ("2. Launch Big MAC tier (DRD+, 150m2, higher targets)", False, WHITE),
    ("3. Introduce Internal MAC for 16-17 branch closures in 2026", False, WHITE),
], size=9)
add_ml(s, Inches(6.5), Inches(5.5), Inches(6.2), Inches(1.2), [
    ("Action items:", True, ORANGE),
    ("• Performance plan for MAC 02 (57%) & MAC 01 (MOC fail)", False, RGBColor(230,100,100)),
    ("• Streamline location approval (current bottleneck: 2-3 months)", False, WHITE),
    ("• Set higher 1mAA threshold for Big MAC (12/mth from M6)", False, WHITE),
    ("", False, None),
    ("Investment to date: ~4.2 bil across 8 MACs → 6.2 bil FYP", False, RGBColor(170,170,187)),
    ("Steady-state cost: ~13% FYP vs. branch model ~16-20%", False, RGBColor(170,170,187)),
], size=9)

ftr(s, sn)

# ════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════
out = "/home/ubuntu/nspaceresearch/na/output/mac-proposal/mac-pilot-review-additional-slides.pptx"
prs.save(out)
print(f"Saved {sn} slides to {out}")
