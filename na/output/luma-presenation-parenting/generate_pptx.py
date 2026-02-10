from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)
ACCENT_LIGHT = RGBColor(0xA5, 0x9E, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ROSE_COLOR = RGBColor(0x4E, 0xC9, 0x7A)
THORN_COLOR = RGBColor(0xFF, 0x6B, 0x6B)
BUD_COLOR = RGBColor(0xFF, 0xD9, 0x3D)
SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf

def add_rounded_rect(slide, left, top, width, height, fill_color, text="", font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Calibri"
    return shape

def add_divider(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 1, 1.5, 11, 1.2, "Coaching Session Presentation", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_divider(slide, 4.5, 2.8, 4.3)
add_text_box(slide, 1, 3.1, 11, 0.8, "Tien Phan", font_size=28, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.0, 11, 0.6, "LUMA Practitioner Program  •  Follow-up Coaching", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.8, "Sharing how I applied Design Thinking to a real-life parenting challenge", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ── SLIDE 2: The Challenge ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "A Few Words About My Challenge", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)

add_text_box(slide, 0.8, 1.6, 11, 0.6, "Planning meaningful after-school activities for my two daughters", font_size=20, color=ACCENT_LIGHT)

items = [
    "👧  Two daughters — ages 4 and 2 — come home at 4 PM every weekday",
    "👩‍👧‍👧  A caregiver supervises while parents are at work",
    "🎓  Both attend a Steiner (Waldorf) school",
    "📋  Activities must be easy for the caregiver to follow and instruct",
    "🌈  Activities should be diverse, age-appropriate, and complement Steiner methods",
    "🧩  May draw from other educational approaches to fill gaps",
]
add_bullet_list(slide, 1.0, 2.5, 10.5, 4.5, items, font_size=18, color=SOFT_WHITE, spacing=Pt(12))

# ── SLIDE 3: Recipe / Methods Used ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "My Recipe — LUMA Methods Used", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)

# Method cards
methods = [
    ("Stakeholder Mapping", "Identified everyone involved:\nchildren, caregiver, parents,\nteachers — and their needs.", ACCENT),
    ("Interviewing", "Talked with the caregiver and\nteachers to understand daily\nroutines and constraints.", RGBColor(0x00, 0x96, 0x88)),
    ("Rose, Thorn, Bud", "Evaluated past activity plans\nto find what worked, what\ndidn't, and new opportunities.", RGBColor(0xE9, 0x1E, 0x63)),
    ("Creative Matrix", "Cross-referenced activity types\nwith developmental goals to\ngenerate diverse ideas.", RGBColor(0xFF, 0x98, 0x00)),
]

for i, (title, desc, color) in enumerate(methods):
    x = 0.8 + i * 3.1
    card = add_rounded_rect(slide, x, 1.8, 2.8, 1.0, color, title, font_size=16, font_color=WHITE)
    card.text_frame.paragraphs[0].font.bold = True
    add_text_box(slide, x + 0.15, 3.0, 2.6, 2.5, desc, font_size=14, color=LIGHT_GRAY)

add_text_box(slide, 0.8, 5.5, 11, 0.8, "Each method built on the previous one — moving from understanding → ideation → planning.", font_size=15, color=ACCENT_LIGHT)

# ── SLIDE 4: Outcomes ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Outcomes", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)

outcomes = [
    "✅  A structured weekly activity plan (Mon–Fri) tailored to each child's age",
    "✅  Simple, visual instruction cards the caregiver can follow independently",
    "✅  A balanced mix of creative play, motor skills, language, and sensory activities",
    "✅  Activities that complement Steiner philosophy while introducing Montessori & Reggio elements",
    "✅  A feedback loop — caregiver notes what works, parents iterate weekly",
]
add_bullet_list(slide, 1.0, 1.8, 10.5, 4.5, outcomes, font_size=18, color=SOFT_WHITE, spacing=Pt(14))

add_text_box(slide, 0.8, 5.5, 11, 0.8, "The biggest win: the caregiver now feels confident running activities without step-by-step parent guidance.", font_size=15, color=ACCENT_LIGHT)

# ── SLIDE 5: Rose, Thorn, Bud Reflection ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Reflection — Rose, Thorn & Bud", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)

# Rose
add_rounded_rect(slide, 0.8, 1.7, 3.6, 0.6, ROSE_COLOR, "🌹  Rose (Positive)", font_size=16)
rose_items = [
    "Design Thinking made the planning process systematic instead of ad-hoc",
    "Stakeholder Mapping ensured no perspective was missed",
    "The caregiver reported higher confidence and engagement from the children",
]
add_bullet_list(slide, 1.0, 2.5, 3.6, 2.5, rose_items, font_size=13, color=SOFT_WHITE, spacing=Pt(6))

# Thorn
add_rounded_rect(slide, 4.8, 1.7, 3.6, 0.6, THORN_COLOR, "🥀  Thorn (Challenge)", font_size=16)
thorn_items = [
    "Time-consuming to prepare detailed instruction cards initially",
    "Hard to predict which activities suit a 2-year-old vs. a 4-year-old perfectly",
    "Caregiver needed an adjustment period to follow new formats",
]
add_bullet_list(slide, 5.0, 2.5, 3.6, 2.5, thorn_items, font_size=13, color=SOFT_WHITE, spacing=Pt(6))

# Bud
add_rounded_rect(slide, 8.8, 1.7, 3.6, 0.6, BUD_COLOR, "🌱  Bud (Potential)", font_size=16, font_color=RGBColor(0x1B, 0x1B, 0x2F))
bud_items = [
    "Could create a reusable activity template library for each age group",
    "Explore involving the children in choosing activities (co-design)",
    "Extend the approach to weekend and holiday planning",
]
add_bullet_list(slide, 9.0, 2.5, 3.6, 2.5, bud_items, font_size=13, color=SOFT_WHITE, spacing=Pt(6))

# Evaluation criteria
add_text_box(slide, 0.8, 5.3, 11, 0.5, "Evaluation Criteria", font_size=16, bold=True, color=ACCENT_LIGHT)
criteria = ["Application: Methods were applied correctly to a real-world scenario",
            "Outcome: Tangible weekly plan and instruction cards produced",
            "Communication: Shared artifacts with caregiver and coaching peers"]
add_bullet_list(slide, 1.0, 5.8, 10.5, 1.5, criteria, font_size=13, color=LIGHT_GRAY, spacing=Pt(4))

# ── SLIDE 6: Learn More ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "I'd Like to Learn More About…", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)

learn_items = [
    "🔍  How to better sequence LUMA methods for personal (non-work) challenges",
    "🔍  Tips for facilitating co-design sessions with very young children",
    "🔍  Ways to create lightweight, visual artifacts that non-designers can use",
    "🔍  How peers have applied Design Thinking outside of traditional work contexts",
]
add_bullet_list(slide, 1.0, 1.8, 10.5, 3.5, learn_items, font_size=20, color=SOFT_WHITE, spacing=Pt(16))

add_text_box(slide, 1, 5.0, 11, 1, "Thank you! Looking forward to your feedback.", font_size=24, bold=True, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

# Save
prs.save("/home/ubuntu/nspaceresearch/na/output/luma-presenation-parenting/coaching-presentation.pptx")
print("PPTX saved.")
