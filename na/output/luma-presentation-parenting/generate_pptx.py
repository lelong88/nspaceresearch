from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)
ACCENT_LIGHT = RGBColor(0xA5, 0x9E, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ROSE_COLOR = RGBColor(0x4E, 0xC9, 0x7A)
THORN_COLOR = RGBColor(0xFF, 0x6B, 0x6B)
BUD_COLOR = RGBColor(0xFF, 0xD9, 0x3D)
SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)
TEAL = RGBColor(0x00, 0x96, 0x88)
PINK = RGBColor(0xE9, 0x1E, 0x63)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
CARD_BG = RGBColor(0x2A, 0x2A, 0x45)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8), bold_items=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        if bold_items and i in bold_items:
            p.font.bold = True
    return tf


def add_sub_bullet_list(slide, left, top, width, height, items, font_size=14, color=LIGHT_GRAY, spacing=Pt(4)):
    """Add a list with sub-bullets (indented items marked with leading spaces)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size - (level * 1))
        p.font.color.rgb = color if level > 0 else WHITE
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = level
    return tf

def add_rounded_rect(slide, left, top, width, height, fill_color, text="", font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Calibri"
    return shape

def add_divider(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_image_placeholder(slide, left, top, width, height, label="[Image Placeholder]"):
    """Add a dashed-border rectangle as an image placeholder."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    shape.line.color.rgb = ACCENT_LIGHT
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 4  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_LIGHT
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(int(height * 72 / 3))
    return shape

def add_footer(slide, slide_num, total=10):
    """Add LUMA-style footer with slide number."""
    add_text_box(slide, 0.5, 7.0, 4, 0.4, "LUMA Institute and its licensors  •  CONFIDENTIAL", font_size=9, color=LIGHT_GRAY)
    add_text_box(slide, 11.5, 7.0, 1.5, 0.4, f"LUMA INSTITUTE", font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, 12.3, 7.0, 0.8, 0.4, str(slide_num), font_size=9, color=ACCENT_LIGHT, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title & Session Introduction
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 1, 1.2, 11.3, 1.2, "Coaching Session Presentation", font_size=42, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_divider(slide, 4.5, 2.5, 4.3)
add_text_box(slide, 1, 2.8, 11.3, 0.8, "Tien Phan", font_size=30, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.7, 11.3, 0.6, "LUMA Practitioner Program  •  Follow-up Coaching", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_image_placeholder(slide, 5.4, 4.5, 2.5, 1.8, "[Photo of presenter\nor LUMA logo]")
add_footer(slide, 1)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — A Few Words About Your Challenge
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "A FEW WORDS ABOUT MY CHALLENGE", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)

add_text_box(slide, 0.8, 1.6, 7.5, 1.2,
    "As a parent of two young daughters (ages 4 and 2), I need to design meaningful, "
    "age-appropriate after-school activities that a caregiver can independently follow and instruct — "
    "while complementing their Steiner (Waldorf) education.",
    font_size=18, color=SOFT_WHITE)

items = [
    "👧  Two daughters — ages 4 and 2 — come home from school at 4 PM every weekday",
    "👩‍👧‍👧  A caregiver supervises while both parents are at work",
    "🎓  Both attend a Steiner (Waldorf) school",
    "📋  Activities must be simple enough for the caregiver to follow independently",
    "🌈  Must be diverse, age-appropriate, and complement Steiner philosophy",
    "🧩  May draw from Montessori & Reggio Emilia to fill developmental gaps",
]
add_bullet_list(slide, 1.0, 3.0, 7.2, 4.0, items, font_size=17, color=SOFT_WHITE, spacing=Pt(10))

add_image_placeholder(slide, 9.0, 1.6, 3.8, 5.0, "[Image: children doing\nafter-school activities\nor Waldorf-style play]")
add_footer(slide, 2)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Methods Used (Overview)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "METHODS USED", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)

# Looking category
add_rounded_rect(slide, 0.8, 1.7, 5.5, 0.55, ACCENT, "Looking: Methods for Observing Human Experience", font_size=15)
looking_items = [
    "•  Walk-a-Mile Immersion — experiencing the caregiver's daily routine firsthand",
    "•  Interviewing — conversations with caregiver and teachers",
]
add_bullet_list(slide, 1.0, 2.5, 5.3, 1.2, looking_items, font_size=15, color=SOFT_WHITE, spacing=Pt(6))

# Understanding category
add_rounded_rect(slide, 0.8, 3.5, 5.5, 0.55, TEAL, "Understanding: Methods for Analyzing Challenges & Opportunities", font_size=15)
understanding_items = [
    "•  Stakeholder Mapping — identifying all people involved and their needs",
    "•  Rose, Thorn, Bud — evaluating previous activity plans",
]
add_bullet_list(slide, 1.0, 4.3, 5.3, 1.2, understanding_items, font_size=15, color=SOFT_WHITE, spacing=Pt(6))

# Making category
add_rounded_rect(slide, 0.8, 5.3, 5.5, 0.55, ORANGE, "Making: Methods for Envisioning Future Possibilities", font_size=15)
making_items = [
    "•  Creative Matrix — cross-referencing activity types with developmental goals",
]
add_bullet_list(slide, 1.0, 6.1, 5.3, 0.8, making_items, font_size=15, color=SOFT_WHITE, spacing=Pt(6))

add_image_placeholder(slide, 7.2, 1.7, 5.5, 5.0, "[Image: LUMA method taxonomy\nor visual recipe diagram]")
add_footer(slide, 3)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Approach: Walk-a-Mile Immersion
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "APPROACH", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)
add_rounded_rect(slide, 0.8, 1.5, 3.5, 0.5, ACCENT, "Looking: Observing Human Experience", font_size=13)

add_text_box(slide, 0.8, 2.3, 7, 0.6, "1. Walk-a-Mile Immersion", font_size=24, bold=True, color=ACCENT_LIGHT)

walk_items = [
    "I spent a full afternoon shadowing the caregiver to experience the daily routine firsthand — from school pickup to bedtime.",
    "Goal: understand what the caregiver actually faces — time constraints, the children's energy levels, available materials, and space limitations.",
    "I observed how the children transition between activities, what captures their attention, and where they lose interest.",
    "This immersion revealed that the caregiver struggles most during the 4:30–5:30 PM window when both children need different types of engagement.",
]
add_bullet_list(slide, 1.0, 3.0, 6.8, 3.5, walk_items, font_size=16, color=SOFT_WHITE, spacing=Pt(10))

add_image_placeholder(slide, 8.5, 1.5, 4.3, 3.0, "[Photo: Walk-a-Mile\nobservation session\nor daily routine snapshot]")
add_image_placeholder(slide, 8.5, 4.8, 4.3, 1.8, "[Photo: Children's play area\nor activity materials]")
add_footer(slide, 4)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — Approach: Interviewing
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "APPROACH", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)
add_rounded_rect(slide, 0.8, 1.5, 3.5, 0.5, ACCENT, "Looking: Observing Human Experience", font_size=13)

add_text_box(slide, 0.8, 2.3, 7, 0.6, "2. Interviewing", font_size=24, bold=True, color=ACCENT_LIGHT)

interview_items = [
    "I interviewed the caregiver to understand her comfort level, what activities she finds easy or difficult to lead, and her observations about each child's preferences.",
    "I also spoke with the Steiner school teachers to learn about the children's classroom activities, developmental milestones, and recommended at-home extensions.",
    "Key insights collected:",
]
add_bullet_list(slide, 1.0, 3.0, 6.8, 2.0, interview_items, font_size=16, color=SOFT_WHITE, spacing=Pt(8))

insight_items = [
    ("The caregiver prefers visual, step-by-step instructions over written paragraphs", 0),
    ("The 4-year-old thrives with creative and imaginative play", 1),
    ("The 2-year-old needs more sensory and motor-skill activities", 1),
    ("Teachers recommended nature-based and hands-on activities", 1),
    ("These insights directly shaped the activity plan design", 0),
]
add_sub_bullet_list(slide, 1.2, 4.8, 6.5, 2.2, insight_items, font_size=15, color=SOFT_WHITE, spacing=Pt(5))

add_image_placeholder(slide, 8.5, 1.5, 4.3, 5.0, "[Photo: Interview notes\nor conversation with\ncaregiver/teachers]")
add_footer(slide, 5)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — Approach: Stakeholder Mapping
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "APPROACH", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)
add_rounded_rect(slide, 0.8, 1.5, 4.5, 0.5, TEAL, "Understanding: Analyzing Challenges & Opportunities", font_size=13)

add_text_box(slide, 0.8, 2.3, 12, 0.6, "3. Stakeholder Mapping", font_size=24, bold=True, color=ACCENT_LIGHT)

# Stakeholder cards
stakeholders = [
    ("Children\n(ages 4 & 2)", "Primary beneficiaries\nDifferent developmental needs\nDifferent attention spans", ACCENT),
    ("Caregiver", "Daily executor of activities\nNeeds simple instructions\nLimited educational background", TEAL),
    ("Parents\n(Tien & wife)", "Activity designers & planners\nWeekly iteration & feedback\nQuality oversight", PINK),
    ("School Teachers", "Educational philosophy input\nDevelopmental milestone guidance\nActivity recommendations", ORANGE),
]

for i, (title, desc, color) in enumerate(stakeholders):
    x = 0.8 + i * 3.1
    card = add_rounded_rect(slide, x, 3.1, 2.8, 1.0, color, title, font_size=14, font_color=WHITE)
    card.text_frame.paragraphs[0].font.bold = True
    add_text_box(slide, x + 0.15, 4.3, 2.6, 2.0, desc, font_size=13, color=LIGHT_GRAY)

add_text_box(slide, 0.8, 6.0, 11, 0.6, "Why Stakeholder Mapping mattered:", font_size=16, bold=True, color=ACCENT_LIGHT)
why_items = [
    "Clarified each person's role and constraints  •  Ensured no perspective was missed  •  Revealed the caregiver as the critical success factor",
]
add_bullet_list(slide, 1.0, 6.4, 11, 0.6, why_items, font_size=14, color=SOFT_WHITE, spacing=Pt(4))
add_footer(slide, 6)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Approach: Creative Matrix
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "APPROACH", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)
add_rounded_rect(slide, 0.8, 1.5, 4.5, 0.5, ORANGE, "Making: Envisioning Future Possibilities", font_size=13)

add_text_box(slide, 0.8, 2.3, 12, 0.6, "4. Creative Matrix", font_size=24, bold=True, color=ACCENT_LIGHT)

add_text_box(slide, 0.8, 3.0, 6.5, 0.8,
    "I cross-referenced activity categories with developmental goals for each age group. "
    "This systematic approach generated a wide range of ideas and ensured balanced coverage.",
    font_size=16, color=SOFT_WHITE)

# Matrix placeholder
add_image_placeholder(slide, 0.8, 4.0, 7.0, 2.8,
    "[Image: Creative Matrix grid\n\n"
    "Columns: Creative Play | Motor Skills | Language | Sensory | Nature\n"
    "Rows: Age 4 goals | Age 2 goals | Caregiver-friendly | Steiner-aligned]")

add_text_box(slide, 8.5, 2.3, 4.3, 0.5, "Key outcomes from the matrix:", font_size=16, bold=True, color=ACCENT_LIGHT)
matrix_outcomes = [
    "Generated 30+ activity ideas across 5 categories",
    "Filtered down to 15 core activities per age group",
    "Identified activities that work for both ages simultaneously",
    "Flagged activities needing special materials or preparation",
    "Created a balanced weekly rotation template",
]
add_bullet_list(slide, 8.7, 3.0, 4.0, 3.5, matrix_outcomes, font_size=15, color=SOFT_WHITE, spacing=Pt(10))
add_footer(slide, 7)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Outcomes (Detailed)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "OUTCOMES", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)

outcomes_data = [
    ("1. Structured weekly activity plan", [
        "Monday–Friday schedule tailored to each child's age",
        "Balanced mix: creative play, motor skills, language, sensory, nature",
        "Activities complement Steiner while introducing Montessori & Reggio elements",
    ]),
    ("2. Visual instruction cards for the caregiver", [
        "Simple, step-by-step cards with illustrations",
        "Caregiver can follow independently — no phone calls needed",
        "Color-coded by activity type for quick reference",
    ]),
    ("3. Weekly feedback loop established", [
        "Caregiver notes what works and what doesn't each day",
        "Parents review and iterate the plan every weekend",
        "Continuous improvement built into the process",
    ]),
    ("4. Caregiver confidence significantly improved", [
        "From uncertain and dependent → confident and autonomous",
        "Children show higher engagement and less screen time",
        "Parents feel peace of mind during work hours",
    ]),
]

for i, (header, bullets) in enumerate(outcomes_data):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.6 + row * 2.8
    add_text_box(slide, x, y, 5.8, 0.5, header, font_size=18, bold=True, color=ACCENT_LIGHT)
    add_bullet_list(slide, x + 0.2, y + 0.5, 5.5, 2.0, bullets, font_size=14, color=SOFT_WHITE, spacing=Pt(5))

add_footer(slide, 8)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Outcomes (Recap / Visual)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "OUTCOMES", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)

add_text_box(slide, 0.8, 1.6, 7, 1.5,
    "The process transformed our after-school routine from ad-hoc and stressful "
    "into a structured, joyful experience for everyone involved.\n\n"
    "The children now look forward to their afternoon activities. The caregiver feels "
    "empowered and confident. And as parents, we have peace of mind knowing our "
    "daughters are engaged in meaningful, developmentally appropriate play — even "
    "when we can't be there.",
    font_size=17, color=SOFT_WHITE)

add_image_placeholder(slide, 8.5, 1.5, 4.3, 2.5, "[Photo: Sample weekly\nactivity plan or\ninstruction cards]")

add_image_placeholder(slide, 0.8, 4.0, 3.8, 2.8, "[Photo: Children engaged\nin a Steiner-style\ncreative activity]")
add_image_placeholder(slide, 5.0, 4.0, 3.8, 2.8, "[Photo: Caregiver using\nvisual instruction cards\nwith the children]")
add_image_placeholder(slide, 9.2, 4.3, 3.6, 2.5, "[Photo: Feedback notes\nor weekly iteration\nprocess]")
add_footer(slide, 9)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Rose, Thorn, Bud Reflection + Evaluation
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "EVALUATE YOUR APPLICATION OF METHODS", font_size=28, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.2, 3)
add_text_box(slide, 0.8, 1.35, 11, 0.4, "Reflect on how it went by recording a Rose, a Thorn, and a Bud:", font_size=15, color=LIGHT_GRAY)

# Rose column
add_rounded_rect(slide, 0.8, 1.9, 3.6, 0.55, ROSE_COLOR, "🌹  Rose (Positive)", font_size=15)
rose_items = [
    "Design Thinking made the planning process systematic instead of ad-hoc",
    "Stakeholder Mapping ensured no perspective was missed",
    "Walk-a-Mile revealed real constraints the caregiver faces daily",
    "The caregiver reported higher confidence and children show more engagement",
]
add_bullet_list(slide, 1.0, 2.65, 3.4, 2.2, rose_items, font_size=12, color=SOFT_WHITE, spacing=Pt(5))

# Thorn column
add_rounded_rect(slide, 4.8, 1.9, 3.6, 0.55, THORN_COLOR, "🥀  Thorn (Challenge)", font_size=15)
thorn_items = [
    "Time-consuming to prepare detailed visual instruction cards initially",
    "Hard to predict which activities suit a 2-year-old vs. a 4-year-old perfectly",
    "Caregiver needed an adjustment period to follow the new structured format",
]
add_bullet_list(slide, 5.0, 2.65, 3.4, 2.2, thorn_items, font_size=12, color=SOFT_WHITE, spacing=Pt(5))

# Bud column
add_rounded_rect(slide, 8.8, 1.9, 3.6, 0.55, BUD_COLOR, "🌱  Bud (Potential)", font_size=15, font_color=RGBColor(0x1B, 0x1B, 0x2F))
bud_items = [
    "Build a reusable activity template library for each age group",
    "Explore co-designing activities with the children themselves",
    "Extend the approach to weekends and holidays",
    "Share the framework with other parents in the school community",
]
add_bullet_list(slide, 9.0, 2.65, 3.4, 2.2, bud_items, font_size=12, color=SOFT_WHITE, spacing=Pt(5))

# Evaluation criteria section
add_divider(slide, 0.8, 5.0, 11.7)
add_text_box(slide, 0.8, 5.2, 11, 0.4, "Evaluation Criteria", font_size=16, bold=True, color=ACCENT_LIGHT)

eval_data = [
    ("Application", "Walk-a-Mile and Interviewing created firsthand insights; Stakeholder Mapping clarified responsibilities; Creative Matrix generated structured ideas"),
    ("Communication", "Visual instruction cards shared with caregiver; Weekly feedback forms; Presentation artifacts for coaching peers"),
    ("Outcome", "Tangible weekly plan and instruction cards produced; Caregiver autonomy achieved; Children's engagement improved measurably"),
]
for i, (label, desc) in enumerate(eval_data):
    y = 5.65 + i * 0.45
    add_text_box(slide, 0.8, y, 1.8, 0.4, f"▸ {label}:", font_size=13, bold=True, color=ACCENT_LIGHT)
    add_text_box(slide, 2.6, y, 10, 0.4, desc, font_size=12, color=SOFT_WHITE)

add_footer(slide, 10)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — I'd Like to Learn More About...
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "I'D LIKE TO LEARN MORE ABOUT...", font_size=32, bold=True, color=WHITE)
add_divider(slide, 0.8, 1.3, 3)

add_text_box(slide, 0.8, 1.6, 11, 0.6, "Input an item or two that you'd like to learn more about during the coaching session.", font_size=16, color=LIGHT_GRAY)

learn_items = [
    "1.  How to better sequence LUMA methods for personal (non-work) challenges — especially when the \"users\" are very young children",
    "2.  Tips for facilitating co-design sessions with toddlers and preschoolers — making LUMA methods accessible for ages 2–4",
    "3.  How to create lightweight, visual artifacts that non-designers (like a caregiver) can use and iterate on independently",
    "4.  How peers have applied Design Thinking outside of traditional work contexts — personal life, family, community",
]
add_bullet_list(slide, 1.0, 2.4, 10.5, 3.5, learn_items, font_size=18, color=SOFT_WHITE, spacing=Pt(16))

add_text_box(slide, 1, 5.5, 11, 1, "Thank you! Looking forward to your feedback and questions.", font_size=26, bold=True, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

add_image_placeholder(slide, 9.5, 5.8, 3.3, 1.2, "[LUMA Institute logo]")
add_footer(slide, 11)

# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output_path = "/home/ubuntu/nspaceresearch/na/output/luma-presentation-parenting/coaching-presentation.pptx"
prs.save(output_path)
print(f"PPTX saved to {output_path}")
