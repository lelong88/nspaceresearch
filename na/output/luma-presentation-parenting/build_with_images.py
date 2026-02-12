"""
Generate images via Bedrock and rebuild the PPTX with real photos replacing placeholders.
"""
import sys, os, asyncio, io
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))

from gen_image_bedrock import generate_image_bedrock
from PIL import Image

IMG_DIR = os.path.join(os.path.dirname(__file__), "images_generated")
os.makedirs(IMG_DIR, exist_ok=True)

# ── Image definitions: (filename, prompt, aspect_ratio) ──
IMAGES = [
    (
        "slide1_presenter.png",
        "A professional flat-lay photo of a modern workspace desk with a laptop showing a presentation slide, a coffee cup, colorful sticky notes, and a small succulent plant. Clean, minimal, top-down view. Warm soft lighting.",
        "16:9",
    ),
    (
        "slide2_waldorf_play.png",
        "A warm, bright photograph of two young Asian girls (ages 2 and 4) playing together at a wooden table with Waldorf-style natural toys, watercolor paints, beeswax crayons, and wooden blocks. Soft natural light from a window. Cozy home environment.",
        "3:4",
    ),
    (
        "slide3_luma_methods.png",
        "A clean flat-lay photograph of a design thinking workshop table with colorful post-it notes arranged in categories, markers, a whiteboard with a flowchart diagram, and printed cards labeled Looking, Understanding, Making. Professional, bright, top-down view.",
        "16:9",
    ),
    (
        "slide4_walk_a_mile.png",
        "A warm candid photograph of an Asian woman caregiver sitting on the floor with two young girls in a bright living room, guiding them through a hands-on craft activity with paper, glue, and natural materials. Afternoon sunlight streaming through windows.",
        "16:9",
    ),
    (
        "slide4_play_area.png",
        "A bright photograph of a children's play corner in a home with wooden shelves holding Waldorf toys, watercolor supplies, nature collection baskets, and a small wooden table with chairs. Organized, warm, inviting. Natural light.",
        "16:9",
    ),
    (
        "slide5_interview.png",
        "A photograph of two women sitting at a kitchen table having a friendly conversation, one taking notes in a notebook. A cup of tea on the table. Warm, natural indoor lighting. Casual and comfortable atmosphere.",
        "3:4",
    ),
    (
        "slide7_creative_matrix.png",
        "A top-down photograph of a large paper grid on a table filled with colorful sticky notes in different colors organized in rows and columns. Markers and pens around the edges. Design thinking brainstorming session. Clean, bright, professional.",
        "16:9",
    ),
    (
        "slide9_activity_plan.png",
        "A bright photograph of a colorful weekly planner pinned to a wall, with illustrated activity cards arranged by day. Each card has simple drawings and icons. Cheerful, organized, educational feel. Soft natural light.",
        "16:9",
    ),
    (
        "slide9_creative_activity.png",
        "A warm photograph of a young Asian girl around age 4 painting with watercolors at a small wooden table, concentrating happily. Natural materials like leaves and flowers on the table. Soft afternoon light. Waldorf-inspired creative activity.",
        "4:3",
    ),
    (
        "slide9_caregiver_cards.png",
        "A photograph of an Asian woman holding a colorful illustrated instruction card while a young child follows along doing a craft activity at a table. Bright, warm home setting. The card has simple step-by-step drawings.",
        "4:3",
    ),
    (
        "slide9_feedback.png",
        "A close-up photograph of a notebook open on a table with handwritten notes, checkmarks, and smiley faces. A pen resting on the page. Soft warm lighting. Simple feedback journal for tracking daily activities.",
        "4:3",
    ),
    (
        "slide11_logo.png",
        "A minimal, clean graphic of an open book with a lightbulb above it, symbolizing learning and innovation. Simple flat design on a dark navy background. Professional, modern, corporate style.",
        "16:9",
    ),
]


def resize_image(img_bytes, max_width=1200, max_height=900):
    """Resize image to reasonable dimensions for PPTX embedding."""
    img = Image.open(io.BytesIO(img_bytes))
    img.thumbnail((max_width, max_height), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    return buf.getvalue()


async def generate_all_images():
    """Generate all images, skipping ones that already exist."""
    for filename, prompt, aspect in IMAGES:
        path = os.path.join(IMG_DIR, filename)
        if os.path.exists(path):
            print(f"  ✓ {filename} already exists, skipping")
            continue
        print(f"  ⏳ Generating {filename}...")
        try:
            raw = await generate_image_bedrock(prompt, aspect_ratio=aspect)
            resized = resize_image(raw)
            with open(path, "wb") as f:
                f.write(resized)
            print(f"  ✓ {filename} saved ({len(resized)//1024} KB)")
        except Exception as e:
            print(f"  ✗ {filename} FAILED: {e}")


# ═══════════════════════════════════════════════════════════════════
# PPTX generation with real images
# ═══════════════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)
ACCENT_LIGHT = RGBColor(0xA5, 0x9E, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ROSE_COLOR = RGBColor(0x4E, 0xC9, 0x7A)
THORN_COLOR = RGBColor(0xFF, 0x6B, 0x6B)
BUD_COLOR = RGBColor(0xFF, 0xD9, 0x3D)
SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)
TEAL = RGBColor(0x00, 0x96, 0x88)
PINK = RGBColor(0xE9, 0x1E, 0x63)
ORANGE = RGBColor(0xFF, 0x98, 0x00)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8), bold_items=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        if bold_items and i in bold_items:
            p.font.bold = True
    return tf

def add_sub_bullet_list(slide, left, top, width, height, items, font_size=14, color=LIGHT_GRAY, spacing=Pt(4)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size - level)
        p.font.color.rgb = color if level > 0 else WHITE
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = level
    return tf

def add_rounded_rect(slide, left, top, width, height, fill_color, text="", font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Calibri"
    return shape

def add_divider(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_footer(slide, slide_num):
    add_text_box(slide, 0.5, 7.0, 4, 0.4, "LUMA Institute and its licensors  •  CONFIDENTIAL", font_size=9, color=LIGHT_GRAY)
    add_text_box(slide, 11.5, 7.0, 1.5, 0.4, "LUMA INSTITUTE", font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, 12.3, 7.0, 0.8, 0.4, str(slide_num), font_size=9, color=ACCENT_LIGHT, alignment=PP_ALIGN.RIGHT)

def img(name):
    """Return path to generated image, or None if missing."""
    p = os.path.join(IMG_DIR, name)
    return p if os.path.exists(p) else None

def add_img(slide, left, top, width, height, filename, fallback_label="[Image]"):
    """Add real image or a placeholder rectangle if image is missing."""
    path = img(filename)
    if path:
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
        shape.line.color.rgb = ACCENT_LIGHT
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = fallback_label
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = ACCENT_LIGHT
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1 — Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 1, 1.2, 11.3, 1.2, "Coaching Session Presentation", font_size=42, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_divider(slide, 4.5, 2.5, 4.3)
    add_text_box(slide, 1, 2.8, 11.3, 0.8, "Tien Phan", font_size=30, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.7, 11.3, 0.6, "LUMA Practitioner Program  •  Follow-up Coaching", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_img(slide, 5.4, 4.5, 2.5, 1.8, "slide1_presenter.png")
    add_footer(slide, 1)

    # ── SLIDE 2 — Challenge ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "A FEW WORDS ABOUT MY CHALLENGE", font_size=32, bold=True, color=WHITE)
    add_divider(slide, 0.8, 1.3, 3)
    add_text_box(slide, 0.8, 1.6, 7.5, 1.2,
        "As a parent of two young daughters (ages 4 and 2), I need to design meaningful, "
        "age-appropriate after-school activities that a caregiver can independently follow and instruct — "
        "while complementing their Steiner (Waldorf) education.",
        font_size=18, color=SOFT_WHITE)
    items = [
        "👧  Two daughters — ages 4 and 2 — come home from school at 4 PM every weekday",
        "👩‍👧‍👧  A caregiver supervises while both parents are at work",
        "🎓  Both attend a Steiner (Waldorf) school",
        "📋  Activities must be simple enough for the caregiver to follow independently",
        "🌈  Must be diverse, age-appropriate, and complement Steiner philosophy",
        "🧩  May draw from Montessori & Reggio Emilia to fill developmental gaps",
    ]
    add_bullet_list(slide, 1.0, 3.0, 7.2, 4.0, items, font_size=17, color=SOFT_WHITE, spacing=Pt(10))
    add_img(slide, 9.0, 1.6, 3.8, 5.0, "slide2_waldorf_play.png")
    add_footer(slide, 2)

    # ── SLIDE 3 — Methods Used ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "METHODS USED", font_size=32, bold=True, color=WHITE)
    add_divider(slide, 0.8, 1.3, 3)
    add_rounded_rect(slide, 0.8, 1.7, 5.5, 0.55, ACCENT, "Looking: Methods for Observing Human Experience", font_size=15)
    add_bullet_list(slide, 1.0, 2.5, 5.3, 1.2, [
        "•  Walk-a-Mile Immersion — experiencing the caregiver's daily routine firsthand",
        "•  Interviewing — conversations with caregiver and teachers",
    ], font_size=15, color=SOFT_WHITE, spacing=Pt(6))
    add_rounded_rect(slide, 0.8, 3.5, 5.5, 0.55, TEAL, "Understanding: Methods for Analyzing Challenges & Opportunities", font_size=15)
    add_bullet_list(slide, 1.0, 4.3, 5.3, 1.2, [
        "•  Stakeholder Mapping — identifying all people involved and their needs",
        "•  Rose, Thorn, Bud — evaluating previous activity plans",
    ], font_size=15, color=SOFT_WHITE, spacing=Pt(6))
    add_rounded_rect(slide, 0.8, 5.3, 5.5, 0.55, ORANGE, "Making: Methods for Envisioning Future Possibilities", font_size=15)
    add_bullet_list(slide, 1.0, 6.1, 5.3, 0.8, [
        "•  Creative Matrix — cross-referencing activity types with developmental goals",
    ], font_size=15, color=SOFT_WHITE, spacing=Pt(6))
    add_img(slide, 7.2, 1.7, 5.5, 5.0, "slide3_luma_methods.png")
    add_footer(slide, 3)

    # ── SLIDE 4 — Walk-a-Mile ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "APPROACH", font_size=32, bold=True, color=WHITE)
    add_divider(slide, 0.8, 1.3, 3)
    add_rounded_rect(slide, 0.8, 1.5, 3.5, 0.5, ACCENT, "Looking: Observing Human Experience", font_size=13)
    add_text_box(slide, 0.8, 2.3, 7, 0.6, "1. Walk-a-Mile Immersion", font_size=24, bold=True, color=ACCENT_LIGHT)
    add_bullet_list(slide, 1.0, 3.0, 6.8, 3.5, [
        "I spent a full afternoon shadowing the caregiver to experience the daily routine firsthand — from school pickup to bedtime.",
        "Goal: understand what the caregiver actually faces — time constraints, the children's energy levels, available materials, and space limitations.",
        "I observed how the children transition between activities, what captures their attention, and where they lose interest.",
        "This immersion revealed that the caregiver struggles most during the 4:30–5:30 PM window when both children need different types of engagement.",
    ], font_size=16, color=SOFT_WHITE, spacing=Pt(10))
    add_img(slide, 8.5, 1.5, 4.3, 3.0, "slide4_walk_a_mile.png")
    add_img(slide, 8.5, 4.8, 4.3, 1.8, "slide4_play_area.png")
    add_footer(slide, 4)

    # ── SLIDE 5 — Interviewing ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "APPROACH", font_size=32, bold=True, color=WHITE)
    add_divider(slide, 0.8, 1.3, 3)
    add_rounded_rect(slide, 0.8, 1.5, 3.5, 0.5, ACCENT, "Looking: Observing Human Experience", font_size=13)
    add_text_box(slide, 0.8, 2.3, 7, 0.6, "2. Interviewing", font_size=24, bold=True, color=ACCENT_LIGHT)
    add_bullet_list(slide, 1.0, 3.0, 6.8, 2.0, [
        "I interviewed the caregiver to understand her comfort level, what activities she finds easy or difficult to lead, and her observations about each child's preferences.",
        "I also spoke with the Steiner school teachers to learn about the children's classroom activities, developmental milestones, and recommended at-home extensions.",
        "Key insights collected:",
    ], font_size=16, color=SOFT_WHITE, spacing=Pt(8))
    add_sub_bullet_list(slide, 1.2, 4.8, 6.5, 2.2, [
        ("The caregiver prefers visual, step-by-step instructions over written paragraphs", 0),
        ("The 4-year-old thrives with creative and imaginative play", 1),
        ("The 2-year-old needs more sensory and motor-skill activities", 1),
        ("Teachers recommended nature-based and hands-on activities", 1),
        ("These insights directly shaped the activity plan design", 0),
    ], font_size=15, color=SOFT_WHITE, spacing=Pt(5))
    add_img(slide, 8.5, 1.5, 4.3, 5.0, "slide5_interview.png")
    add_footer(slide, 5)

    # ── SLIDE 6 — Stakeholder Mapping ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "APPROACH", font_size=32, bold=True, color=WHITE)
    add_divider(slide, 0.8, 1.3, 3)
    add_rounded_rect(slide, 0.8, 1.5, 4.5, 0.5, TEAL, "Understanding: Analyzing Challenges & Opportunities", font_size=13)
    add_text_box(slide, 0.8, 2.3, 12, 0.6, "3. Stakeholder Mapping", font_size=24, bold=True, color=ACCENT_LIGHT)
    stakeholders = [
        ("Children\n(ages 4 & 2)", "Primary beneficiaries\nDifferent developmental needs\nDifferent attention spans", ACCENT),
        ("Caregiver", "Daily executor of activities\nNeeds simple instructions\nLimited educational background", TEAL),
        ("Parents\n(Tien & wife)", "Activity designers & planners\nWeekly iteration & feedback\nQuality oversight", PINK),
        ("School Teachers", "Educational philosophy input\nDevelopmental milestone guidance\nActivity recommendations", ORANGE),
    ]
    for i, (title, desc, color) in enumerate(stakeholders):
        x = 0.8 + i * 3.1
        card = add_rounded_rect(slide, x, 3.1, 2.8, 1.0, color, title, font_size=14, font_color=WHITE)
        card.text_frame.paragraphs[0].font.bold = True
        add_text_box(slide, x + 0.15, 4.3, 2.6, 2.0, desc, font_size=13, color=LIGHT_GRAY)
    add_text_box(slide, 0.8, 6.0, 11, 0.6, "Why Stakeholder Mapping mattered:", font_size=16, bold=True, color=ACCENT_LIGHT)
    add_bullet_list(slide, 1.0, 6.4, 11, 0.6, [
        "Clarified each person's role and constraints  •  Ensured no perspective was missed  •  Revealed the caregiver as the critical success factor",
    ], font_size=14, color=SOFT_WHITE, spacing=Pt(4))
    add_footer(slide, 6)

    # ── SLIDE 7 — Creative Matrix ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "APPROACH", font_size=32, bold=True, color=WHITE)
    add_divider(slide, 0.8, 1.3, 3)
    add_rounded_rect(slide, 0.8, 1.5, 4.5, 0.5, ORANGE, "Making: Envisioning Future Possibilities", font_size=13)
    add_text_box(slide, 0.8, 2.3, 12, 0.6, "4. Creative Matrix", font_size=24, bold=True, color=ACCENT_LIGHT)
    add_text_box(slide, 0.8, 3.0, 6.5, 0.8,
        "I cross-referenced activity categories with developmental goals for each age group. "
        "This systematic approach generated a wide range of ideas and ensured balanced coverage.",
        font_size=16, color=SOFT_WHITE)
    add_img(slide, 0.8, 4.0, 7.0, 2.8, "slide7_creative_matrix.png")
    add_text_box(slide, 8.5, 2.3, 4.3, 0.5, "Key outcomes from the matrix:", font_size=16, bold=True, color=ACCENT_LIGHT)
    add_bullet_list(slide, 8.7, 3.0, 4.0, 3.5, [
        "Generated 30+ activity ideas across 5 categories",
        "Filtered down to 15 core activities per age group",
        "Identified activities that work for both ages simultaneously",
        "Flagged activities needing special materials or preparation",
        "Created a balanced weekly rotation template",
    ], font_size=15, color=SOFT_WHITE, spacing=Pt(10))
    add_footer(slide, 7)

    return prs


def build_pptx_remaining(prs):
    """Add slides 8-11 to the presentation."""

    # ── SLIDE 8 — Outcomes (Detailed) ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "OUTCOMES", font_size=32, bold=True, color=WHITE)
    add_divider(slide, 0.8, 1.3, 3)
    outcomes_data = [
        ("1. Structured weekly activity plan", [
            "Monday–Friday schedule tailored to each child's age",
            "Balanced mix: creative play, motor skills, language, sensory, nature",
            "Activities complement Steiner while introducing Montessori & Reggio elements",
        ]),
        ("2. Visual instruction cards for the caregiver", [
            "Simple, step-by-step cards with illustrations",
            "Caregiver can follow independently — no phone calls needed",
            "Color-coded by activity type for quick reference",
        ]),
        ("3. Weekly feedback loop established", [
            "Caregiver notes what works and what doesn't each day",
            "Parents review and iterate the plan every weekend",
            "Continuous improvement built into the process",
        ]),
        ("4. Caregiver confidence significantly improved", [
            "From uncertain and dependent → confident and autonomous",
            "Children show higher engagement and less screen time",
            "Parents feel peace of mind during work hours",
        ]),
    ]
    for i, (header, bullets) in enumerate(outcomes_data):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 6.2
        y = 1.6 + row * 2.8
        add_text_box(slide, x, y, 5.8, 0.5, header, font_size=18, bold=True, color=ACCENT_LIGHT)
        add_bullet_list(slide, x + 0.2, y + 0.5, 5.5, 2.0, bullets, font_size=14, color=SOFT_WHITE, spacing=Pt(5))
    add_footer(slide, 8)

    # ── SLIDE 9 — Outcomes (Recap / Visual) ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "OUTCOMES", font_size=32, bold=True, color=WHITE)
    add_divider(slide, 0.8, 1.3, 3)
    add_text_box(slide, 0.8, 1.6, 7, 1.5,
        "The process transformed our after-school routine from ad-hoc and stressful "
        "into a structured, joyful experience for everyone involved.\n\n"
        "The children now look forward to their afternoon activities. The caregiver feels "
        "empowered and confident. And as parents, we have peace of mind knowing our "
        "daughters are engaged in meaningful, developmentally appropriate play — even "
        "when we can't be there.",
        font_size=17, color=SOFT_WHITE)
    add_img(slide, 8.5, 1.5, 4.3, 2.5, "slide9_activity_plan.png")
    add_img(slide, 0.8, 4.0, 3.8, 2.8, "slide9_creative_activity.png")
    add_img(slide, 5.0, 4.0, 3.8, 2.8, "slide9_caregiver_cards.png")
    add_img(slide, 9.2, 4.3, 3.6, 2.5, "slide9_feedback.png")
    add_footer(slide, 9)

    # ── SLIDE 10 — Rose, Thorn, Bud ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "EVALUATE YOUR APPLICATION OF METHODS", font_size=28, bold=True, color=WHITE)
    add_divider(slide, 0.8, 1.2, 3)
    add_text_box(slide, 0.8, 1.35, 11, 0.4, "Reflect on how it went by recording a Rose, a Thorn, and a Bud:", font_size=15, color=LIGHT_GRAY)

    add_rounded_rect(slide, 0.8, 1.9, 3.6, 0.55, ROSE_COLOR, "🌹  Rose (Positive)", font_size=15)
    add_bullet_list(slide, 1.0, 2.65, 3.4, 2.2, [
        "Design Thinking made the planning process systematic instead of ad-hoc",
        "Stakeholder Mapping ensured no perspective was missed",
        "Walk-a-Mile revealed real constraints the caregiver faces daily",
        "The caregiver reported higher confidence and children show more engagement",
    ], font_size=12, color=SOFT_WHITE, spacing=Pt(5))

    add_rounded_rect(slide, 4.8, 1.9, 3.6, 0.55, THORN_COLOR, "🥀  Thorn (Challenge)", font_size=15)
    add_bullet_list(slide, 5.0, 2.65, 3.4, 2.2, [
        "Time-consuming to prepare detailed visual instruction cards initially",
        "Hard to predict which activities suit a 2-year-old vs. a 4-year-old perfectly",
        "Caregiver needed an adjustment period to follow the new structured format",
    ], font_size=12, color=SOFT_WHITE, spacing=Pt(5))

    add_rounded_rect(slide, 8.8, 1.9, 3.6, 0.55, BUD_COLOR, "🌱  Bud (Potential)", font_size=15, font_color=RGBColor(0x1B, 0x1B, 0x2F))
    add_bullet_list(slide, 9.0, 2.65, 3.4, 2.2, [
        "Build a reusable activity template library for each age group",
        "Explore co-designing activities with the children themselves",
        "Extend the approach to weekends and holidays",
        "Share the framework with other parents in the school community",
    ], font_size=12, color=SOFT_WHITE, spacing=Pt(5))

    add_divider(slide, 0.8, 5.0, 11.7)
    add_text_box(slide, 0.8, 5.2, 11, 0.4, "Evaluation Criteria", font_size=16, bold=True, color=ACCENT_LIGHT)
    eval_data = [
        ("Application", "Walk-a-Mile and Interviewing created firsthand insights; Stakeholder Mapping clarified responsibilities; Creative Matrix generated structured ideas"),
        ("Communication", "Visual instruction cards shared with caregiver; Weekly feedback forms; Presentation artifacts for coaching peers"),
        ("Outcome", "Tangible weekly plan and instruction cards produced; Caregiver autonomy achieved; Children's engagement improved measurably"),
    ]
    for i, (label, desc) in enumerate(eval_data):
        y = 5.65 + i * 0.45
        add_text_box(slide, 0.8, y, 1.8, 0.4, f"▸ {label}:", font_size=13, bold=True, color=ACCENT_LIGHT)
        add_text_box(slide, 2.6, y, 10, 0.4, desc, font_size=12, color=SOFT_WHITE)
    add_footer(slide, 10)

    # ── SLIDE 11 — Learn More ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, 0.8, 0.5, 11, 0.8, "I'D LIKE TO LEARN MORE ABOUT...", font_size=32, bold=True, color=WHITE)
    add_divider(slide, 0.8, 1.3, 3)
    add_text_box(slide, 0.8, 1.6, 11, 0.6, "Input an item or two that you'd like to learn more about during the coaching session.", font_size=16, color=LIGHT_GRAY)
    add_bullet_list(slide, 1.0, 2.4, 10.5, 3.5, [
        "1.  How to better sequence LUMA methods for personal (non-work) challenges — especially when the \"users\" are very young children",
        "2.  Tips for facilitating co-design sessions with toddlers and preschoolers — making LUMA methods accessible for ages 2–4",
        "3.  How to create lightweight, visual artifacts that non-designers (like a caregiver) can use and iterate on independently",
        "4.  How peers have applied Design Thinking outside of traditional work contexts — personal life, family, community",
    ], font_size=18, color=SOFT_WHITE, spacing=Pt(16))
    add_text_box(slide, 1, 5.5, 11, 1, "Thank you! Looking forward to your feedback and questions.", font_size=26, bold=True, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
    add_img(slide, 9.5, 5.8, 3.3, 1.2, "slide11_logo.png")
    add_footer(slide, 11)

    return prs


async def main():
    import time
    t0 = time.time()

    print("Step 1: Generating images...")
    await generate_all_images()
    print(f"  Images done in {time.time() - t0:.1f}s\n")

    print("Step 2: Building PPTX with images...")
    prs = build_pptx()
    prs = build_pptx_remaining(prs)

    out = "/home/ubuntu/nspaceresearch/na/output/luma-presentation-parenting/coaching-presentation.pptx"
    prs.save(out)
    print(f"  PPTX saved to {out}")
    print(f"\nTotal time: {time.time() - t0:.1f}s")


if __name__ == "__main__":
    asyncio.run(main())
