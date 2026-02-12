from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

W = Inches(13.333)
H = Inches(7.5)
GREEN = RGBColor(0, 167, 88)
DARK = RGBColor(26, 26, 46)
BLUE = RGBColor(0, 115, 207)
RED = RGBColor(230, 57, 70)
ORANGE = RGBColor(244, 162, 97)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(102, 102, 102)
LGRAY = RGBColor(244, 245, 247)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, w, h, fill=None, text="", font_size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return shape

def add_text(slide, left, top, w, h, text, size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, w, h, lines, size=11, color=DARK, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bold
        p.space_after = Pt(4)
    return txBox

def add_kpi(slide, left, top, val, label, accent=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(1.1))
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = LGRAY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(32)
    p.font.color.rgb = DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False

def add_table(slide, left, top, w, rows_data, col_widths=None, row_height=0.35):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(row_height * rows))
    table = table_shape.table
    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LGRAY
    return table_shape

def slide_header(slide, tag, title, tag_color=GREEN):
    bg_map = {
        GREEN: RGBColor(230,247,239),
        BLUE: RGBColor(224,240,255),
        RED: RGBColor(253,232,234),
        ORANGE: RGBColor(255,243,224),
    }
    add_box(slide, Inches(0.5), Inches(0.3), Inches(1.8), Inches(0.3),
            fill=bg_map.get(tag_color, LGRAY), text=tag, font_size=8,
            color=tag_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.5),
             title, size=22, bold=True, color=DARK)

def slide_footer(slide, num):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(2), Inches(0.3),
             "III Manulife", size=8, color=GRAY)
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3),
             str(num), size=8, color=GRAY, align=PP_ALIGN.RIGHT)

# ============ SLIDE 1: TITLE ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
add_text(s, Inches(1), Inches(0.8), Inches(11), Inches(0.3),
         "CONFIDENTIAL — INTERNAL USE ONLY", size=10, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
         "Agency VAR Threshold\nAnalysis & Proposal", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(3.5), Inches(9), Inches(1.0),
         "Comprehensive review of Value-at-Risk exposure for the agency salesforce,\nhistorical trend analysis, and proposed threshold recalibration for 2026",
         size=14, color=RGBColor(170,170,187), align=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(5.0), Inches(9), Inches(0.3),
         "3rd Draft — Updated with VAR History Performance Context & High-Exposure Deep Dive",
         size=11, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(5.5), Inches(9), Inches(0.3),
         "Agency Risk & Compensation Review  |  February 2026", size=10, color=GRAY, align=PP_ALIGN.CENTER)
slide_footer(s, 1)

# ============ SLIDE 2: EXECUTIVE SUMMARY ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "EXECUTIVE SUMMARY", "Key Findings at a Glance")
add_kpi(s, Inches(0.5), Inches(1.3), "$0.3M", "CURRENT VAR (2025 ACTUAL)", GREEN)
add_kpi(s, Inches(3.5), Inches(1.3), "$3.0M", "APPROVED TRIGGER (LEGACY)", BLUE)
add_kpi(s, Inches(6.5), Inches(1.3), "$0.85M", "FORECAST VAR 2026", ORANGE)
add_kpi(s, Inches(9.5), Inches(1.3), "$1.5M", "PROPOSED TRIGGER POINT", RED)

add_text(s, Inches(0.5), Inches(2.7), Inches(6), Inches(0.35),
         "Why VAR Was Low in 2025", size=14, bold=True, color=BLUE)
add_multiline(s, Inches(0.5), Inches(3.1), Inches(6), Inches(2.5), [
    ("• Conduct reviews cleaned out high-risk teams, leaving a healthier salesforce", False, DARK),
    ("• Business downturn: recruitment declined ~60%; high persistency kept FYP strong", False, DARK),
    ("• After excluding timing-gap cases, clean VAR ≈ $0.24M", False, DARK),
    ("• Aligned with CAFR arbitrage risk assessment back-test results", False, DARK),
], size=10)

add_text(s, Inches(6.8), Inches(2.7), Inches(6), Inches(0.35),
         "Why the Threshold Matters", size=14, bold=True, color=BLUE)
add_multiline(s, Inches(6.8), Inches(3.1), Inches(6), Inches(2.5), [
    ("• Breaching triggers mandatory risk protocol — deep-dive and escalation", False, DARK),
    ("• New compensation scheme pushes VAR from $0.3M → ~$0.85M", False, DARK),
    ("• Legacy $3.0M too high; $1.5M = proactive 50% tightening", False, DARK),
], size=10)

add_box(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.2), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(5.35), Inches(12), Inches(0.25),
         "3rd Draft Updates", size=11, bold=True, color=GREEN)
add_multiline(s, Inches(0.65), Inches(5.65), Inches(11.8), Inches(0.8), [
    ("• Added full agency performance context (FYP, commission rates, production mix) to VAR History", False, DARK),
    ("• Deep dive on 18 high-exposure units (VAR ≥$4K); verified defer payment as main driver for special segments", False, DARK),
    ("• Trigger point: $1.5M (50% reduction from legacy $3.0M)", False, DARK),
], size=9)
slide_footer(s, 2)

print("Slides 1-2 done")

# ============ SLIDE 3: CHANNEL VAR HISTORY WITH PERFORMANCE CONTEXT ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "HISTORICAL ANALYSIS", "Channel VAR History with Agency Performance Context")

var_data = [
    ["Month","2020","2021","2022","2023","2024","2025"],
    ["Jan","2.2","1.7","2.0","3.5","2.4","0.3"],
    ["Feb","2.1","1.8","2.1","3.5","2.0","0.3"],
    ["Mar","2.0","1.8","2.1","3.5","1.7","0.3"],
    ["Apr","1.9","1.8","2.1","3.4","1.5","0.2"],
    ["May","1.9","1.7","2.3","3.7","1.3","0.2"],
    ["Jun","1.9","1.7","2.2","3.5","0.8","0.2"],
    ["Jul","1.9","1.8","2.0","3.7","0.7","0.2"],
    ["Aug","1.8","1.8","2.4","3.6","0.7","0.2"],
    ["Sep","1.8","1.8","2.9","3.3","0.4","0.3"],
    ["Oct","1.8","1.8","3.0","3.3","0.4","0.3"],
    ["Nov","1.8","1.9","3.2","3.4","0.4","0.3"],
    ["Dec","1.7","2.0","3.8","2.6","0.4","0.3"],
]
add_table(s, Inches(0.3), Inches(1.15), Inches(7.5), var_data, row_height=0.28)

# Performance context table (NEW)
perf_data = [
    ["Metrics","2020","2021","2022","2023","2024","2025","2026F"],
    ["FYP (Mil USD)","222","240","204","72","50","53","72"],
    ["Avg Commission","—","—","—","32%","—","26%","~26%"],
    ["% FYP Core Agents","25%","27%","32%","33%","56%","65%","—"],
    ["% FYP New Agents","56%","51%","43%","39%","15%","13%","—"],
    ["#UM+ ending Dec","6,325","9,096","7,176","5,999","2,600","—","—"],
    ["Avg VaR (Mn)","1.9","1.8","2.5","3.5","1.1","0.3","0.85"],
    ["End %VaR/FYP","0.8%","0.8%","1.9%","3.6%","0.7%","0.6%","—"],
]
add_table(s, Inches(0.3), Inches(5.0), Inches(12.7), perf_data, row_height=0.25)

# Narrative boxes on the right
add_box(s, Inches(8.0), Inches(1.15), Inches(5.0), Inches(1.2), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(8.15), Inches(1.2), Inches(4.7), Inches(0.25),
         "2020–2022: Stable (~$1.8–2.0M)", size=9, bold=True, color=GREEN)
add_text(s, Inches(8.15), Inches(1.45), Inches(4.7), Inches(0.85),
         "Mass recruitment ~50K-60K/yr. FYP ~$200M+. New agents ~50% of FYP. Top agents ~30%. Core agents 25-32%. 6K-9K UM+.",
         size=8, color=DARK)

add_box(s, Inches(8.0), Inches(2.45), Inches(5.0), Inches(1.2), fill=RGBColor(255,243,224), text="")
add_text(s, Inches(8.15), Inches(2.5), Inches(4.7), Inches(0.25),
         "⚠ 2022–2023: Spike to $3.8M", size=9, bold=True, color=ORANGE)
add_text(s, Inches(8.15), Inches(2.75), Inches(4.7), Inches(0.85),
         "FYP collapsed 67% ($204M→$72M). Commission unchanged (FYC=40%, avg=32%). Deferred payments from 2022 peak. VAR/FYP peaked 4.0%.",
         size=8, color=DARK)

add_box(s, Inches(8.0), Inches(3.75), Inches(5.0), Inches(1.15), fill=RGBColor(224,240,255), text="")
add_text(s, Inches(8.15), Inches(3.8), Inches(4.7), Inches(0.25),
         "2024–2025: Decline to $0.3M", size=9, bold=True, color=BLUE)
add_text(s, Inches(8.15), Inches(4.05), Inches(4.7), Inches(0.8),
         "Commission reduced to 26%. Core agents surged to 65%. UM+ dropped to 2,600. Conduct reviews cleaned salesforce. New agents only 13%.",
         size=8, color=DARK)

slide_footer(s, 3)

# ============ SLIDE 4: HISTORICALLY LOW 2025 VAR ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "DEEP DIVE", "Historically Low 2025 VaR of $0.29 mil USD")
add_kpi(s, Inches(0.5), Inches(1.3), "$0.29M", "TOTAL VAR (GROSS)", GREEN)
add_kpi(s, Inches(3.5), Inches(1.3), "$0.05M", "TIMING-GAP CASES", BLUE)
add_kpi(s, Inches(6.5), Inches(1.3), "$0.24M", "CLEAN VAR (NET)", GREEN)
add_kpi(s, Inches(9.5), Inches(1.3), "260", "UNITS WITH VAR > 0", ORANGE)

add_multiline(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(2.5), [
    ("• Current VaR exposure ending 2025 (rolling 12 months) is $0.29 mil USD, equivalent to ~0.6% of Agency FYP — significantly below approved trigger of $3 mil USD", False, DARK),
    ("", False, DARK),
    ("• ~$0.05 mil USD VaR from timing-gap cases (FYP = 0 or negative, payout positive):", False, DARK),
    ("    — Deferred payment of previously earned retention bonuses", False, GRAY),
    ("    — Payout includes indirect team income vs FYP captures only direct team production", False, GRAY),
    ("", False, DARK),
    ("• After excluding timing-gap cases, clean VaR = ~$0.24 mil USD — aligned with CAFR arbitrage risk assessment", False, DARK),
    ("", False, DARK),
    ("• 260 UM+ contributing to total positive VaR of $0.29 mil USD", False, DARK),
], size=10)

add_box(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.3), fill=RGBColor(224,240,255), text="")
add_text(s, Inches(0.65), Inches(5.35), Inches(12), Inches(0.25),
         "Monthly Monitoring Controls", size=11, bold=True, color=BLUE)
add_text(s, Inches(0.65), Inches(5.65), Inches(11.8), Inches(0.9),
         "Line1B conducts monthly screening for compensation payout ratios with other risk indicators to identify suspicious cases; conducts in-depth analyses when triggered; suspension of compensation payment for high risk cases until investigation is completed.",
         size=10, color=DARK)
slide_footer(s, 4)

print("Slides 3-4 done")

# ============ SLIDE 5: UNITS PROFILE & PERFORMANCE ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "PROFILE ANALYSIS", "Units with VaR > 0: Profile & Performance", BLUE)

# Three donut-style KPIs
add_kpi(s, Inches(0.5), Inches(1.3), "63%", "ARE NON-GTF", BLUE)
add_kpi(s, Inches(3.5), Inches(1.3), "95%", "ARE NON-MBA PRO", ORANGE)
add_kpi(s, Inches(6.5), Inches(1.3), "72%", "LOW VAR (<$1K L12M)", GREEN)
add_kpi(s, Inches(9.5), Inches(1.3), "18", "HIGH EXPOSURE (≥$4K)", RED)

# Distribution table
dist_data = [
    ["% Payout/FYP","<$1K","<$2K","<$4K","≥$4K","Total"],
    ["FYP=0, Payout>0","67","8","3","—","78"],
    ["FYP<0, Payout>0","4","3","1","—","8"],
    ["FYP<0, Payout<0","4","2","—","—","6"],
    ["FYP>0, Payout>0","86","26","12","18","168"],
    ["Total","187","39","16","18","260"],
]
add_table(s, Inches(0.5), Inches(2.7), Inches(6.2), dist_data, row_height=0.3)

# VaR amount table
var_amt = [
    ["Segment","<$1K","<$2K","<$4K","≥$4K","Total $VaR"],
    ["FYP=0, Payout>0","18","10","8","1","37"],
    ["FYP<0, Payout>0","2","4","4","—","11"],
    ["FYP<0, Payout<0","2","3","—","—","4"],
    ["FYP>0, Payout>0","78","44","34","34","241"],
    ["Total","66","51","46","130","293"],
]
add_table(s, Inches(0.5), Inches(4.7), Inches(6.2), var_amt, row_height=0.3)

# Performance summary on right
perf_tbl = [
    ["Metric","Non MBA","Silver","Gold","Platinum","Total"],
    ["# UM VaR>0","166","4","4","4","168"],
    ["GTF","40","1","4","—","45"],
    ["Non GTF","115","3","—","4","123"],
    ["Avg FYP DT/mo (mil VND)","32","107","48","314","44"],
    ["Avg FYP personal/mo","13","24","4","60","14"],
    ["Avg 1mAA/mo","1.5","3.3","1.8","5.0","1.7"],
    ["Avg NR 2025 YTD","2.3","3.8","1.0","6.3","2.5"],
]
add_table(s, Inches(6.9), Inches(2.7), Inches(6.1), perf_tbl, row_height=0.28)

# Key insight box
add_box(s, Inches(6.9), Inches(5.3), Inches(6.1), Inches(1.3), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(7.05), Inches(5.35), Inches(5.8), Inches(0.25),
         "Key Insight: 72% Low-VaR Units", size=10, bold=True, color=GREEN)
add_text(s, Inches(7.05), Inches(5.65), Inches(5.8), Inches(0.85),
         "187 units (72%) have VaR <$1K, contributing only $66K (22% of total). Acceptable risk. The 18 units with VaR ≥$4K account for $130K (44% of total) — see deep dive next slide.",
         size=9, color=DARK)
slide_footer(s, 5)

# ============ SLIDE 6: DEEP DIVE 18 HIGH-EXPOSURE UNITS (NEW) ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "NEW: DEEP DIVE", "18 High-Exposure Units (VaR ≥ $4,000)", RED)

add_kpi(s, Inches(0.5), Inches(1.3), "$130K", "TOTAL VAR (44% OF ALL)", RED)
add_kpi(s, Inches(3.5), Inches(1.3), "18", "UNITS", ORANGE)
add_kpi(s, Inches(6.5), Inches(1.3), "100%", "FYP>0 & PAYOUT>0", BLUE)
add_kpi(s, Inches(9.5), Inches(1.3), "0", "CONDUCT FLAGS", GREEN)

for i, (title, body) in enumerate([
    ("Compensation Structure Driver",
     "High UM+ compensation (overrides + activation) relative to direct team FYP. Structural formula asymmetry: income from full hierarchy vs FYP from direct team only."),
    ("Large Indirect Teams",
     "Some units have large indirect teams generating override income not offset by FYP in the VAR formula. This creates positive VAR even without fraud signal."),
    ("Risk Assessment",
     "%VaR/FYP ranges from moderate to high. None flagged for conduct issues. VAR driven by compensation structure, not fraud. Individual monitoring recommended."),
]):
    x = Inches(0.5 + i * 4.2)
    add_box(s, x, Inches(2.7), Inches(3.9), Inches(1.8), fill=LGRAY, text="")
    add_text(s, x + Inches(0.15), Inches(2.8), Inches(3.6), Inches(0.3),
             title, size=10, bold=True, color=DARK)
    add_text(s, x + Inches(0.15), Inches(3.15), Inches(3.6), Inches(1.2),
             body, size=9, color=GRAY)

add_box(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.7), fill=RGBColor(255,243,224), text="")
add_text(s, Inches(0.65), Inches(4.85), Inches(12), Inches(0.25),
         "⚠ Recommended Actions for 18 High-Exposure Units", size=11, bold=True, color=ORANGE)
add_multiline(s, Inches(0.65), Inches(5.15), Inches(11.8), Inches(1.2), [
    ("• Individual case monitoring on monthly basis by Line1B", False, DARK),
    ("• Cross-reference with multi-indicator risk screening (persistency, sale patterns, phone/bank sharing)", False, DARK),
    ("• Review if compensation structure adjustments needed for outlier cases", False, DARK),
    ("• These 18 units contribute $130K (44%) of total $293K VaR — disproportionate but structurally explained", False, DARK),
], size=10)
slide_footer(s, 6)

print("Slides 5-6 done")

# ============ SLIDE 7: BONUS CONTRIBUTION ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "COMPOSITION ANALYSIS", "Bonus Contribution Breakdown — Units with VAR > 0", BLUE)

bonus_data = [
    ["Component","Amount (USD K)","% of Total"],
    ["Total Payout","1,214","100%"],
    ["1) FYC","305","25%"],
    ["2) Agent Income","392","32%"],
    ["   Agent Compensation","192","16%"],
    ["   Agent Contest","107","~15%"],
    ["3) UM+ Income","517","43%"],
    ["   UM+ Compensation","372","31%"],
    ["      Override","228","19%"],
    ["      Team Activation","125","10%"],
    ["      Other Comp Bonuses","19","2%"],
    ["   GTF Subsidy","45","4%"],
    ["   UM+ Contest","85","7%"],
]
add_table(s, Inches(0.5), Inches(1.3), Inches(7), bonus_data)

add_box(s, Inches(7.8), Inches(1.3), Inches(5.2), Inches(3.5), fill=LGRAY, text="")
add_text(s, Inches(7.95), Inches(1.35), Inches(5.0), Inches(0.3),
         "Income Composition", size=12, bold=True, color=DARK)
add_multiline(s, Inches(7.95), Inches(1.7), Inches(5.0), Inches(3.0), [
    ("FYC: 25%", True, BLUE),
    ("Agent Income: 32%", True, ORANGE),
    ("  — Agent Compensation: 16%", False, GRAY),
    ("  — Agent Contest: ~15%", False, GRAY),
    ("UM+ Income: 43%", True, RED),
    ("  — Override: 19% ($228K)", False, GRAY),
    ("  — Team Activation: 10% ($125K)", False, GRAY),
    ("  — GTF Subsidy: 4% ($45K)", False, GRAY),
    ("  — UM+ Contest: 7% ($85K)", False, GRAY),
], size=10)

add_box(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.0), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(5.65), Inches(12), Inches(0.25),
         "Key Finding: UM+ Income Dominates VAR Exposure (43%)", size=11, bold=True, color=GREEN)
add_text(s, Inches(0.65), Inches(5.95), Inches(11.8), Inches(0.6),
         "Leader overrides ($228K, 19%) and leader activation ($125K, 10%) are the largest contributors. Arbitrage risk is most likely at UM+ layer where managers can influence their unit production patterns.",
         size=10, color=DARK)
slide_footer(s, 7)

# ============ SLIDE 8: HIGH-RISK SEGMENTS ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "RISK SEGMENTS", "Zero-FYP & Negative-FYP Units — Defer Payment Verified", RED)

add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3),
         "FYP = 0 & Payout > 0  (78 units, $36.6K VaR)", size=12, bold=True, color=BLUE)
z_data = [
    ["Component","Amount","% ","Note"],
    ["Total VaR / Payout","$36.6K","100%",""],
    ["FYC","$0.1K","0%",""],
    ["Agent Income","$19.3K","53%",""],
    ["  Agent Contest","$15.0K","41%","Defer payment"],
    ["  Agent_NAB_FT","$3.2K","9%","Defer payment"],
    ["UM+ Income","$17.1K","47%",""],
    ["  GTF Subsidy (9 units)","$8.1K","22%","FYC direct & indirect"],
    ["  Leader Contest","$5.4K","15%","Defer payment"],
]
add_table(s, Inches(0.5), Inches(1.55), Inches(6), z_data, row_height=0.3)

add_text(s, Inches(6.8), Inches(1.2), Inches(6), Inches(0.3),
         "FYP < 0 & Payout > 0  (7 units, $10.9K VaR)", size=12, bold=True, color=BLUE)
n_data = [
    ["Component","Amount","% ","Note"],
    ["Total VaR","$10.9K","—",""],
    ["Total Payout","$5.4K","100%",""],
    ["FYC (negative)","($2.0K)","-36%",""],
    ["Agent Income","$4.2K","77%",""],
    ["  Agent Contest","$3.2K","59%","Defer payment"],
    ["UM+ Income","$3.2K","60%",""],
    ["  Leader Override","$0.8K","15%","FYC direct & indirect"],
    ["  Leader Contest","$1.8K","34%","Defer payment"],
]
add_table(s, Inches(6.8), Inches(1.55), Inches(6), n_data, row_height=0.3)

add_box(s, Inches(0.5), Inches(4.5), Inches(5.9), Inches(2.0), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(4.55), Inches(5.6), Inches(0.25),
         "✓ Verified: Defer Payment is Main Driver", size=10, bold=True, color=GREEN)
add_multiline(s, Inches(0.65), Inches(4.85), Inches(5.6), Inches(1.5), [
    ("• 78+7+6 = 91 special cases verified", False, DARK),
    ("• Majority driven by deferred payments from prior year production", False, DARK),
    ("• 66% of units have payout/FYP < 150%", False, DARK),
    ("• Most belong to MBA Pro group (not top UM)", False, DARK),
], size=9)

add_box(s, Inches(6.8), Inches(4.5), Inches(6.0), Inches(2.0), fill=RGBColor(255,243,224), text="")
add_text(s, Inches(6.95), Inches(4.55), Inches(5.7), Inches(0.25),
         "GTF Subsidy: Risk Component — With Claw-Back", size=10, bold=True, color=ORANGE)
add_multiline(s, Inches(6.95), Inches(4.85), Inches(5.7), Inches(1.5), [
    ("• GTF subsidy = 22% of VaR in zero-FYP units ($8.1K of $36.6K)", False, DARK),
    ("• Calculated on net production — no free-look requirement", False, DARK),
    ("• Claw-back: GTF terminates within 12-18 months → repay subsidy", False, DARK),
    ("• Residual risk: submit-cancel-resubmit loops possible if GTF active", False, DARK),
], size=9)
slide_footer(s, 8)

print("Slides 7-8 done")

# ============ SLIDE 9: PERFORMANCE NON-GTF ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "PERFORMANCE", "Non-GTF Units: FYP DT > 0, Payout > 0", BLUE)

ngtf_data = [
    ["%Payout/FYP DT",">0-<4K","≥4-<20K","≥20-<60K","≥60K","Total"],
    ["≥100% - <120%","16","17","6","2","33%"],
    ["≥120% - <150%","17","9","1","1","23%"],
    ["≥150% - <180%","12","—","—","—","10%"],
    ["≥180% - <200%","3","—","1","—","2%"],
    ["≥200% - <250%","3","—","—","—","2%"],
    ["≥250% - <300%","2","—","—","—","2%"],
    ["≥300%","34","—","—","—","28%"],
    ["Total","5%","21%","6%","2%","100%"],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(7.5), ngtf_data, row_height=0.28)

ngtf_seg = [
    ["FYP DT L12M","Non MBA","Silver","Gold","Platinum","Total"],
    [">0-<4K","87","—","—","—","71%"],
    ["≥4-<20K","24","1","—","1","21%"],
    ["≥20-<60K","4","2","—","1","6%"],
    ["≥60K","1","—","—","2","2%"],
    ["Total","94%","2%","0%","3%","100%"],
]
add_table(s, Inches(0.5), Inches(4.2), Inches(7.5), ngtf_seg, row_height=0.28)

add_box(s, Inches(8.3), Inches(1.2), Inches(4.7), Inches(2.5), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(8.45), Inches(1.25), Inches(4.4), Inches(0.25),
         "By %Payout/FYP", size=10, bold=True, color=GREEN)
add_multiline(s, Inches(8.45), Inches(1.55), Inches(4.4), Inches(2.0), [
    ("• 66% of units have ratio below 150%", False, DARK),
    ("• 4% with ratio >250% mostly low-production (FYP <4K USD)", False, DARK),
    ("• Payouts driven by deferred-payment components", False, DARK),
], size=9)

add_box(s, Inches(8.3), Inches(4.0), Inches(4.7), Inches(2.5), fill=RGBColor(224,240,255), text="")
add_text(s, Inches(8.45), Inches(4.05), Inches(4.4), Inches(0.25),
         "By Leader Segment", size=10, bold=True, color=BLUE)
add_multiline(s, Inches(8.45), Inches(4.35), Inches(4.4), Inches(2.0), [
    ("• 94% are Non MBA Pro", False, DARK),
    ("• 6% are Top Leaders (MBA Pro): FYP DT ≥4K, ≥2 1mAA/mo, avg 1 NR/mo", False, DARK),
    ("• 71% have FYP DT <4K USD — small production units", False, DARK),
], size=9)
slide_footer(s, 9)

# ============ SLIDE 10: PERFORMANCE GTF ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "PERFORMANCE", "GTF Units: FYP DT > 0, Payout > 0", BLUE)

gtf_data = [
    ["%Payout/FYP DT",">0-<4K","≥4-<20K","≥20-<60K","≥60K","Total"],
    ["≥100% - <120%","7","5","—","—","27%"],
    ["≥120% - <150%","12","1","—","—","29%"],
    ["≥150% - <180%","3","—","—","—","7%"],
    ["≥180% - <200%","3","—","—","—","7%"],
    ["≥200% - <250%","1","1","—","—","4%"],
    ["≥250% - <300%","2","—","—","—","4%"],
    ["≥300%","10","—","—","—","22%"],
    ["Total","84%","16%","0%","0%","100%"],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(7.5), gtf_data, row_height=0.28)

gtf_seg = [
    ["FYP DT L12M","Non MBA","Silver","Gold","Platinum","Total"],
    [">0-<4K","35","—","3","—","84%"],
    ["≥4-<20K","5","1","1","—","16%"],
    ["≥20-<60K","—","—","—","—","0%"],
    ["≥60K","—","—","—","—","0%"],
    ["Total","89%","2%","9%","0%","100%"],
]
add_table(s, Inches(0.5), Inches(4.2), Inches(7.5), gtf_seg, row_height=0.28)

add_box(s, Inches(8.3), Inches(1.2), Inches(4.7), Inches(2.5), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(8.45), Inches(1.25), Inches(4.4), Inches(0.25),
         "By %Payout/FYP", size=10, bold=True, color=GREEN)
add_multiline(s, Inches(8.45), Inches(1.55), Inches(4.4), Inches(2.0), [
    ("• 56% of units have ratio below 150%", False, DARK),
    ("• 12 units (~26%) with ratio >250% — low-production groups", False, DARK),
    ("• Payouts driven by non-deferred-payment components", False, DARK),
], size=9)

add_box(s, Inches(8.3), Inches(4.0), Inches(4.7), Inches(2.5), fill=RGBColor(224,240,255), text="")
add_text(s, Inches(8.45), Inches(4.05), Inches(4.4), Inches(0.25),
         "By Leader Segment", size=10, bold=True, color=BLUE)
add_multiline(s, Inches(8.45), Inches(4.35), Inches(4.4), Inches(2.0), [
    ("• 89% are Non MBA Pro (~40 units)", False, DARK),
    ("• 11% are Top Leaders (MBA Pro, ~5 units)", False, DARK),
    ("• 84% have FYP DT <4K USD — very small production", False, DARK),
], size=9)
slide_footer(s, 10)

print("Slides 9-10 done")

# ============ SLIDE 11: BACK-TEST & FORECAST (UPDATED) ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "FORWARD-LOOKING", "2026 Compensation Impact: Back-Test & Forecast", ORANGE)
add_kpi(s, Inches(0.5), Inches(1.3), "$0.3M", "2025 ACTUAL (OLD SCHEME)", GREEN)
add_kpi(s, Inches(3.5), Inches(1.3), "$0.48M", "2025 BACK-TEST (NEW SCHEME)", ORANGE)
add_kpi(s, Inches(6.5), Inches(1.3), "$0.85M", "2026 FORECAST (NEW SCHEME)", RED)
add_kpi(s, Inches(9.5), Inches(1.3), "+35%", "FYP GROWTH TARGET", BLUE)

add_text(s, Inches(0.5), Inches(2.7), Inches(6), Inches(0.3),
         "Back-Test 2025 (New Scheme on Actual Data)", size=11, bold=True, color=BLUE)
bt_data = [
    ["Segment","Non-MBA","Silver","Gold","Platinum","Total"],
    ["VaR (K USD)","378","83","18","2","481"],
    ["# Units VAR>0","190","6","3","3","202"],
    ["  GTF","45","2","1","—","48"],
    ["  Non-GTF","145","4","2","3","154"],
]
add_table(s, Inches(0.5), Inches(3.05), Inches(6), bt_data)

add_text(s, Inches(6.8), Inches(2.7), Inches(6), Inches(0.3),
         "Forecast 2026 (New Scheme + Growth)", size=11, bold=True, color=BLUE)
fc_data = [
    ["Segment","Non-MBA","Silver","Gold","Platinum","Total"],
    ["VaR (K USD)","421","224","94","117","856"],
    ["# Units VAR>0","260","18","15","13","306"],
    ["  GTF","108","4","—","—","112"],
    ["  Non-GTF","96","9","7","7","119"],
]
add_table(s, Inches(6.8), Inches(3.05), Inches(6), fc_data)

add_box(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.6), fill=RGBColor(255,243,224), text="")
add_text(s, Inches(0.65), Inches(5.05), Inches(12), Inches(0.25),
         "⚠ Divergent Trends by Segment (Updated Analysis)", size=11, bold=True, color=ORANGE)
add_multiline(s, Inches(0.65), Inches(5.35), Inches(11.8), Inches(1.2), [
    ("• Non-MBA Pro VaR drops 19%: income only +3% while FYP grows 35% → VAR/FYP improves", False, DARK),
    ("• MBA Pro VaR surges 3.2×: new scheme boosts top-performer compensation significantly", False, DARK),
    ("• GTF units with VAR>0 increase 18% (105→124): production increase causes tier jumps → disproportionate income increase", False, DARK),
    ("• Root cause confirmed: GTF compensation scheme unchanged but higher production → higher tier rates", False, RED),
], size=10)
slide_footer(s, 11)

# ============ SLIDE 12: 2026 VAR FORECAST MODEL (NEW) ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "NEW: FORECAST MODEL", "2026 VAR Forecast — Why Not Back to 2023 Peak?", BLUE)

# Comparison table
comp_data = [
    ["Factor","2023 (Peak Year)","2026 (Forecast)","Impact on VAR"],
    ["FYP Level","$72M","$72M (target)","Neutral"],
    ["Commission Rate","32% avg (FYC=40%)","~26% avg","↓ Lower payout"],
    ["Production Mix","New agents ~39%","Core agents >65%","↓ More stable"],
    ["#UM+","5,999","~2,600-3,000","↓ Fewer units"],
    ["Compensation Scheme","Old (higher rates)","New (higher for MBA)","↑ MBA Pro impact"],
    ["Salesforce Quality","Pre-conduct review","Post-conduct review","↓ Cleaner teams"],
    ["Deferred Payments","High (from 2022 peak)","Normal levels","↓ Less timing noise"],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), comp_data, row_height=0.3)

add_box(s, Inches(0.5), Inches(3.8), Inches(5.9), Inches(2.8), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(3.85), Inches(5.6), Inches(0.25),
         "Why 2026 VAR ($0.85M) << 2023 VAR ($3.5M)", size=10, bold=True, color=GREEN)
add_multiline(s, Inches(0.65), Inches(4.15), Inches(5.6), Inches(2.3), [
    ("Despite same FYP level ($72M):", True, DARK),
    ("• Commission rate 6pp lower (26% vs 32%)", False, DARK),
    ("• No mass deferred payments from prior peak year", False, DARK),
    ("• Core agent contribution >65% (vs 33% in 2023)", False, DARK),
    ("• UM+ count ~2,600 (vs 5,999 in 2023)", False, DARK),
    ("• Post-conduct review: high-risk teams removed", False, DARK),
    ("• No special retention schemes inflating payouts", False, DARK),
], size=9)

add_box(s, Inches(6.7), Inches(3.8), Inches(6.1), Inches(2.8), fill=RGBColor(255,243,224), text="")
add_text(s, Inches(6.85), Inches(3.85), Inches(5.8), Inches(0.25),
         "⚠ Upside Risk Factors", size=10, bold=True, color=ORANGE)
add_multiline(s, Inches(6.85), Inches(4.15), Inches(5.8), Inches(2.3), [
    ("• New compensation scheme triples VAR on same production for MBA Pro", False, DARK),
    ("• New ILP product: higher case size (2.5× min premium) → higher productivities", False, DARK),
    ("• If FYP growth exceeds 35% target → more units cross VAR threshold", False, DARK),
    ("• Core agents with income/FYP >100% may have elevated APE track risk", False, DARK),
    ("• GTF tier-jump effect: production growth → disproportionate income increase", False, DARK),
], size=9)
slide_footer(s, 12)

print("Slides 11-12 done")

# ============ SLIDE 13: TRIGGER POINT PROPOSAL ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "PROPOSAL", "Proposed Threshold: $1.5M Trigger Point")

add_box(s, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.9), fill=DARK,
        text="$1.5M", font_size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(3), Inches(2.15), Inches(7.3), Inches(0.3),
         "Recommended Trigger Point  |  Reduced from legacy $3.0M (50% tightening)",
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(2.6), Inches(6.5), Inches(0.3),
         "Rationale for $1.5M", size=13, bold=True, color=BLUE)
add_multiline(s, Inches(0.5), Inches(2.95), Inches(6.3), Inches(3.5), [
    ("• Proactive 50% reduction from $3.0M — strongest tightening signal", False, DARK),
    ("• ~1.76× buffer above $0.85M forecast absorbs normal fluctuations", False, DARK),
    ("• Calibrated for new compensation scheme impact on same production", False, DARK),
    ("• Cleaner salesforce post 2024-2025 conduct reviews", False, DARK),
    ("• 2023 peak ($3.5M) was anomalous — no longer applicable baseline", False, DARK),
], size=10)

add_text(s, Inches(7.2), Inches(2.6), Inches(5.5), Inches(0.3),
         "Scenario Analysis", size=13, bold=True, color=BLUE)
sc_data = [
    ["Scenario","Est. VAR","Triggers?"],
    ["Base case (2026 forecast)","$0.85M","✗ No"],
    ["+50% stress on forecast","$1.28M","✗ No"],
    ["One-off retention ($0.5M)","$1.35M","✗ No"],
    ["Recruitment surge + abuse","$1.5M+","✓ Yes"],
    ["Return to 2023 levels","$3.5M","✓ Yes"],
]
add_table(s, Inches(7.2), Inches(2.95), Inches(5.5), sc_data)

add_box(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.1), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(5.55), Inches(12), Inches(0.25),
         "Buffer Absorbs Known Noise Factors", size=10, bold=True, color=GREEN)
add_multiline(s, Inches(0.65), Inches(5.85), Inches(11.8), Inches(0.7), [
    ("• Structural formula asymmetry (income from full hierarchy vs FYP from direct team only)", False, DARK),
    ("• New ILP product impact (2.5× min premium → higher productivities → higher compensation rates)", False, DARK),
    ("• Seasonal deferred payments and normal business fluctuations", False, DARK),
], size=9)
slide_footer(s, 13)

# ============ SLIDE 14: MONITORING FRAMEWORK ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "GOVERNANCE", "VAR Monitoring & Control Framework", BLUE)

mon_data = [
    ["Action Item","Owner","Frequency","Timeline"],
    ["Dashboard Channel VaR report development","ASBD","—","2025 Q1"],
    ["Monitor total payout incl. PSM+ for significant fluctuations","ASBD","Monthly","Ongoing"],
    ["Review VaR monitoring tool for enhancement","ASBD","Quarterly","Ongoing"],
    ["Deep-dive payout data: compensation components, agency segments","ASBD","Monthly","From Mar '25"],
    ["Red-flag investigation: multi-indicator screening","Line1B","Monthly","Already in place"],
    ["Update tracking in Local ARC; escalate to ARC/RCSC when breached","Line1B","Monthly","From Mar '25"],
    ["Embed VAR impact review in annual compensation review","ASBD","Annual","2025 Q4"],
]
add_table(s, Inches(0.5), Inches(1.3), Inches(12.3), mon_data)

add_box(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.3), fill=RGBColor(230,247,239), text="")
add_text(s, Inches(0.65), Inches(5.25), Inches(12), Inches(0.25),
         "Escalation Protocol When Threshold Is Breached", size=11, bold=True, color=GREEN)
add_text(s, Inches(0.65), Inches(5.55), Inches(11.8), Inches(0.9),
         "When monthly VAR exceeds $1.5M, the local Compensation team and Line1B (front risk sales) must provide a deep-dive analysis. Line1B already manages monthly protocols for other risk controls — this is integrated into existing workflows, not a new standalone process. Multi-indicator screening includes: persistency ratios, sale patterns, phone/bank sharing, payout anomalies.",
         size=10, color=DARK)
slide_footer(s, 14)

print("Slides 13-14 done")

# ============ SLIDE 15: CONCLUSION ============
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
add_box(s, Inches(0.5), Inches(0.3), Inches(1.5), Inches(0.3), fill=RGBColor(50,50,70),
        text="SUMMARY", font_size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.5),
         "Conclusion & Recommendation", size=22, bold=True, color=WHITE)

panels = [
    ("Current State", "$0.3M",
     "2025 actual VAR — historically low due to high persistency, reduced manpower, and business downtrend. Clean VAR (excl. timing gaps) = $0.24M.",
     GREEN),
    ("Expected State", "$0.85M",
     "2026 forecast under new compensation scheme with 35% FYP growth. Same FYP as 2023 but lower commission rate (26% vs 32%) and cleaner salesforce.",
     ORANGE),
    ("Proposed Trigger", "$1.5M",
     "Reduced from legacy $3.0M — a 50% tightening. Provides ~1.76× buffer above forecast. Absorbs formula asymmetry, new ILP impact, and normal fluctuations.",
     BLUE),
]
for i, (title, val, desc, accent) in enumerate(panels):
    x = Inches(0.5 + i * 4.2)
    add_box(s, x, Inches(1.5), Inches(3.9), Inches(2.5), fill=RGBColor(40,40,60), text="")
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(0.06), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_text(s, x + Inches(0.2), Inches(1.6), Inches(3.5), Inches(0.3),
             title, size=13, color=accent, bold=True)
    add_text(s, x + Inches(0.2), Inches(1.95), Inches(3.5), Inches(0.5),
             val, size=28, color=WHITE, bold=True)
    add_text(s, x + Inches(0.2), Inches(2.55), Inches(3.5), Inches(1.3),
             desc, size=9, color=RGBColor(170,170,187))

add_box(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(1.5), fill=RGBColor(0,40,25), text="")
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.3), Inches(0.06), Inches(1.5))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN
bar.line.fill.background()
bar.shadow.inherit = False
add_text(s, Inches(0.7), Inches(4.4), Inches(11.8), Inches(1.3),
         "The $1.5M trigger point represents a proactive 50% tightening from the legacy $3.0M threshold, demonstrating strengthened risk controls. It is calibrated against the 2026 forecast ($0.85M), stress-tested under the new compensation scheme, and supported by a structurally cleaner salesforce following 2024–2025 conduct reviews. The comprehensive monthly monitoring framework with multi-indicator screening remains in place regardless of threshold level.",
         size=11, color=RGBColor(221,221,238))

add_text(s, Inches(2), Inches(6.2), Inches(9), Inches(0.3),
         "All data as of January 2026  |  3rd Draft", size=9, color=GRAY, align=PP_ALIGN.CENTER)
slide_footer(s, 15)

# ============ SAVE ============
import os
out_dir = os.path.dirname(os.path.abspath(__file__))
out = os.path.join(out_dir, "var-report-3rd-draft.pptx")
prs.save(out)
print(f"\nSaved to {out}")
print("Total slides: 15")
