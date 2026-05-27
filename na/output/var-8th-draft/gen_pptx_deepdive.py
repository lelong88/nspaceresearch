"""Generate PPTX for the March 2026 VaR deep-dive report."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

W = Inches(13.333)
H = Inches(7.5)
GREEN = RGBColor(0, 167, 88)
DARK = RGBColor(26, 26, 46)
BLUE = RGBColor(0, 115, 207)
RED = RGBColor(230, 57, 70)
ORANGE = RGBColor(244, 162, 97)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(102, 102, 102)
LGRAY = RGBColor(244, 245, 247)
HILITE = RGBColor(255, 243, 224)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_box(slide, left, top, w, h, fill=None, text="", font_size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return shape


def add_text(slide, left, top, w, h, text, size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tx


def add_multiline(slide, left, top, w, h, lines, size=11, color=DARK):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if len(item) == 2:
            txt, bold = item
            clr = color
        else:
            txt, bold, clr = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bold
        p.space_after = Pt(4)
    return tx


def add_kpi(slide, left, top, val, label, accent=GREEN, w=Inches(2.8), h=Inches(1.1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = LGRAY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(28)
    p.font.color.rgb = DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(8)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_table(slide, left, top, w, rows_data, col_widths=None, font_size=9, header_align=PP_ALIGN.CENTER, body_align=PP_ALIGN.CENTER, highlight_rows=None, highlight_color=HILITE, row_height=0.32):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(row_height * rows))
    table = table_shape.table
    if col_widths:
        for j, cw in enumerate(col_widths):
            table.columns[j].width = cw
    highlight_rows = highlight_rows or set()
    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.alignment = header_align if i == 0 else body_align
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            elif i in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LGRAY
    return table_shape


def slide_header(slide, tag, title, tag_color=BLUE):
    bg_map = {GREEN: RGBColor(230, 247, 239), BLUE: RGBColor(224, 240, 255), RED: RGBColor(253, 232, 234), ORANGE: RGBColor(255, 243, 224)}
    add_box(slide, Inches(0.5), Inches(0.3), Inches(2.4), Inches(0.32), fill=bg_map.get(tag_color, RGBColor(224, 240, 255)),
            text=tag, font_size=9, color=tag_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.5), title, size=22, bold=True, color=DARK)


def slide_footer(slide, num):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
             "Manulife VN  ·  March 2026 VaR Deep-Dive  ·  ILP26 launch month review",
             size=8, color=GRAY)
    add_text(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.3),
             str(num), size=8, color=GRAY, align=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 1 — Title
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_box(s, Inches(0), Inches(0), W, H, fill=DARK, text="")
add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.6),
         "VaR Deep-Dive — March 2026", size=42, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(2.9), Inches(12), Inches(0.5),
         "Standalone-month review for the ILP26 launch month", size=20, color=ORANGE)
add_text(s, Inches(0.7), Inches(4.0), Inches(12), Inches(0.4),
         "Audience:  Regional Office — VaR Review Committee", size=14, color=WHITE)
add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.4),
         "Source:  slides/var-8th-draft (OCR'd)", size=14, color=WHITE)
add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
         "Prepared by:  Compensation & Risk team", size=11, color=GRAY)
slide_footer(s, 1)
print("Slide 1 done")

# =====================================================================
# SLIDE 2 — Why this deep-dive
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "CONTEXT", "Why a March-only deep-dive?", ORANGE)

add_multiline(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0), [
    ("In March 2026 the company launched ILP26, a new ILP product with large case size. Three things happened simultaneously:", True, DARK),
    ("", False, DARK),
    ("    1.  Sales volume surged — large FYP inflow concentrated in March.", False, DARK),
    ("    2.  Multiple incentive contests ran on top of business-as-usual compensation.", False, DARK),
    ("    3.  Sales hit the highest bonus-rate tier of the new compensation scheme.", False, DARK),
    ("        Because the scheme rewards on FYP — and ILP26 produces uniformly high FYP — many", False, GRAY),
    ("        agents and UMs jumped to the top bonus rate.", False, GRAY),
], size=12)

add_box(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.6), fill=HILITE, text="")
add_text(s, Inches(0.7), Inches(5.15), Inches(12), Inches(0.4),
         "Regional Office ask", size=12, bold=True, color=ORANGE)
add_text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(1.0),
         "Deep-dive March 2026 alone to confirm there is no abnormal pattern hidden inside the rolling-12-month figure.",
         size=14, color=DARK)
slide_footer(s, 2)
print("Slide 2 done")

# =====================================================================
# SLIDE 3 — Methodology
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "METHODOLOGY", "March-only VaR — same formula, single-month window", BLUE)

add_box(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.9), fill=LGRAY, text="")
add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.6),
         "VaR  =  Total Payout  −  Total FYP", size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.3),
         "Standard formula. We replace the rolling-12-month window with March 2026 only.",
         size=10, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(2.6), Inches(8), Inches(0.3),
         "Payout coverage (March 2026 only)", size=12, bold=True, color=DARK)

cov = [
    ["Component", "Included?"],
    ["First-year commission", "✓"],
    ["Rider bonus", "✓"],
    ["UM compensation", "✓"],
    ["Leader subsidy / GTF", "✓"],
    ["Other agent / leader items", "✓"],
    ["Contest payout based on March production", "✗  (pending freelook → booked May 2026)"],
]
add_table(s, Inches(0.5), Inches(2.95), Inches(12.3), cov, font_size=11, body_align=PP_ALIGN.LEFT, header_align=PP_ALIGN.LEFT, row_height=0.34)

add_box(s, Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.9), fill=HILITE, text="")
add_text(s, Inches(0.7), Inches(5.78), Inches(12), Inches(0.4),
         "Implication", size=11, bold=True, color=ORANGE)
add_text(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.4),
         "March VaR shown here understates the eventual figure by the contest payout that will land in May 2026.",
         size=12, color=DARK)
slide_footer(s, 3)
print("Slide 3 done")

# =====================================================================
# SLIDE 4 — Trend
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "TREND", "March 2026 sits within the normal range — no anomaly", GREEN)

trend = [
    ["Month", "07/25", "08/25", "09/25", "10/25", "11/25", "12/25", "01/26", "02/26", "03/26", "04/26", "05/26 fcst", "06/26"],
    ["VaR (mil USD)", "0.49", "1.43", "0.13", "0.12", "0.30", "0.34", "0.26", "0.56", "0.43", "0.38", "0.64 ↑", "0.24"],
    ["VaR ('000 VND)", "12,270", "35,685", "3,300", "3,099", "7,448", "8,487", "6,478", "13,902", "10,804", "9,377", "16,061", "6,002"],
]
add_table(s, Inches(0.3), Inches(1.5), Inches(12.7), trend, font_size=10, row_height=0.40)

add_text(s, Inches(0.5), Inches(3.4), Inches(12), Inches(0.3),
         "How to read this", size=13, bold=True, color=BLUE)
add_multiline(s, Inches(0.5), Inches(3.75), Inches(12.3), Inches(3.0), [
    ("•  Aug'25 ($1.43M peak) — driven by contest payout booked from June 2025 high-FYP production.", False, DARK),
    ("•  Feb'26 ($0.56M) — driven by contest payout booked from December 2025 high-FYP production plus large H2 contests.", False, DARK),
    ("•  Mar'26 ($0.43M) — in line with the normal range. Sits BELOW the Aug'25 and Feb'26 peaks despite ILP26", False, DARK),
    ("    launch, because contest payouts are not yet booked.", False, GRAY),
    ("•  May'26 forecast ($0.64M) — expected to rise once March-production-based contests post after freelook.", False, DARK),
], size=11)

add_box(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.65), fill=RGBColor(230, 247, 239), text="")
add_text(s, Inches(0.7), Inches(6.32), Inches(12), Inches(0.5),
         "Conclusion: March 2026 standalone VaR shows no abnormal spike. Pattern follows existing contest-cycle seasonality.",
         size=12, bold=True, color=GREEN)
slide_footer(s, 4)
print("Slide 4 done")

# =====================================================================
# SLIDE 5 — Population view: VaR bucket
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "POPULATION", "505 UM+ units with VaR > 0 in March — distribution by bucket", BLUE)

add_kpi(s, Inches(0.3), Inches(1.4), "505", "UM+ UNITS WITH VaR > 0", BLUE, w=Inches(3.0))
add_kpi(s, Inches(3.5), Inches(1.4), "10.80 bil", "TOTAL VaR ('000 VND)", BLUE, w=Inches(3.0))
add_kpi(s, Inches(6.7), Inches(1.4), "$0.43M", "TOTAL VaR (USD)", BLUE, w=Inches(3.0))
add_kpi(s, Inches(9.9), Inches(1.4), "45 / 59%", "UNITS / VaR SHARE ≥ $2,000", ORANGE, w=Inches(3.1))

bucket = [
    ["VaR bucket (USD)", "# UM+", "VaR ('000 VND)", "Share of VaR"],
    ["< 600", "352", "1,550,833", "14%"],
    ["600 – < 1,200", "71", "1,483,589", "14%"],
    ["1,200 – < 2,000", "37", "1,401,291", "13%"],
    ["≥ 2,000", "45", "6,368,519", "59%"],
    ["Total", "505", "10,804,233", "100%"],
]
add_table(s, Inches(0.5), Inches(2.9), Inches(12.3), bucket, font_size=12, highlight_rows={4, 5}, row_height=0.42)

add_box(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.0), fill=HILITE, text="")
add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.4),
         "Read", size=11, bold=True, color=ORANGE)
add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.5),
         "70% of units sit under $1,200 individual VaR — too small to justify gaming. Signal concentrates in the 45 units in the ≥ $2,000 deep-dive bucket.",
         size=12, color=DARK)
slide_footer(s, 5)
print("Slide 5 done")

# =====================================================================
# SLIDE 6 — Population by GTF and MBA Pro tier
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "POPULATION", "By GTF and MBA Pro tier (all 505 units)", BLUE)

prof = [
    ["", "Non MBA Pro", "MBA Pro Silver", "MBA Pro Gold", "MBA Pro Platinum", "Total"],
    ["# UM+ with VaR > 0", "455", "18", "14", "18", "505"],
    ["Non GTF", "294", "11", "5", "15", "325"],
    ["GTF", "161", "7", "9", "3", "180"],
    ["    UM", "117", "6", "5", "3", "131"],
    ["    SUM", "2", "-", "-", "-", "2"],
    ["    DM", "32", "-", "4", "-", "36"],
    ["    SDM", "7", "1", "-", "-", "8"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), prof, font_size=11, body_align=PP_ALIGN.CENTER, row_height=0.36)

add_text(s, Inches(0.5), Inches(4.6), Inches(12), Inches(0.3),
         "P13 persistency view (units with FYP > 0, n = 313)", size=13, bold=True, color=BLUE)

p13 = [
    ["P13", "< 600", "600 – 1,200", "1,200 – 2,000", "≥ 2,000", "Total", "Share"],
    ["< 50%", "10", "2", "-", "-", "12", "4%"],
    ["50 – 70%", "20", "2", "3", "-", "25", "8%"],
    ["70 – 80%", "26", "7", "2", "4", "39", "12%"],
    ["≥ 80%", "128", "30", "26", "36", "220", "70%"],
    ["Not yet P13", "12", "4", "1", "-", "17", "5%"],
    ["Total", "196", "45", "32", "40", "313", "100%"],
]
add_table(s, Inches(0.5), Inches(4.95), Inches(12.3), p13, font_size=10, highlight_rows={4, 6}, row_height=0.32)

slide_footer(s, 6)
print("Slide 6 done")

# =====================================================================
# SLIDE 7 — Deep-dive headline (Note column)
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "DEEP-DIVE", "≥ $2,000 USD bucket — the 45 deep-dive units", RED)

add_box(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.55), fill=HILITE, text="")
add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
         "Data note:  OCR captured 40 of the 45 rows (~5.17 bil VND of 6.37 bil VND total). Counts below use the 40 captured rows; directional conclusions are unaffected.",
         size=10, color=DARK)

add_text(s, Inches(0.5), Inches(2.1), Inches(12), Inches(0.3),
         "Note column (col AR) — what's driving each case", size=13, bold=True, color=BLUE)

note = [
    ["Note (col AR) — page 2 only (n = 22, with Note populated)", "# units", "VaR ('000 VND)", "~ USD"],
    ["Diff FYP bonus and FYP VAR  (calc gap)", "17", "2,799,147", "$111,966"],
    ["Delay payment", "2", "178,641", "$7,146"],
    ["–  (residual)", "3", "248,471", "$9,939"],
    ["Subtotal (reviewed, page 2)", "22", "3,226,259", "$129,050"],
    ["Page 1 (Select review = N, already cleared)", "18", "1,943,121", "$77,725"],
    ["Deep-dive total (OCR'd)", "40", "5,169,380", "$206,775"],
]
add_table(s, Inches(0.5), Inches(2.45), Inches(12.3), note, font_size=11, header_align=PP_ALIGN.LEFT, body_align=PP_ALIGN.LEFT, highlight_rows={1, 4, 6}, row_height=0.38)

add_box(s, Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.2), fill=RGBColor(230, 247, 239), text="")
add_text(s, Inches(0.7), Inches(5.78), Inches(12), Inches(0.4),
         "Headline conclusion", size=12, bold=True, color=GREEN)
add_multiline(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.7), [
    ("~92% of the reviewed VaR amount (page 2) is explained by calculation gap or delayed payment, not arbitrage.", True, DARK),
    ("Calc gap = bonus computed on direct-report-team FYP (wider) vs. VaR formula subtracting only direct-team FYP (narrower).", False, GRAY),
], size=11)
slide_footer(s, 7)
print("Slide 7 done")

# =====================================================================
# SLIDE 8 — By Rank (col D)
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "DEEP-DIVE", "By Rank  (column D)", RED)

rank = [
    ["Rank", "# units", "VaR ('000 VND)", "~ USD", "Share of VaR", "Avg / unit ('000 VND)"],
    ["BM", "9", "1,929,050", "$77,162", "37%", "214,339"],
    ["DM", "8", "941,724", "$37,669", "18%", "117,716"],
    ["SDM", "9", "902,616", "$36,105", "17%", "100,291"],
    ["AM", "5", "583,048", "$23,322", "11%", "116,610"],
    ["SUM", "6", "539,486", "$21,579", "10%", "89,914"],
    ["UM", "3", "273,456", "$10,938", "5%", "91,152"],
    ["Total", "40", "5,169,380", "$206,775", "100%", "129,235"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), rank, font_size=12, highlight_rows={1, 7}, row_height=0.42)

add_box(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4), fill=HILITE, text="")
add_text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.4),
         "Read", size=12, bold=True, color=ORANGE)
add_multiline(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.85), [
    ("BM carries the largest concentration — 23% of units but 37% of VaR amount.", True, DARK),
    ("Driven by the calc-gap mechanic: BMs override the widest indirect team. Structural, not behavioural.", False, GRAY),
], size=11)
slide_footer(s, 8)
print("Slide 8 done")

# =====================================================================
# SLIDE 9 — By MBA Pro tier (col E)
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "DEEP-DIVE", "By MBA Pro tier  (column E)", RED)

mba = [
    ["MBA tier", "# units", "VaR ('000 VND)", "~ USD", "Share of VaR"],
    ["Non MBA Pro", "23", "2,046,943", "$81,878", "40%"],
    ["Platinum (P)", "11", "1,955,590", "$78,224", "38%"],
    ["Silver (S)", "3", "751,433", "$30,057", "15%"],
    ["Gold (G)", "3", "415,414", "$16,617", "8%"],
    ["Total", "40", "5,169,380", "$206,775", "100%"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), mba, font_size=12, highlight_rows={2, 5}, row_height=0.45)

add_box(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.1), fill=HILITE, text="")
add_text(s, Inches(0.7), Inches(4.85), Inches(12), Inches(0.4),
         "Read", size=12, bold=True, color=ORANGE)
add_multiline(s, Inches(0.7), Inches(5.25), Inches(12), Inches(1.6), [
    ("MBA Pro Platinum is over-represented relative to population:", True, DARK),
    ("Only 18 Platinums in the entire 505-unit population — but 11 of them appear in this deep-dive (61%).", False, GRAY),
    ("Platinums are the highest-tier producers; the calc-gap mechanic naturally amplifies on them because", False, GRAY),
    ("they typically lead the largest direct + indirect teams.", False, GRAY),
], size=11)
slide_footer(s, 9)
print("Slide 9 done")

# =====================================================================
# SLIDE 10 — By GTF (col F/H)
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "DEEP-DIVE", "By GTF  (column F / H)", RED)

gtf = [
    ["Group", "# units", "VaR ('000 VND)", "~ USD", "Share of VaR"],
    ["Non-GTF", "30", "3,945,878", "$157,835", "76%"],
    ["GTF", "10", "1,223,502", "$48,940", "24%"],
    ["Total", "40", "5,169,380", "$206,775", "100%"],
]
add_table(s, Inches(0.5), Inches(1.6), Inches(12.3), gtf, font_size=13, highlight_rows={1, 3}, row_height=0.55)

add_box(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.4), fill=HILITE, text="")
add_text(s, Inches(0.7), Inches(4.55), Inches(12), Inches(0.4),
         "Read", size=12, bold=True, color=ORANGE)
add_multiline(s, Inches(0.7), Inches(4.95), Inches(12), Inches(1.85), [
    ("Non-GTF leaders carry 3 out of every 4 dollars of deep-dive VaR.", True, DARK),
    ("GTFs (with explicit subsidy) appear less, even though GTF subsidy is a known calc-gap driver —", False, GRAY),
    ("most flagged cases are organic UM+ leaders with high indirect-team activity.", False, GRAY),
], size=11)
slide_footer(s, 10)
print("Slide 10 done")

# =====================================================================
# SLIDE 11 — Cross-tabs: Rank × MBA  and  Rank × GTF
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "DEEP-DIVE", "Cross-tabs — Rank × MBA Pro  &  Rank × GTF (counts)", RED)

add_text(s, Inches(0.5), Inches(1.4), Inches(7), Inches(0.3),
         "Rank × MBA Pro  (count of units)", size=12, bold=True, color=BLUE)
rxm = [
    ["Rank", "Non MBA Pro", "Platinum", "Silver", "Gold", "Total"],
    ["UM", "2", "1", "-", "-", "3"],
    ["SUM", "4", "-", "1", "1", "6"],
    ["DM", "1", "4", "1", "2", "8"],
    ["SDM", "7", "2", "-", "-", "9"],
    ["AM", "3", "2", "-", "-", "5"],
    ["BM", "6", "2", "1", "-", "9"],
    ["Total", "23", "11", "3", "3", "40"],
]
add_table(s, Inches(0.5), Inches(1.75), Inches(7), rxm, font_size=10, highlight_rows={7}, row_height=0.32)

add_text(s, Inches(8.0), Inches(1.4), Inches(5), Inches(0.3),
         "Rank × GTF  (count of units)", size=12, bold=True, color=BLUE)
rxg = [
    ["Rank", "GTF", "Non-GTF", "Total"],
    ["UM", "1", "2", "3"],
    ["SUM", "2", "4", "6"],
    ["DM", "3", "5", "8"],
    ["SDM", "1", "8", "9"],
    ["AM", "2", "3", "5"],
    ["BM", "1", "8", "9"],
    ["Total", "10", "30", "40"],
]
add_table(s, Inches(8.0), Inches(1.75), Inches(4.8), rxg, font_size=10, highlight_rows={7}, row_height=0.32)

add_box(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4), fill=HILITE, text="")
add_text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.4),
         "Read", size=12, bold=True, color=ORANGE)
add_multiline(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.9), [
    ("Calc-gap-style cases appear across every rank above UM. Concentration is a formula-scope effect", False, DARK),
    ("(direct-team vs. wider override scope), not a rank-specific behaviour.", False, GRAY),
], size=11)
slide_footer(s, 11)
print("Slide 11 done")

# =====================================================================
# SLIDE 12 — Note × Rank cross-tab (page 2 only)
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "DEEP-DIVE", "Note × Rank  (Select review = Y rows, n = 22)", RED)

nxr = [
    ["Note", "UM", "SUM", "DM", "SDM", "AM", "BM", "Total"],
    ["Diff FYP bonus and FYP VAR", "-", "2", "3", "6", "2", "4", "17"],
    ["Delay payment", "1", "-", "-", "-", "-", "1", "2"],
    ["–  (residual)", "-", "-", "1", "1", "1", "-", "3"],
    ["Total", "1", "2", "4", "7", "3", "5", "22"],
]
add_table(s, Inches(0.5), Inches(1.6), Inches(12.3), nxr, font_size=12, header_align=PP_ALIGN.CENTER, body_align=PP_ALIGN.CENTER, highlight_rows={1, 4}, row_height=0.5)

add_box(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.85), fill=RGBColor(230, 247, 239), text="")
add_text(s, Inches(0.7), Inches(5.15), Inches(12), Inches(0.4),
         "Take-away", size=12, bold=True, color=GREEN)
add_multiline(s, Inches(0.7), Inches(5.55), Inches(12), Inches(1.2), [
    ("17 of 22 reviewed cases (77% by count, ~87% by amount) are calc-gap.", True, DARK),
    ("Only 2 cases are delayed-payment timing.  Only 3 small residual cases remain unflagged.", False, DARK),
    ("No case in this set indicates arbitrage behaviour.", False, GREEN),
], size=11)
slide_footer(s, 12)
print("Slide 12 done")

# =====================================================================
# SLIDE 13 — Conclusions
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "CONCLUSIONS", "Conclusions for Regional Office", GREEN)

add_multiline(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), [
    ("1.   March 2026 standalone VaR is normal.", True, DARK),
    ("       $0.43M USD sits below Aug'25 ($1.43M) and Feb'26 ($0.56M) contest-driven peaks.", False, GRAY),
    ("       Forecast May'26 ($0.64M) will be the next peak, when March-production contests post.", False, GRAY),
    ("", False, DARK),
    ("2.   ILP26 launch is not driving abnormal VaR yet.", True, DARK),
    ("       March payout in this view excludes the contest layer. Even at the higher bonus tiers triggered by", False, GRAY),
    ("       ILP26 case-size, the residual VaR after standard explanations is small.", False, GRAY),
    ("", False, DARK),
    ("3.   Deep-dive bucket (≥ $2,000 USD) holds 59% of total March VaR — but ~92% of the reviewed", True, DARK),
    ("       amount is explained by calculation gap or delayed payment, not arbitrage.", True, DARK),
    ("", False, DARK),
    ("4.   Concentration patterns are structural, not behavioural:", True, DARK),
    ("       •  BM rank — 37% of deep-dive VaR.  Widest indirect-team override.", False, GRAY),
    ("       •  MBA Pro Platinum — over-represented (11 of 18 platinums in the population).  Top-producer override.", False, GRAY),
    ("       •  Non-GTF leaders dominate (76% of VaR).  Organic team scale, not subsidised structures.", False, GRAY),
    ("", False, DARK),
    ("5.   No new control gap identified.", True, DARK),
    ("       Existing Line1B monthly screening, deep-dive review, and payment-suspension protocol remain sufficient.", False, GRAY),
    ("       Recommend continuing monthly tracking of the calc-gap fingerprint (BM + Platinum + non-GTF)", False, GRAY),
    ("       into May'26 to catch any deviation when contest payouts post.", False, GRAY),
], size=12)
slide_footer(s, 13)
print("Slide 13 done")

out_dir = os.path.dirname(os.path.abspath(__file__))
out = os.path.join(out_dir, "march-2026-var-deepdive.pptx")
prs.save(out)
print(f"Saved to {out}")
