from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

W = Inches(13.333)
H = Inches(7.5)
GREEN = RGBColor(0, 167, 88)
DARK = RGBColor(26, 26, 46)
BLUE = RGBColor(0, 115, 207)
RED = RGBColor(230, 57, 70)
ORANGE = RGBColor(244, 162, 97)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(102, 102, 102)
LGRAY = RGBColor(244, 245, 247)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, w, h, fill=None, text="", font_size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return shape

def add_text(slide, left, top, w, h, text, size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, w, h, lines, size=11, color=DARK, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bold
        p.space_after = Pt(4)
    return txBox

def add_kpi(slide, left, top, val, label, accent=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(1.1))
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = LGRAY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(32)
    p.font.color.rgb = DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False

def add_table(slide, left, top, w, rows_data, col_widths=None):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(0.35 * rows))
    table = table_shape.table
    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LGRAY
    return table_shape

def slide_header(slide, tag, title, tag_color=GREEN):
    add_box(slide, Inches(0.5), Inches(0.3), Inches(1.8), Inches(0.3), fill=RGBColor(230,247,239) if tag_color==GREEN else RGBColor(224,240,255) if tag_color==BLUE else RGBColor(253,232,234) if tag_color==RED else RGBColor(255,243,224), text=tag, font_size=8, color=tag_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.5), title, size=22, bold=True, color=DARK)

def slide_footer(slide, num):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(2), Inches(0.3), "III Manulife", size=8, color=GRAY)
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3), str(num), size=8, color=GRAY, align=PP_ALIGN.RIGHT)

# =====================================================================
# SLIDE 1: Compensation Change Impact — Back-test & Forecast
# (Updated per 4th draft feedback: revised segment breakdown with
#  Non-MBA / MBA split, updated numbers $0.29M→$0.48M→$0.86M)
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "2026 VAR PROPOSAL", "Compensation Change Impact: Back-test & Forecast", ORANGE)

# KPI banners
add_kpi(s, Inches(0.3), Inches(1.2), "$0.29 mil", "2025 ACTUAL (OLD SCHEME)", GREEN)
add_kpi(s, Inches(3.3), Inches(1.2), "$0.48 mil", "2025 BACK-TEST (NEW SCHEME)", ORANGE)
add_kpi(s, Inches(6.3), Inches(1.2), "$0.86 mil", "2026 FORECAST (NEW SCHEME)", RED)
add_kpi(s, Inches(9.3), Inches(1.2), "+35%", "FYP GROWTH TARGET", BLUE)

# Main data table — segment breakdown
main_data = [
    ["Scenario", "Segment", "Non-MBA", "MBA", "Silver", "Gold", "Platinum", "Total"],
    ["Current 2025", "VaR ('k USD)", "248", "45", "39", "6", "0", "293"],
    ["", "# Units VAR>0", "156", "12", "4", "4", "4", "168"],
    ["Back-Test 2025", "VaR ('k USD)", "378", "103", "83", "18", "2", "481"],
    ["(New comp scheme on", "# Units VAR>0", "190", "12", "6", "3", "3", "202"],
    ["actual perf. 2025)", "", "", "", "", "", "", ""],
    ["Forecast 2026", "VaR ('k USD)", "421", "435", "224", "94", "117", "856"],
    ["(New comp + FYP", "# Units VAR>0", "260", "46", "18", "15", "13", "306"],
    ["growth 2026)", "", "", "", "", "", "", ""],
]
add_table(s, Inches(0.3), Inches(2.5), Inches(12.7), main_data)

# Insight bullets
add_multiline(s, Inches(0.3), Inches(5.8), Inches(12.5), Inches(1.0), [
    ("• Backtest 2025: +64% VaR ($293K → $481K) and rise in positive VaR units (168 → 202) driven by new compensation scheme, with same production as 2025 actual.", False, DARK),
    ("• 2026 FYP growth forecast further pushes VaR +78% ($481K → $856K) and expands positive VaR units (202 → 306), reflecting higher productivity among top performers and substantial inflow of new recruits.", False, DARK),
], size=9)

# Divergent trends callout
add_box(s, Inches(0.3), Inches(6.4), Inches(12.7), Inches(0.5), fill=RGBColor(255,243,224), text="")
add_text(s, Inches(0.45), Inches(6.42), Inches(12.3), Inches(0.45),
    "⚠ Divergent trends: Non-MBA Pro VaR +52% ($248K→$378K) under new scheme  |  MBA Pro VaR surges ×2.3 ($45K→$103K): new scheme boosts top-performer MBA UM+ compensation significantly",
    size=9, bold=False, color=ORANGE)

slide_footer(s, "")
print("Slide 1 (Compensation Impact) done")

# =====================================================================
# SLIDE 2: New ILP Product Impact
# (New slide per 4th draft feedback: ILP product with min 50M VND,
#  contributing ~16% FYP, pushing VaR from $0.86M → $1.27M)
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "2026 VAR PROPOSAL", "New ILP Product Impact", ORANGE)

# KPI banners
add_kpi(s, Inches(0.3), Inches(1.2), "$0.86 mil", "2026 FORECAST (NEW COMP SCHEME)", ORANGE)
add_kpi(s, Inches(5.2), Inches(1.2), "$1.27 mil", "2026 FORECAST (+ NEW ILP PRODUCT)", RED)

# Forecast table — VaR by segment and size bucket
add_text(s, Inches(0.3), Inches(2.5), Inches(6), Inches(0.3), "Forecast 2026 — VaR Distribution", size=11, bold=True, color=BLUE)

forecast_data = [
    ["", "Non MBA Pro", "MBA Pro Silver", "MBA Pro Gold", "MBA Pro Platinum", "Total"],
    ["Total VaR ('k USD)", "616", "348", "153", "155", "1,272"],
    ["<2000 USD", "117", "3", "2", "4", "126"],
    ["<4000 USD", "76", "4", "5", "-", "85"],
    [">=4000 USD", "424", "341", "146", "151", "1,061"],
]
add_table(s, Inches(0.3), Inches(2.85), Inches(12.7), forecast_data)

# Unit count table
add_text(s, Inches(0.3), Inches(4.8), Inches(6), Inches(0.3), "Units with VaR > 0 by Segment", size=11, bold=True, color=BLUE)

unit_data = [
    ["", "Non MBA Pro", "MBA Pro Silver", "MBA Pro Gold", "MBA Pro Platinum", "Total"],
    ["# UM with VaR DT L12M > 0", "219", "14", "13", "10", "256"],
    ["Non GTF", "109", "11", "12", "9", "141"],
    ["GTF", "110", "3", "1", "1", "115"],
    ["  UM", "43", "1", "1", "1", "46"],
    ["  SUM", "-", "-", "-", "-", "-"],
    ["  DM", "16", "2", "-", "-", "18"],
    ["  SDM", "8", "-", "-", "-", "8"],
]
add_table(s, Inches(0.3), Inches(5.15), Inches(12.7), unit_data)

# Impact narrative
add_box(s, Inches(0.3), Inches(2.5), Inches(12.7), Inches(0.01), fill=None, text="")  # spacer
add_multiline(s, Inches(0.3), Inches(6.85), Inches(12.5), Inches(0.6), [
    ("➤ Impact of new ILP product: minimum premium 2.5× higher than current case size, expected to contribute 16% total FYP, increase avg case size from 20M→30M VND, improve productivity by 15%, drive higher payout rates for top MBA producers → VaR increases from $0.86M to $1.27M.", False, DARK),
], size=9)

slide_footer(s, "")
print("Slide 2 (ILP Impact) done")

# =====================================================================
# SLIDE 3: $1.5M Trigger Point Proposal
# (Updated per 4th draft feedback: threshold changed from $2.0M to $1.5M,
#  50% tightening from legacy $3.0M, ~17% buffer over $1.27M forecast)
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "PROPOSAL", "$1.5 mil Trigger Point", RED)

# Three-column summary boxes
col_data = [
    ("Current State", "$0.29 mil", "2025 actual VaR — historically low due to reduced manpower, and business downtrend.", GREEN),
    ("Expected State", "$1.27 mil", "2026 forecast under new compensation scheme with 35% FYP growth, new compensation scheme effective from Jan'26 and new ILP product launched from Mar'26.", ORANGE),
    ("Proposed Trigger", "$1.5 mil", "Reduced from legacy $3.0M — a 50% tightening.", RED),
]
for i, (title, val, desc, accent) in enumerate(col_data):
    x = Inches(0.3 + i * 4.25)
    add_box(s, x, Inches(1.2), Inches(4.0), Inches(2.0), fill=LGRAY, text="")
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.2), Inches(0.06), Inches(2.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit = False
    add_text(s, x + Inches(0.2), Inches(1.3), Inches(3.6), Inches(0.25), title, size=11, bold=True, color=accent)
    add_text(s, x + Inches(0.2), Inches(1.6), Inches(3.6), Inches(0.5), val, size=28, color=DARK, bold=True)
    add_text(s, x + Inches(0.2), Inches(2.15), Inches(3.6), Inches(0.9), desc, size=9, color=GRAY)

# Rationale section
add_text(s, Inches(0.3), Inches(3.5), Inches(12), Inches(0.3), "2026 VaR Threshold Rationale", size=14, bold=True, color=BLUE)

add_multiline(s, Inches(0.3), Inches(3.9), Inches(12.5), Inches(3.0), [
    ("• $1.5 mil trigger point represents a proactive 50% tightening from the legacy $3.0M threshold, demonstrating tightening signal and strengthened risk controls.", False, DARK),
    ("• Calibrated against the 2026 forecast ($1.27 mil) under the new compensation scheme and new ILP product launch, supported by a structurally cleaner salesforce following 2024–2025 conduct reviews.", False, DARK),
    ("• Proposed $1.5 mil (+17% buffer from $1.27 mil forecast) to absorb:", False, DARK),
    ("    ➤ Calculation-gap cases due to timing difference between production and payment period, as well as a structural asymmetry in the VaR formula where Income includes earnings from the entire team hierarchy, but FYP only captures direct-team production (excludes UM+ in direct team).", False, GRAY),
    ("    ➤ A UM+ with a large indirect team or with many high-productive UM+ in direct report can show positive VAR even without any risk — because their income reflects a wider scope than the FYP offset.", False, GRAY),
    ("    ➤ This ratio is consistent with 2025 observation, in which $0.05 mil/$0.29 mil in VaR (~17%) resulted from calculation-gap cases.", False, GRAY),
    ("• The comprehensive monthly monitoring framework with multi-indicator screening remains in place regardless of threshold level (details in monitoring slide).", False, DARK),
], size=10)

slide_footer(s, "")
print("Slide 3 (Trigger Point Proposal) done")

# ============ SAVE ============
import os
out_dir = os.path.dirname(os.path.abspath(__file__))
out = os.path.join(out_dir, "gen_pptx_5th_draft_partial.pptx")
prs.save(out)
print(f"Saved to {out}")
