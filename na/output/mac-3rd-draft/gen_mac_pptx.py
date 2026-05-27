"""Generate the MAC Program PPTX deck from the knowledge base."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Layout
W = Inches(13.333)
H = Inches(7.5)

# Manulife palette
GREEN = RGBColor(0, 167, 88)
DARK = RGBColor(26, 26, 46)
BLUE = RGBColor(0, 115, 207)
RED = RGBColor(230, 57, 70)
ORANGE = RGBColor(244, 162, 97)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(102, 102, 102)
LGRAY = RGBColor(244, 245, 247)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, w, h, fill=None, text="", font_size=12, color=DARK,
            bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return shape


def add_text(slide, left, top, w, h, text, size=12, color=DARK, bold=False,
             align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tx


def add_multiline(slide, left, top, w, h, lines, size=11, color=DARK):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            txt, bold, clr = item
        else:
            txt, bold, clr = item, False, color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bold
        p.space_after = Pt(4)
    return tx


def add_kpi(slide, left, top, val, label, accent=GREEN, w=Inches(2.8), h=Inches(1.1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = LGRAY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(28)
    p.font.color.rgb = DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_table(slide, left, top, w, rows_data, font_size=9, header_color=DARK):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(0.32 * rows))
    table = table_shape.table
    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.alignment = PP_ALIGN.CENTER if (i == 0 or j > 0) else PP_ALIGN.LEFT
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LGRAY
    return table_shape


def slide_header(slide, tag, title, tag_color=GREEN):
    bg_color = (RGBColor(230, 247, 239) if tag_color == GREEN
                else RGBColor(224, 240, 255) if tag_color == BLUE
                else RGBColor(253, 232, 234) if tag_color == RED
                else RGBColor(255, 243, 224))
    add_box(slide, Inches(0.5), Inches(0.3), Inches(2.4), Inches(0.3), fill=bg_color,
            text=tag, font_size=8, color=tag_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.5), title,
             size=22, bold=True, color=DARK)


def slide_footer(slide, num):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(4), Inches(0.3),
             "III Manulife  |  MAC Program Knowledge Base", size=8, color=GRAY)
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3),
             str(num), size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])
