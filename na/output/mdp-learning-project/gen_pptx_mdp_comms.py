"""
Generate a 4-slide PPTX for MDP Learning Project communication plan.
Slide 1: Xanh Tuong Lai interactive explainer demo vs raw slides
Slide 2: Insurance Payout Simulator demo vs raw slides
Slide 3: Communication plan - email draft
Slide 4: Communication plan - game show details
"""

import asyncio
import os
import sys

# Add project root to path for gen_image_bedrock import
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", ".."))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Dimensions & Colors ──
W = Inches(13.333)
H = Inches(7.5)
GREEN = RGBColor(0, 167, 88)
DARK = RGBColor(26, 26, 46)
BLUE = RGBColor(0, 115, 207)
RED = RGBColor(230, 57, 70)
ORANGE = RGBColor(244, 162, 97)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(102, 102, 102)
LGRAY = RGBColor(244, 245, 247)
MANULIFE_GREEN = RGBColor(0, 107, 63)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


# ── Helper functions ──
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, w, h, fill=None, text="", font_size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return shape


def add_text(slide, left, top, w, h, text, size=12, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, w, h, lines, size=11, color=DARK, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bold
        p.space_after = Pt(4)
    return txBox



def slide_header(slide, tag, title, tag_color=GREEN):
    tag_bg = (
        RGBColor(230, 247, 239) if tag_color == GREEN else
        RGBColor(224, 240, 255) if tag_color == BLUE else
        RGBColor(253, 232, 234) if tag_color == RED else
        RGBColor(255, 243, 224)
    )
    add_box(slide, Inches(0.5), Inches(0.3), Inches(2.2), Inches(0.3),
            fill=tag_bg, text=tag, font_size=8, color=tag_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.5),
             title, size=22, bold=True, color=DARK)


def slide_footer(slide, num):
    add_text(slide, Inches(0.5), Inches(7.0), Inches(2), Inches(0.3),
             "MDP Learning Project", size=8, color=GRAY)
    add_text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3),
             str(num), size=8, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image_placeholder(slide, left, top, w, h, label="[Screenshot placeholder]"):
    """Add a dashed-border placeholder box for an image to be pasted later."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = LGRAY
    shape.line.color.rgb = GRAY
    shape.line.dash_style = 4  # dash
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.font.italic = True
    # Vertical center
    tf.paragraphs[0].space_before = Pt(h.inches * 72 / 2 - 10)
    return shape


def add_accent_bar(slide, left, top, h, color=GREEN):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


# ============ SLIDE 1: XANH TUONG LAI EXPLAINER ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "DEMO #1 — PRODUCT EXPLAINER", "Xanh Tương Lai: Interactive Demo vs. Raw Slides", GREEN)

# Left: Raw slides (the problem)
add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3),
         "❌  Traditional Approach: 36 Slides of Text", size=13, bold=True, color=RED)
add_accent_bar(s, Inches(0.5), Inches(1.6), Inches(3.0), RED)
add_box(s, Inches(0.6), Inches(1.6), Inches(5.9), Inches(3.0), fill=RGBColor(253, 232, 234))
add_image_placeholder(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(1.8),
                      "[Paste raw slide screenshot here]")
add_multiline(s, Inches(0.8), Inches(3.6), Inches(5.5), Inches(1.0), [
    ("• 36 dense pages of legal/product text", False, DARK),
    ("• No personalization — same content for everyone", False, DARK),
    ("• Staff skim or skip entirely", False, DARK),
    ("• Zero engagement, low retention", False, DARK),
], size=9)

# Right: Interactive demo (the solution)
add_text(s, Inches(6.8), Inches(1.2), Inches(6), Inches(0.3),
         "✅  Interactive Explainer Demo", size=13, bold=True, color=GREEN)
add_accent_bar(s, Inches(6.8), Inches(1.6), Inches(3.0), GREEN)
add_box(s, Inches(6.9), Inches(1.6), Inches(5.9), Inches(3.0), fill=RGBColor(230, 247, 239))
add_image_placeholder(s, Inches(7.1), Inches(1.7), Inches(5.5), Inches(1.8),
                      "[Paste interactive demo screenshot here]")
add_multiline(s, Inches(7.1), Inches(3.6), Inches(5.5), Inches(1.0), [
    ("• Personalized story based on user's age, budget, family", False, DARK),
    ("• Interactive scenarios: death, illness, maturity, market crash", False, DARK),
    ("• Real-time calculations with actual product formulas", False, DARK),
    ("• FAQ, glossary, fund info — all in one page", False, DARK),
], size=9)

# Bottom insight box
add_box(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.8), fill=RGBColor(230, 247, 239))
add_accent_bar(s, Inches(0.5), Inches(4.9), Inches(1.8), GREEN)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.3),
         "Why This Matters", size=14, bold=True, color=MANULIFE_GREEN)
add_multiline(s, Inches(0.7), Inches(5.35), Inches(11.8), Inches(1.2), [
    ("Staff don't need to memorize 36 slides. They experience the product through their own profile — age, budget, family situation — and see exactly how Xanh Tương Lai works for a real customer.", False, DARK),
    ("", False, None),
    ("The interactive demo turns passive reading into active exploration. Staff can play with scenarios (\"what if the customer dies in year 3?\", \"what happens in a market crash?\") and build genuine product intuition.", False, DARK),
], size=10)
slide_footer(s, 1)

print("Slide 1 done")


# ============ SLIDE 2: INSURANCE PAYOUT SIMULATOR ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "DEMO #2 — BENEFIT SIMULATOR", "Insurance Payout Simulator: Interactive vs. Raw Slides", GREEN)

# Left: Raw slides (the problem)
add_text(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.3),
         "❌  Traditional: 35-Page Illustration Document", size=13, bold=True, color=RED)
add_accent_bar(s, Inches(0.5), Inches(1.6), Inches(3.0), RED)
add_box(s, Inches(0.6), Inches(1.6), Inches(5.9), Inches(3.0), fill=RGBColor(253, 232, 234))
add_image_placeholder(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(1.8),
                      "[Paste raw illustration slide here]")
add_multiline(s, Inches(0.8), Inches(3.6), Inches(5.5), Inches(1.0), [
    ("• 35 pages of tables, numbers, legal disclaimers", False, DARK),
    ("• Fixed to one customer profile — can't explore \"what if\"", False, DARK),
    ("• Staff struggle to explain payout mechanics", False, DARK),
    ("• Customers confused by wall of numbers", False, DARK),
], size=9)

# Right: Interactive simulator (the solution)
add_text(s, Inches(6.8), Inches(1.2), Inches(6), Inches(0.3),
         "✅  Interactive Payout Simulator", size=13, bold=True, color=GREEN)
add_accent_bar(s, Inches(6.8), Inches(1.6), Inches(3.0), GREEN)
add_box(s, Inches(6.9), Inches(1.6), Inches(5.9), Inches(3.0), fill=RGBColor(230, 247, 239))
add_image_placeholder(s, Inches(7.1), Inches(1.7), Inches(5.5), Inches(1.8),
                      "[Paste simulator demo screenshot here]")
add_multiline(s, Inches(7.1), Inches(3.6), Inches(5.5), Inches(1.0), [
    ("• Adjust age, premium, fund — see results instantly", False, DARK),
    ("• Visual story: \"what happens if you die in year 3?\"", False, DARK),
    ("• Year-by-year account value projection chart", False, DARK),
    ("• Staff can demo live to customers in meetings", False, DARK),
], size=9)

# Bottom insight box
add_box(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.8), fill=RGBColor(230, 247, 239))
add_accent_bar(s, Inches(0.5), Inches(4.9), Inches(1.8), GREEN)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.3),
         "Why This Matters", size=14, bold=True, color=MANULIFE_GREEN)
add_multiline(s, Inches(0.7), Inches(5.35), Inches(11.8), Inches(1.2), [
    ("A 35-page illustration document is a compliance artifact, not a learning tool. Staff can't internalize payout mechanics by reading tables — they need to play with the numbers.", False, DARK),
    ("", False, None),
    ("The simulator lets staff input any customer profile and instantly see how benefits work across scenarios. It builds confidence for real sales conversations and replaces rote memorization with hands-on understanding.", False, DARK),
], size=10)
slide_footer(s, 2)

print("Slide 2 done")


# ============ SLIDE 3: COMMUNICATION PLAN — EMAIL ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "COMMUNICATION PLAN", "Step 1: Internal Email Launch", BLUE)

# Email mockup area
add_box(s, Inches(0.5), Inches(1.2), Inches(7.5), Inches(5.5), fill=WHITE)
# Email border
border = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(7.5), Inches(5.5))
border.shadow.inherit = False
border.fill.background()
border.line.color.rgb = RGBColor(200, 200, 200)
border.line.width = Pt(1)

# Email header bar
add_box(s, Inches(0.5), Inches(1.2), Inches(7.5), Inches(0.5), fill=MANULIFE_GREEN,
        text="", font_size=10)
add_text(s, Inches(0.7), Inches(1.25), Inches(7), Inches(0.4),
         "📧  From: MDP Learning Team  |  To: All Internal Staff", size=9, color=WHITE, bold=True)

# Subject line
add_text(s, Inches(0.7), Inches(1.8), Inches(7), Inches(0.3),
         "Subject: 🎮 Học sản phẩm kiểu mới — Chơi mà hiểu, hiểu mà nhớ!", size=11, bold=True, color=DARK)

# Banner placeholder
add_image_placeholder(s, Inches(0.7), Inches(2.2), Inches(7.1), Inches(1.2),
                      "[Paste generated banner image here — see gen_images.py]")

# Email body
add_multiline(s, Inches(0.7), Inches(3.5), Inches(7.1), Inches(3.0), [
    ("Xin chào các anh chị! 👋", True, MANULIFE_GREEN),
    ("", False, None),
    ("Bạn có 5 phút? Đó là tất cả những gì bạn cần để hiểu sản phẩm Xanh Tương Lai — không cần đọc 36 trang slide!", False, DARK),
    ("", False, None),
    ("Chúng tôi vừa ra mắt 2 công cụ interactive hoàn toàn mới:", False, DARK),
    ("🔹 Product Explainer — nhập thông tin, xem câu chuyện bảo hiểm cá nhân hóa", False, BLUE),
    ("🔹 Payout Simulator — thử các kịch bản chi trả, hiểu ngay cơ chế sản phẩm", False, BLUE),
    ("", False, None),
    ("Và đặc biệt... 🎯", True, DARK),
    ("", False, None),
    ("Tháng này chúng tôi sẽ tổ chức GAME SHOW trực tuyến — kiểm tra kiến thức sản phẩm theo phong cách Confetti! 15 câu hỏi, trả lời nhanh, có thưởng. Chi tiết bên dưới 👇", False, DARK),
    ("", False, None),
    ("📅 Fuel-up Friday | 1:30 – 2:00 PM | Livestream qua FB \"Manulife Tốt hơn mỗi ngày\"", True, MANULIFE_GREEN),
], size=9)

# Right side: key details
add_text(s, Inches(8.3), Inches(1.2), Inches(4.8), Inches(0.3),
         "Email Highlights", size=14, bold=True, color=BLUE)

for i, (icon, title, desc) in enumerate([
    ("🎯", "Hook", "Fun subject line with emoji — stands out in inbox"),
    ("🖼️", "Visual Banner", "AI-generated banner (Bedrock Nova) — eye-catching, on-brand"),
    ("📱", "Interactive Links", "Direct links to both demo tools — one click to try"),
    ("🎮", "Game Show Teaser", "Creates anticipation for the upcoming live event"),
    ("⏰", "Clear CTA", "Date, time, platform — no ambiguity"),
]):
    y = Inches(1.6 + i * 1.05)
    add_box(s, Inches(8.3), y, Inches(4.5), Inches(0.9), fill=LGRAY)
    add_text(s, Inches(8.45), y + Inches(0.05), Inches(0.4), Inches(0.3), icon, size=16)
    add_text(s, Inches(8.9), y + Inches(0.05), Inches(3.8), Inches(0.25), title, size=10, bold=True, color=DARK)
    add_text(s, Inches(8.9), y + Inches(0.35), Inches(3.8), Inches(0.5), desc, size=8, color=GRAY)

slide_footer(s, 3)

print("Slide 3 done")


# ============ SLIDE 4: COMMUNICATION PLAN — GAME SHOW ============
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, "COMMUNICATION PLAN", "Step 2: Online Game Show — Confetti Format", ORANGE)

# Main game show info panel
add_box(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.8), fill=RGBColor(255, 243, 224))
add_accent_bar(s, Inches(0.5), Inches(1.2), Inches(2.8), ORANGE)
add_text(s, Inches(0.7), Inches(1.3), Inches(5.4), Inches(0.3),
         "🎮  Game Show: Confetti Format", size=16, bold=True, color=DARK)
add_text(s, Inches(0.7), Inches(1.7), Inches(5.4), Inches(0.25),
         "Fun | Engaging | Interactive | Knowledge Sharing", size=10, color=ORANGE, bold=True)

add_multiline(s, Inches(0.7), Inches(2.1), Inches(5.4), Inches(1.8), [
    ("Format: Livestream → 15 multiple-choice questions → real-time answers", False, DARK),
    ("When: 1 round/month — Fuel-up Friday, 1:30 – 2:00 PM", False, DARK),
    ("Where: FB Livestream \"Manulife Tốt hơn mỗi ngày\"", False, DARK),
    ("Target: Internal staff (Phase 1), then Sales force (Phase 2)", False, DARK),
], size=10)

# Flow diagram
add_text(s, Inches(0.5), Inches(4.2), Inches(5.8), Inches(0.3),
         "Game Flow", size=13, bold=True, color=BLUE)

flow_items = [
    ("📢", "Topics\nAnnounced", RGBColor(224, 240, 255)),
    ("📚", "Materials\nShared", RGBColor(230, 247, 239)),
    ("🎬", "Livestream\nStarts", RGBColor(255, 243, 224)),
    ("❓", "15 MCQ\nQuestions", RGBColor(253, 232, 234)),
    ("🏆", "Winners\nAnnounced", RGBColor(224, 240, 255)),
]
for i, (icon, label, bg) in enumerate(flow_items):
    x = Inches(0.5 + i * 1.25)
    add_box(s, x, Inches(4.6), Inches(1.1), Inches(1.1), fill=bg,
            text="", font_size=9)
    add_text(s, x, Inches(4.65), Inches(1.1), Inches(0.4), icon, size=20, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.05), Inches(1.1), Inches(0.6), label, size=8, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    # Arrow between items
    if i < len(flow_items) - 1:
        add_text(s, x + Inches(1.1), Inches(4.9), Inches(0.15), Inches(0.3), "→", size=14, color=GRAY)

# Right side: Key design principles
add_text(s, Inches(6.6), Inches(1.2), Inches(6.2), Inches(0.3),
         "Design Principles", size=14, bold=True, color=BLUE)

principles = [
    ("🤝", "Collaboration over Competition",
     "Participants can discuss answers with each other — embrace knowledge sharing, not rivalry"),
    ("🎲", "Random from Test Bank",
     "Questions are randomized from a pool — prevents copying answers without actually learning"),
    ("📋", "Pre-shared Materials",
     "Summary/keynote shared before each round — easy to review, lowers barrier to entry"),
    ("🎯", "Phase 1: Product Features",
     "Start with Xanh Tương Lai product knowledge for internal staff"),
    ("🚀", "Phase 2: Cross-functional",
     "Expand to all MVL functions: Operations, Distribution, Legal & Compliance, CX, Campaigns"),
]
for i, (icon, title, desc) in enumerate(principles):
    y = Inches(1.6 + i * 1.05)
    add_box(s, Inches(6.6), y, Inches(6.2), Inches(0.9), fill=LGRAY)
    add_text(s, Inches(6.75), y + Inches(0.05), Inches(0.4), Inches(0.3), icon, size=16)
    add_text(s, Inches(7.2), y + Inches(0.05), Inches(5.5), Inches(0.25), title, size=10, bold=True, color=DARK)
    add_text(s, Inches(7.2), y + Inches(0.35), Inches(5.5), Inches(0.5), desc, size=8, color=GRAY)

# Visual reference note
add_text(s, Inches(6.6), Inches(6.85), Inches(6.2), Inches(0.3),
         "Visual style inspired by existing game show slide deck — board game / Confetti theme", size=7, color=GRAY, align=PP_ALIGN.RIGHT)

slide_footer(s, 4)

print("Slide 4 done")


# ============ SAVE ============
OUT_DIR = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(OUT_DIR, "mdp-learning-comms.pptx")
prs.save(out_path)
print(f"\nSaved to {out_path}")

