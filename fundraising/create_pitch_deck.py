#!/usr/bin/env python3
"""Generate Angel Round Pitch Deck for Step - Language Learning App"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)      # Dark navy
ACCENT = RGBColor(0x6C, 0x63, 0xFF)        # Purple/indigo
ACCENT2 = RGBColor(0x00, 0xD2, 0xFF)       # Cyan
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xC0)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x3F)
GREEN = RGBColor(0x00, 0xC9, 0x7B)
WARM_BG = RGBColor(0xF8, 0xF7, 0xFF)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, radius=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=16, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p

# ============================================================
# SLIDE 1: TITLE / COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, PRIMARY)

# Decorative accent bar at top
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

# Company name
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             "STEP", 28, ACCENT2, True, PP_ALIGN.LEFT)

# Main title
add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
             "Learn a Language From Content You Can't Stop Reading",
             44, WHITE, True, PP_ALIGN.LEFT)

# Subtitle
add_text_box(slide, Inches(1), Inches(3.6), Inches(9), Inches(1.0),
             "AI-powered language learning through engaging content — stories, news, movies, music & more.",
             22, LIGHT_GRAY, False, PP_ALIGN.LEFT)

# Round info box
box = add_shape(slide, Inches(1), Inches(5.2), Inches(4.5), Inches(1.2), RGBColor(0x25, 0x25, 0x45), 0.05)
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Angel Round  ·  $200K  ·  $2M Pre-Money"
p.font.size = Pt(20)
p.font.color.rgb = ACCENT2
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Website
add_text_box(slide, Inches(1), Inches(6.6), Inches(4), Inches(0.5),
             "step.is", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(4), Inches(0.5),
             "THE PROBLEM", 14, ACCENT2, True)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "Language learning is broken for busy people",
             36, WHITE, True)

# Problem cards
problems = [
    ("😴", "Boring Content", "Traditional apps use repetitive drills and generic sentences. Learners lose interest within weeks."),
    ("⏰", "No Time", "Busy adults can't dedicate 30+ min/day to language study. They need learning that fits into their life."),
    ("📉", "Low Retention", "Without engaging context, vocabulary doesn't stick. 80% of Duolingo users quit within 2 weeks."),
    ("🎯", "No Discipline", "Self-study requires enormous willpower. Most people want to learn but can't sustain the habit."),
]

for i, (emoji, title, desc) in enumerate(problems):
    left = Inches(1 + i * 2.9)
    box = add_shape(slide, left, Inches(2.5), Inches(2.6), Inches(3.8), RGBColor(0x25, 0x25, 0x45), 0.05)
    add_text_box(slide, left + Inches(0.3), Inches(2.7), Inches(2.0), Inches(0.7),
                 emoji, 40, WHITE, False, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.3), Inches(3.4), Inches(2.0), Inches(0.5),
                 title, 20, WHITE, True, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.3), Inches(4.0), Inches(2.0), Inches(2.0),
                 desc, 14, LIGHT_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 3: THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(4), Inches(0.5),
             "THE SOLUTION", 14, ACCENT2, True)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "Step: Learn a language while doing something you enjoy",
             36, WHITE, True)

add_text_box(slide, Inches(1), Inches(2.0), Inches(10), Inches(0.8),
             "Step turns entertainment into fluency. Users learn through structured curriculums built around "
             "content they actually want to consume — fiction, news, movies, music, podcasts.",
             18, LIGHT_GRAY, False, PP_ALIGN.LEFT)

# Feature columns
features = [
    ("📖", "Reading", "Ebooks & read-along\nwith inline translations\nand vocabulary highlights"),
    ("🗣️", "Speaking", "Pronunciation drills\n& AI-powered roleplay\nconversations"),
    ("✍️", "Writing", "Sentence & paragraph\nexercises with instant\nAI feedback"),
    ("🃏", "Vocabulary", "Smart flashcards with\nspaced repetition\nalgorithm"),
    ("🎧", "Listening", "Audio exercises from\nreal-world content\nat your level"),
]

for i, (emoji, title, desc) in enumerate(features):
    left = Inches(1 + i * 2.3)
    box = add_shape(slide, left, Inches(3.2), Inches(2.1), Inches(3.2), RGBColor(0x25, 0x25, 0x45), 0.05)
    add_text_box(slide, left + Inches(0.25), Inches(3.4), Inches(1.6), Inches(0.6),
                 emoji, 36, WHITE, False, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.25), Inches(4.0), Inches(1.6), Inches(0.4),
                 title, 18, ACCENT2, True, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.25), Inches(4.5), Inches(1.6), Inches(1.5),
                 desc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: HOW IT WORKS / SECRET SAUCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(6), Inches(0.5),
             "HOW IT WORKS", 14, ACCENT2, True)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "AI-Powered Curriculum Engine: 1 Template → 1,000s of Curriculums",
             34, WHITE, True)

add_text_box(slide, Inches(1), Inches(1.9), Inches(10), Inches(0.8),
             "A human expert designs one high-quality template curriculum for a specific learning need. "
             "Our AI engine then replicates it into thousands of derivative curriculums — each tailored "
             "to different topics, levels, and interests — instantly.",
             17, LIGHT_GRAY, False, PP_ALIGN.LEFT)

# Process flow
steps = [
    ("1", "Discover\nLearning Need", "Research real learner\nneeds & pain points"),
    ("→", "", ""),
    ("2", "Design\nTemplate", "Human expert creates\nhigh-quality template"),
    ("→", "", ""),
    ("3", "Test &\nRefine", "Real-world testing\n& iteration cycles"),
    ("→", "", ""),
    ("4", "AI\nReplicates", "1 template becomes\n1000s of curriculums"),
]

x = Inches(0.5)
for step_num, title, desc in steps:
    if step_num == "→":
        add_text_box(slide, x, Inches(3.6), Inches(0.5), Inches(1.0),
                     "→", 36, ACCENT2, False, PP_ALIGN.CENTER)
        x += Inches(0.5)
    else:
        box = add_shape(slide, x, Inches(3.2), Inches(2.7), Inches(2.5), RGBColor(0x25, 0x25, 0x45), 0.05)
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.0), Inches(3.4), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.paragraphs[0].text = step_num
        ctf.paragraphs[0].font.size = Pt(18)
        ctf.paragraphs[0].font.color.rgb = WHITE
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ctf.word_wrap = False

        add_text_box(slide, x + Inches(0.2), Inches(4.1), Inches(2.3), Inches(0.7),
                     title, 16, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(4.8), Inches(2.3), Inches(0.8),
                     desc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)
        x += Inches(2.9)

# Key insight box
box = add_shape(slide, Inches(1), Inches(6.1), Inches(11.3), Inches(0.9), RGBColor(0x1E, 0x1B, 0x4B), 0.03)
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "💡 Key Insight: Every curriculum is pre-configured for a specific audience — all settings, content, and exercises are hidden so learners focus exclusively on learning."
p.font.size = Pt(15)
p.font.color.rgb = ACCENT2
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 5: MARKET OPPORTUNITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(6), Inches(0.5),
             "MARKET OPPORTUNITY", 14, ACCENT2, True)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "A massive, growing market with room for disruption",
             36, WHITE, True)

# TAM SAM SOM
market_data = [
    ("TAM", "$70B+", "Global Language\nLearning Market\nby 2028", "The total addressable\nmarket for all language\nlearning products"),
    ("SAM", "$15B+", "Mobile Language\nLearning Apps", "Digital-first learners\nwho prefer mobile\napp-based solutions"),
    ("SOM", "$150M+", "Content-Driven\nLanguage Learners", "Busy adults who want\nengaging, entertainment-\nbased learning"),
]

for i, (label, value, segment, desc) in enumerate(market_data):
    left = Inches(1 + i * 3.8)
    box = add_shape(slide, left, Inches(2.3), Inches(3.4), Inches(4.2), RGBColor(0x25, 0x25, 0x45), 0.05)
    add_text_box(slide, left + Inches(0.3), Inches(2.5), Inches(2.8), Inches(0.5),
                 label, 16, ACCENT2, True, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.3), Inches(3.0), Inches(2.8), Inches(0.8),
                 value, 42, GREEN, True, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.3), Inches(3.9), Inches(2.8), Inches(1.0),
                 segment, 17, WHITE, True, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.3), Inches(5.0), Inches(2.8), Inches(1.2),
                 desc, 14, LIGHT_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 6: COMPETITIVE ADVANTAGE / WHY US
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(6), Inches(0.5),
             "WHY STEP WINS", 14, ACCENT2, True)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "Content-first approach with unbeatable unit economics",
             36, WHITE, True)

# Comparison table header
add_shape(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(0.6), ACCENT)
headers = ["", "Duolingo", "Babbel", "Busuu", "Step"]
widths = [2.2, 2.2, 2.2, 2.2, 2.5]
x = Inches(1)
for h, w in zip(headers, widths):
    add_text_box(slide, x + Inches(0.1), Inches(2.35), Inches(w - 0.2), Inches(0.5),
                 h, 15, WHITE, True, PP_ALIGN.CENTER)
    x += Inches(w)

# Comparison rows
rows = [
    ("Content Type", "Generic drills", "Scripted dialogues", "Community content", "Real entertainment"),
    ("Engagement", "Gamification-dependent", "Moderate", "Social features", "Content you love"),
    ("Personalization", "Adaptive difficulty", "Course-based", "Level-based", "AI-curated per niche"),
    ("Scalability", "High (generic)", "Low (manual)", "Medium", "Very High (AI engine)"),
    ("Margin", "High", "Medium", "Medium", "Very High (premade)"),
]

for j, (feature, *vals) in enumerate(rows):
    y = Inches(3.0 + j * 0.7)
    bg_color = RGBColor(0x20, 0x20, 0x3A) if j % 2 == 0 else RGBColor(0x25, 0x25, 0x45)
    add_shape(slide, Inches(1), y, Inches(11.3), Inches(0.65), bg_color)
    x = Inches(1)
    all_vals = [feature] + list(vals)
    for k, (val, w) in enumerate(zip(all_vals, widths)):
        clr = WHITE if k == 0 else (GREEN if k == 4 else LIGHT_GRAY)
        bld = True if k == 0 or k == 4 else False
        add_text_box(slide, x + Inches(0.1), y + Inches(0.05), Inches(w - 0.2), Inches(0.55),
                     val, 13, clr, bld, PP_ALIGN.CENTER)
        x += Inches(w)

# ============================================================
# SLIDE 7: BUSINESS MODEL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(6), Inches(0.5),
             "BUSINESS MODEL", 14, ACCENT2, True)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "High-margin, scalable content platform",
             36, WHITE, True)

# Left side: Revenue model
add_text_box(slide, Inches(1), Inches(2.2), Inches(5), Inches(0.5),
             "Revenue Streams", 20, ACCENT2, True)

rev_items = [
    ("Freemium Subscription", "Free tier with limited access. Premium unlocks full curriculum library, all activity types, and advanced features."),
    ("B2B Partnerships", "White-label curriculum solutions for language schools, corporate training, and educational institutions."),
    ("Long-Tail Niche Strategy", "Cover hundreds of specific learning niches. Each niche attracts its own organic audience and word-of-mouth."),
]

y = Inches(2.8)
for title, desc in rev_items:
    box = add_shape(slide, Inches(1), y, Inches(5.5), Inches(1.2), RGBColor(0x25, 0x25, 0x45), 0.03)
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(5.0), Inches(0.4),
                 title, 16, WHITE, True)
    add_text_box(slide, Inches(1.3), y + Inches(0.5), Inches(5.0), Inches(0.7),
                 desc, 12, LIGHT_GRAY, False)
    y += Inches(1.35)

# Right side: Unit economics
add_text_box(slide, Inches(7.2), Inches(2.2), Inches(5), Inches(0.5),
             "Why Margins Are Exceptional", 20, ACCENT2, True)

box = add_shape(slide, Inches(7.2), Inches(2.8), Inches(5.1), Inches(3.8), RGBColor(0x25, 0x25, 0x45), 0.03)

metrics = [
    ("All content is premade", "No per-user generation cost"),
    ("Shared content pool", "All users consume the same content"),
    ("AI replication", "1 template → 1000s of curriculums"),
    ("Near-zero marginal cost", "Adding users costs almost nothing"),
]

y_m = Inches(3.0)
for title, desc in metrics:
    add_text_box(slide, Inches(7.5), y_m, Inches(4.5), Inches(0.35),
                 "✓  " + title, 15, GREEN, True)
    add_text_box(slide, Inches(7.9), y_m + Inches(0.35), Inches(4.1), Inches(0.35),
                 desc, 12, LIGHT_GRAY, False)
    y_m += Inches(0.8)

# ============================================================
# SLIDE 8: TRACTION & VALIDATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(6), Inches(0.5),
             "TRACTION & VALIDATION", 14, ACCENT2, True)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "From NSpace to Step: Lessons learned, momentum building",
             36, WHITE, True)

# NSpace story (left)
add_text_box(slide, Inches(1), Inches(2.2), Inches(5.5), Inches(0.5),
             "NSpace (Previous Product)", 20, ORANGE, True)

box = add_shape(slide, Inches(1), Inches(2.8), Inches(5.5), Inches(2.5), RGBColor(0x25, 0x25, 0x45), 0.03)
nspace_points = [
    "• Creator-led vocabulary learning platform",
    "• Peaked at ~3,000 users",
    "• Proved the core tech: AI-generated deep vocab drills",
    "• Users loved the drill quality",
    "• But: low retention → product-market fit issue at concept level",
    "• Key learning: users need pre-built paths, not DIY tools",
]
y = Inches(2.9)
for point in nspace_points:
    add_text_box(slide, Inches(1.3), y, Inches(5.0), Inches(0.35),
                 point, 13, LIGHT_GRAY, False)
    y += Inches(0.35)

# Step story (right)
add_text_box(slide, Inches(7.2), Inches(2.2), Inches(5.5), Inches(0.5),
             "Step (New Product) — In Beta", 20, GREEN, True)

box = add_shape(slide, Inches(7.2), Inches(2.8), Inches(5.1), Inches(2.5), RGBColor(0x25, 0x25, 0x45), 0.03)
step_points = [
    "• Consumer-led with premade curriculums",
    "• Every curriculum designed for a specific audience",
    "• All settings pre-configured — learner just learns",
    "• Very positive early beta feedback",
    "• Strong early traction signals",
    "• Extremely high margins (shared content pool)",
]
y = Inches(2.9)
for point in step_points:
    add_text_box(slide, Inches(7.5), y, Inches(4.5), Inches(0.35),
                 point, 13, LIGHT_GRAY, False)
    y += Inches(0.35)

# Traction metrics placeholder
add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.5),
             "Key Metrics (Beta)", 20, ACCENT2, True)

metrics_row = [
    ("[TBD]", "Beta Users"),
    ("[TBD]", "Retention Rate"),
    ("[TBD]", "Curriculums Live"),
    ("[TBD]", "Avg. Session Time"),
    ("[TBD]", "NPS Score"),
]

for i, (val, label) in enumerate(metrics_row):
    left = Inches(1 + i * 2.3)
    box = add_shape(slide, left, Inches(6.1), Inches(2.1), Inches(1.0), RGBColor(0x25, 0x25, 0x45), 0.05)
    add_text_box(slide, left + Inches(0.1), Inches(6.15), Inches(1.9), Inches(0.5),
                 val, 22, GREEN, True, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(6.55), Inches(1.9), Inches(0.4),
                 label, 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: TEAM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(4), Inches(0.5),
             "THE TEAM", 14, ACCENT2, True)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "Technical founder with global education & engineering background",
             36, WHITE, True)

# Founder card
box = add_shape(slide, Inches(1), Inches(2.3), Inches(7), Inches(4.2), RGBColor(0x25, 0x25, 0x45), 0.03)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(6), Inches(0.6),
             "Long Le  —  Founder & CEO", 26, WHITE, True)

credentials = [
    "🎓  MBA — Melbourne Business School, Australia",
    "🎓  BA Computer Science & Economics (Double Major) — Gettysburg College, USA",
    "💼  4 years as Software Developer Team Lead at Epic (epic.com)",
    "     Leading Electronic Health Record company",
    "🚀  Built NSpace from 0 → 3,000 users, now building Step",
    "🔧  Full-stack technical founder: product, engineering, AI, business",
]

y = Inches(3.3)
for cred in credentials:
    add_text_box(slide, Inches(1.5), y, Inches(6.0), Inches(0.4),
                 cred, 15, LIGHT_GRAY, False)
    y += Inches(0.45)

# Hiring card
box = add_shape(slide, Inches(8.5), Inches(2.3), Inches(3.8), Inches(4.2), RGBColor(0x1E, 0x1B, 0x4B), 0.03)

add_text_box(slide, Inches(8.8), Inches(2.5), Inches(3.2), Inches(0.5),
             "🔜 Hiring Next", 20, ACCENT2, True)

add_text_box(slide, Inches(8.8), Inches(3.1), Inches(3.2), Inches(0.5),
             "Learning Content Manager", 17, WHITE, True)

hire_details = [
    "IELTS 8.0+ or equivalent",
    "5+ years language teaching",
    "User research experience",
    "TESOL preferred",
    "Creative & independent",
]

y = Inches(3.7)
for detail in hire_details:
    add_text_box(slide, Inches(8.8), y, Inches(3.2), Inches(0.35),
                 "•  " + detail, 13, LIGHT_GRAY, False)
    y += Inches(0.35)

add_text_box(slide, Inches(8.8), Inches(5.6), Inches(3.2), Inches(0.8),
             "This role is the bottleneck — drives the core curriculum engine that powers the entire product.",
             12, ORANGE, False)

# ============================================================
# SLIDE 10: THE ASK — USE OF FUNDS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(4), Inches(0.5),
             "THE ASK", 14, ACCENT2, True)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
             "Raising $200K at $2M pre-money valuation",
             36, WHITE, True)

# Deal terms
add_text_box(slide, Inches(1), Inches(2.2), Inches(5.5), Inches(0.5),
             "Deal Terms", 20, ACCENT2, True)

box = add_shape(slide, Inches(1), Inches(2.8), Inches(5.5), Inches(2.0), RGBColor(0x25, 0x25, 0x45), 0.03)

terms = [
    ("Round:", "Angel"),
    ("Raising:", "$200,000"),
    ("Pre-Money Valuation:", "$2,000,000"),
    ("Previous Round:", "$20K pre-seed (2023) at $1M cap (SAFE)"),
    ("Runway:", "24 months"),
]

y = Inches(2.9)
for label, value in terms:
    add_text_box(slide, Inches(1.3), y, Inches(2.2), Inches(0.35),
                 label, 14, LIGHT_GRAY, False, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(3.5), y, Inches(2.8), Inches(0.35),
                 value, 14, WHITE, True, PP_ALIGN.LEFT)
    y += Inches(0.35)

# Use of funds
add_text_box(slide, Inches(7.2), Inches(2.2), Inches(5.5), Inches(0.5),
             "Use of Funds — $8,000/month for 24 months", 18, ACCENT2, True)

fund_items = [
    ("Learning Content Manager", "$3,000/mo", "37.5%", ACCENT),
    ("Marketing & Growth", "$4,000/mo", "50.0%", GREEN),
    ("Operations & Legal", "$1,000/mo", "12.5%", ORANGE),
]

y = Inches(2.9)
for title, amount, pct, color in fund_items:
    # Color bar
    add_shape(slide, Inches(7.2), y, Inches(0.15), Inches(0.9), color)
    box = add_shape(slide, Inches(7.4), y, Inches(4.9), Inches(0.9), RGBColor(0x25, 0x25, 0x45), 0.02)
    add_text_box(slide, Inches(7.6), y + Inches(0.05), Inches(3.0), Inches(0.4),
                 title, 15, WHITE, True)
    add_text_box(slide, Inches(7.6), y + Inches(0.45), Inches(2.0), Inches(0.35),
                 amount, 13, LIGHT_GRAY, False)
    add_text_box(slide, Inches(10.8), y + Inches(0.15), Inches(1.3), Inches(0.5),
                 pct, 20, color, True, PP_ALIGN.RIGHT)
    y += Inches(1.05)

# Milestones
add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             "24-Month Milestones", 20, ACCENT2, True)

milestones = [
    ("📚", "100+", "Niche\nCurriculums"),
    ("👥", "[TBD]", "Active\nUsers"),
    ("🌍", "5+", "Target\nLanguages"),
    ("💰", "[TBD]", "Monthly\nRevenue"),
    ("🤝", "[TBD]", "B2B\nPartners"),
]

for i, (emoji, val, label) in enumerate(milestones):
    left = Inches(1 + i * 2.3)
    box = add_shape(slide, left, Inches(5.8), Inches(2.1), Inches(1.3), RGBColor(0x25, 0x25, 0x45), 0.05)
    add_text_box(slide, left + Inches(0.1), Inches(5.85), Inches(0.4), Inches(0.4),
                 emoji, 18, WHITE, False, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.5), Inches(5.85), Inches(1.4), Inches(0.4),
                 val, 20, GREEN, True, PP_ALIGN.LEFT)
    add_text_box(slide, left + Inches(0.5), Inches(6.25), Inches(1.4), Inches(0.7),
                 label, 12, LIGHT_GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 11: VISION & CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT)

# Vision statement
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Make language learning feel like entertainment,\nnot education.",
             40, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.3), Inches(9), Inches(1.0),
             "Step is building the Netflix of language learning — where every piece of content "
             "is a lesson, every story teaches you something new, and fluency happens naturally.",
             20, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# CTA box
box = add_shape(slide, Inches(3.5), Inches(4.8), Inches(6.3), Inches(1.5), ACCENT, 0.05)
tf = box.text_frame
tf.word_wrap = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = "Let's talk"
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"

p2 = tf.add_paragraph()
p2.text = "step.is  ·  Long Le  ·  Founder & CEO"
p2.font.size = Pt(16)
p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xFF)
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(8)

# Bottom tagline
add_text_box(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.5),
             "NSPACE TECHNOLOGY PTE. LTD.", 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
prs.save("Step_Angel_Round_Pitch_Deck.pptx")
print("✅ Pitch deck saved: Step_Angel_Round_Pitch_Deck.pptx")
