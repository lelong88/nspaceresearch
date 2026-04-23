"""Build single-slide title decks for Step, parameterized by locale.

Run directly to produce both ``title-slide-vi.pptx`` and
``title-slide-en.pptx`` next to this file. This is a one-off standalone
artifact (local file output by explicit user request) rather than a full
deck that follows the R2 pipeline rules.
"""
from __future__ import annotations

import io
import sys
from dataclasses import dataclass
from pathlib import Path

import cairosvg
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Paths & brand
# ---------------------------------------------------------------------------
_WORKSPACE_ROOT = Path(__file__).resolve().parents[1]
_LOGO_SVG = _WORKSPACE_ROOT / "step-logo.svg"
_OUT_DIR = Path(__file__).resolve().parent

# Palette (shared with decks/general-audience-vi).
ACCENT = RGBColor(0x0F, 0x5F, 0x5C)  # deep teal
BG = RGBColor(0xFA, 0xF7, 0xF2)      # warm cream
BODY = RGBColor(0x2B, 0x2B, 0x2B)    # near-black
MUTED = RGBColor(0x6B, 0x6B, 0x6B)   # secondary

# 16:9 slide size, in inches.
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ---------------------------------------------------------------------------
# Locale content
# ---------------------------------------------------------------------------
@dataclass(frozen=True)
class TitleSlideContent:
    locale: str          # short code used in the output filename
    tagline: str
    presenter: str
    company: str


VI = TitleSlideContent(
    locale="vi",
    tagline="Học ngoại ngữ qua nội dung bạn yêu thích",
    presenter="Long Le",
    company="Công Ty TNHH NSPACE",
)

EN = TitleSlideContent(
    locale="en",
    tagline="Learn a language through content you enjoy",
    presenter="Long Le",
    company="NSPACE TECHNOLOGY PTE. LTD.",
)


# ---------------------------------------------------------------------------
# Low-level drawing helpers
# ---------------------------------------------------------------------------
def _render_logo_png(height_px: int = 384) -> bytes:
    """Rasterize the Step SVG logo to a PNG byte string."""
    return cairosvg.svg2png(
        bytestring=_LOGO_SVG.read_bytes(),
        output_height=height_px,
    )


def _add_bg(slide, color: RGBColor) -> None:
    """Paint the full slide rectangle with a solid fill."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send the background rectangle behind everything else.
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def _add_accent_bar(slide, *, left_in: float, top_in: float,
                    width_in: float, height_in: float,
                    color: RGBColor = ACCENT) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _add_text(slide, text: str, *, left_in: float, top_in: float,
              width_in: float, height_in: float,
              font_size: int, bold: bool = False,
              color: RGBColor = BODY,
              align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE,
              font_name: str = "Calibri") -> None:
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


# ---------------------------------------------------------------------------
# Slide composition
# ---------------------------------------------------------------------------
def build_title_slide(prs: Presentation, content: TitleSlideContent) -> None:
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    blank_layout = prs.slide_layouts[6]  # 6 = Blank
    slide = prs.slides.add_slide(blank_layout)

    _add_bg(slide, BG)

    # Thin teal accent bar on the left edge for brand consistency.
    _add_accent_bar(
        slide,
        left_in=0.0, top_in=2.5,
        width_in=0.12, height_in=2.5,
    )

    # Logo + "STEP." wordmark, centered as a single cluster near the top.
    logo_png = _render_logo_png(height_px=384)
    logo_h_in = 1.4
    # Square viewBox on the source SVG, so width == height at display size.
    logo_w_in = logo_h_in

    wordmark_text = "STEP."
    wordmark_w_in = 3.2
    gap_in = 0.25
    cluster_w_in = logo_w_in + gap_in + wordmark_w_in
    cluster_left_in = (SLIDE_W_IN - cluster_w_in) / 2
    cluster_top_in = 1.15

    slide.shapes.add_picture(
        io.BytesIO(logo_png),
        Inches(cluster_left_in),
        Inches(cluster_top_in),
        height=Inches(logo_h_in),
        width=Inches(logo_w_in),
    )
    _add_text(
        slide,
        wordmark_text,
        left_in=cluster_left_in + logo_w_in + gap_in,
        top_in=cluster_top_in,
        width_in=wordmark_w_in,
        height_in=logo_h_in,
        font_size=88, bold=True, color=ACCENT,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Tagline — the emotional hook.
    _add_text(
        slide,
        content.tagline,
        left_in=1.0, top_in=3.3,
        width_in=SLIDE_W_IN - 2.0, height_in=1.2,
        font_size=40, bold=True, color=ACCENT,
    )

    # Divider dot row under the tagline.
    dot_y = 4.55
    dot_size = 0.12
    for dx in (-0.35, 0.0, 0.35):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(SLIDE_W_IN / 2 + dx - dot_size / 2),
            Inches(dot_y),
            Inches(dot_size), Inches(dot_size),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()

    # Presenter + company block, footer-style.
    _add_text(
        slide,
        content.presenter,
        left_in=1.0, top_in=5.1,
        width_in=SLIDE_W_IN - 2.0, height_in=0.55,
        font_size=22, bold=True, color=BODY,
    )
    _add_text(
        slide,
        content.company,
        left_in=1.0, top_in=5.65,
        width_in=SLIDE_W_IN - 2.0, height_in=0.5,
        font_size=18, bold=False, color=MUTED,
    )


def build_and_save(content: TitleSlideContent) -> Path:
    prs = Presentation()
    build_title_slide(prs, content)
    out_path = _OUT_DIR / f"title-slide-{content.locale}.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")
    return out_path


def main() -> int:
    for content in (VI, EN):
        build_and_save(content)
    return 0


if __name__ == "__main__":
    sys.exit(main())
