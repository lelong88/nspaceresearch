"""Generate the Vietnamese-language pitch deck for the product Step.

This script produces a 7-slide `.pptx` file at
``decks/general-audience-vi/step-pitch-vi.pptx`` using the ``python-pptx``
library. The deck targets a general (non-technical) audience at a
technology exhibition.

Dependency
----------
- ``python-pptx`` (tested against 1.0.x). Install with::

    pip install python-pptx

How to run
----------
From the workspace root::

    python decks/general-audience-vi/generate_deck.py

The script builds the deck entirely in memory, uploads the resulting
``.pptx`` to the Cloudflare R2 bucket ``step-pitch-decks`` via the
``assets`` module, and prints the public URL on ``slides.nspace.is``.
It does NOT write the ``.pptx`` to local disk. Illustrations are generated
via Vertex Gemini on first run and cached under
``decks/general-audience-vi/.asset_cache/`` for subsequent runs.
All Vietnamese strings are UTF-8 literals in source; no external content
files are read.
"""

import io
import sys
from datetime import datetime, timezone
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# Ensure the deck directory is on sys.path so the ``assets`` sibling module
# can be imported directly whether this file is run as a script or imported
# by the test suite.
_DECK_DIR = Path(__file__).resolve().parent
if str(_DECK_DIR) not in sys.path:
    sys.path.insert(0, str(_DECK_DIR))

# ---------------------------------------------------------------------------
# Palette (RGB hex). See design.md for the rationale behind each color.
# ---------------------------------------------------------------------------
ACCENT = RGBColor(0x0F, 0x5F, 0x5C)  # deep teal
BG = RGBColor(0xFA, 0xF7, 0xF2)  # warm neutral background
BODY = RGBColor(0x2B, 0x2B, 0x2B)  # near-black body text
MUTED = RGBColor(0x6B, 0x6B, 0x6B)  # secondary text
ON_ACCENT = RGBColor(0xFA, 0xF7, 0xF2)  # text on accent fills

# Toggle for emoji glyphs inside icon circles. If False, icon helpers
# substitute short ASCII letters (A, B, C, ...) instead of unicode glyphs.
USE_GLYPH_EMOJI = True

# ---------------------------------------------------------------------------
# Typography
# ---------------------------------------------------------------------------
HEADLINE_FONT = "Inter"
BODY_FONT = "Inter"
FALLBACK_FONT = "Calibri"

# Font sizes (points)
SIZE_TITLE_XL = 54  # slide 1 hook, slide 7 CTA headline
SIZE_TITLE_L = 44  # slides 2, 3, 5, 6
SIZE_TITLE_M = 40  # slide 4 (content-heavy slide; keep titles compact)
SIZE_TAGLINE = 32  # emphasized taglines on slide 3, slide 5
SIZE_BODY_L = 24  # primary body
SIZE_BODY_M = 20  # secondary body / bullets
SIZE_BODY_S = 18  # meta/footnote (e.g., CTA next-step line)

# ---------------------------------------------------------------------------
# Slide geometry (inches). Origin is top-left.
# ---------------------------------------------------------------------------
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
MARGIN_IN = 0.6  # >= 0.5 required by the design; 0.6 gives breathing room

# ---------------------------------------------------------------------------
# Output: R2 object key template
# ---------------------------------------------------------------------------
# No local .pptx is written. Deck bytes are uploaded to R2 bucket
# ``step-pitch-decks`` and served via the ``slides.nspace.is`` CDN domain.
# Objects live at the bucket root (no per-feature prefix) per the deck
# pipeline steering rule, so the shareable URL is just
# ``https://slides.nspace.is/<name>.pptx``. Each run writes two keys: a
# timestamped snapshot (immutable) and a ``<name>.pptx`` alias that
# overwrites on every run.
_R2_KEY_TEMPLATE = "step-pitch-vi-{stamp}.pptx"
_R2_KEY_LATEST = "step-pitch-vi.pptx"

# ---------------------------------------------------------------------------
# Slide content (data model)
#
# ``SLIDES`` is the single source of truth for every piece of Vietnamese text
# that appears in the deck. Each entry mirrors the ``SlideContent`` structure
# from ``design.md`` and is consumed by the per-slide builder functions.
#
# Keys:
#   - ``index``     : 1..7 slide number.
#   - ``title``     : Vietnamese headline.
#   - ``subtitle``  : Optional single supporting line (or ``None``).
#   - ``body``      : List of strings (bullets / paragraphs) OR, for slide 6,
#                     a list of 4 ``{"label", "text"}`` dicts describing each
#                     differentiator cell.
#   - ``emphasis``  : Tagline / pull-quote rendered with visual emphasis
#                     (or ``None``).
#   - ``closing``   : Optional closing line (slide 4 uses this).
#   - ``visual``    : ``{"kind": ..., "params": {...}}`` driving which helper
#                     the builder invokes.
#   - ``notes``     : 2-3 sentence Vietnamese speaker notes.
# ---------------------------------------------------------------------------
SLIDES: list[dict] = [
    {
        "index": 1,
        "title": "Bạn đã học tiếng Anh nhiều năm — nhưng vẫn chưa thực sự 'chạm' vào nó?",
        "subtitle": None,
        "body": [
            "Bạn thuộc ngữ pháp, nhưng vẫn chưa cảm được bài hát tiếng Anh mình yêu.",
        ],
        "emphasis": None,
        "closing": None,
        "visual": {"kind": "lyric_motif", "params": {}},
        "notes": (
            "Step là nơi bạn học ngoại ngữ theo cách mới — không mày mò nhiều năm rồi mới hi vọng dùng được. "
            "Bạn vừa học vừa dùng ngoại ngữ để tiếp cận những kiến thức thú vị, ngay từ ngày đầu."
        ),
    },
    {
        "index": 2,
        "title": "Vì sao cách học cũ không còn đủ?",
        "subtitle": None,
        "body": [
            "Cách học truyền thống tách ngôn ngữ khỏi những nội dung bạn thật sự yêu thích.",
            "Từ vựng trong sách nói về sân bay, thời tiết — hiếm khi chạm tới sở thích thật của bạn.",
            "Khi việc học biến thành nghĩa vụ, động lực tắt dần và bạn bỏ cuộc.",
        ],
        "emphasis": None,
        "closing": None,
        "visual": {"kind": "accent_bar_only", "params": {}},
        "notes": (
            "Cách học cũ tách người học khỏi nội dung họ thật sự yêu. "
            "Khi ngôn ngữ không gắn với điều mình quan tâm, động lực tắt rất nhanh. "
            "Step chọn con đường ngược lại: gắn ngôn ngữ với nội dung người học muốn sống cùng."
        ),
    },
    {
        "index": 3,
        "title": "Step — Một cách học khác",
        "subtitle": None,
        "body": [
            "Phương pháp nào để học tiếng Anh?",
            "Nội dung nào bạn muốn sống cùng trong tiếng Anh?",
            "Step là ứng dụng AI xây hành trình học quanh nội dung bạn đã yêu: bài hát, tiểu thuyết, truyện tranh, podcast, phim.",
        ],
        "emphasis": "Step — Học ngoại ngữ qua những điều bạn thực sự yêu thích.",
        "closing": None,
        "visual": {"kind": "reframe_block", "params": {}},
        "notes": (
            "Ở trình độ cơ bản, bạn học qua những câu chuyện và chủ đề đơn giản. "
            "Lên trình độ cao hơn, bạn tiếp cận những ý tưởng hay và sâu sắc. "
            "Điều quan trọng: trong suốt quá trình, bạn luôn đang sử dụng ngoại ngữ chứ không chỉ học về nó."
        ),
    },
    {
        "index": 4,
        "title": "Ví dụ: học tiếng Anh qua 'Heal the World'",
        "subtitle": "Một người yêu bài hát chọn Heal the World làm điểm bắt đầu.",
        "body": [
            "Đọc lời bài hát với từ vựng ngữ cảnh.",
            "Nghe từng câu, luyện phát âm theo giọng hát.",
            "Khám phá văn hóa: 'heal' sâu hơn 'cure'.",
            "Flashcards sinh ra từ lời bài hát vừa học.",
            "Viết lại thông điệp bài hát bằng lời bạn.",
        ],
        "emphasis": None,
        "closing": "Bạn không chỉ học tiếng Anh — bạn đang sống cùng bài hát mình yêu.",
        "visual": {"kind": "activity_grid", "params": {}},
        "notes": (
            "Tôi tên là Long Lê. Tôi xây Step vì với tôi, ngoại ngữ không chỉ là công cụ — nó là một cánh cửa ẩn giấu những câu chuyện và kiến thức sâu sắc. "
            "Ví dụ từ 'decision': gốc Latin nghĩa đen là 'cắt bỏ' — người ta chỉ thật sự quyết định khi sẵn sàng đánh mất một điều gì đó quan trọng. "
            "Nếu ai đó dịch sẵn cho tôi, tôi sẽ mất câu chuyện đằng sau từ ấy — và câu chuyện mới là cánh cửa."
        ),
    },
    {
        "index": 5,
        "title": "Mỗi từ là một cách nhìn thế giới",
        "subtitle": None,
        "body": [
            "Mỗi từ trong một ngôn ngữ đều mang theo một thế giới quan riêng.",
            "'Heal the world' không chỉ là ba từ; đó là cách nhìn chữa lành như điều cộng đồng, giàu cảm xúc và đầy hy vọng — khác với 'cure' vốn thiên về kỹ thuật y khoa.",
        ],
        "emphasis": "Mỗi ngôn ngữ mới là một phiên bản mới của chính bạn.",
        "closing": None,
        "visual": {"kind": "insight_mark", "params": {}},
        "notes": (
            "Mỗi từ trong một ngôn ngữ mang theo một thế giới quan riêng. "
            "Bên trong 'decision' là lòng can đảm, sự chịu đựng mất mát và cảm xúc — không chỉ là 'sự quyết định' đơn thuần. "
            "Khi hiểu được điều đó, ngoại ngữ trở thành cánh cửa bước sang một thế giới quan khác, và đó là lý do tôi học ngoại ngữ."
        ),
    },
    {
        "index": 6,
        "title": "Điều gì làm Step khác biệt?",
        "subtitle": None,
        "body": [
            {
                "label": "AI thông minh",
                "text": "Điều chỉnh độ khó theo đúng trình độ và tiến bộ của bạn.",
            },
            {
                "label": "Đa dạng nội dung",
                "text": "Bài hát, truyện tranh, tiểu thuyết, podcast — bạn chọn điểm bắt đầu.",
            },
            {
                "label": "Kết nối sâu",
                "text": "Mỗi bài tập gắn trực tiếp với nội dung bạn đang yêu thích.",
            },
            {
                "label": "Không gamification rỗng",
                "text": "Không điểm số ảo, không chuỗi streak gây áp lực — chỉ có tiến bộ thật.",
            },
        ],
        "emphasis": None,
        "closing": None,
        "visual": {"kind": "icon_bullets", "params": {}},
        "notes": (
            "Step không làm cho việc học trở nên dễ dàng — trong thời đại AI, học vẫn luôn là một quá trình đầy khó khăn và cần nhiều nỗ lực. "
            "Học là mang những điều bên ngoài chuyển hoá thành bên trong, và quá trình ấy cần sự tự vận động của chính người học. "
            "Điều Step hứa là luôn nuôi dưỡng sự tò mò và động lực tự thân — nền tảng để tiến bộ lâu dài."
        ),
    },
    {
        "index": 7,
        "title": "Bắt đầu với bài hát, cuốn sách, bộ phim bạn yêu nhất.",
        "subtitle": None,
        "body": [
            "Ngôn ngữ không phải là môn học. Nó là cánh cửa vào một thế giới mới — và một phiên bản mới của chính bạn.",
            "Truy cập https://step.is/vi để tải ứng dụng, đăng ký hoặc tìm hiểu thêm về Step.",
        ],
        "emphasis": "Step",
        "closing": None,
        "visual": {"kind": "qr_block", "params": {}},
        "notes": (
            "Câu nói tôi yêu thích trong giáo dục: 'Giáo dục không phải là đổ đầy một cái bình, mà là thắp lên một ngọn lửa'. "
            "Tôi hi vọng Step trở thành một phần của ngọn lửa đó bên trong bạn. "
            "Tải ứng dụng, hoặc ghé https://step.is/vi để bắt đầu."
        ),
    },
]


# ---------------------------------------------------------------------------
# Helper functions
#
# Each helper mutates the given ``slide`` object in place and returns
# ``None``. Geometry is expressed in inches and converted via ``Inches(...)``.
# ---------------------------------------------------------------------------
def add_bg(slide, color: RGBColor) -> None:
    """Paint the full slide rectangle with the given solid fill color.

    Draws a full-bleed rectangle at (0, 0) covering the entire slide area
    (``SLIDE_WIDTH_IN`` x ``SLIDE_HEIGHT_IN``). The rectangle has no stroke
    so the fill reads as a clean background.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_WIDTH_IN),
        Inches(SLIDE_HEIGHT_IN),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_accent_bar(
    slide,
    *,
    left_in: float,
    top_in: float,
    width_in: float = 0.1,
    height_in: float = 1.2,
    color: RGBColor = ACCENT,
) -> None:
    """Draw a thin decorative accent rectangle at the given inch-geometry.

    Used as a visual anchor next to titles (vertical bar) or above them
    (horizontal bar). The rectangle is filled with ``color`` (default:
    the module-level ``ACCENT``) and has no stroke.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------
_ALIGN_MAP: dict[str, PP_ALIGN] = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

_ANCHOR_MAP: dict[str, MSO_ANCHOR] = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


def _set_bullet(paragraph) -> None:
    """Enable a simple '•' bullet on the given paragraph via its pPr XML.

    ``python-pptx`` does not expose bullet formatting on the Python API, so we
    manipulate the underlying ``a:pPr`` element directly: appending an
    ``<a:buChar char="•"/>`` child marks the paragraph as bulleted with the
    given character. A matching ``<a:buFont typeface="Arial"/>`` keeps the
    glyph stable across renderers that otherwise fall back to a theme bullet.
    """
    pPr = paragraph._p.get_or_add_pPr()

    # Remove any pre-existing bullet children to avoid duplicates when the
    # helper is invoked twice on the same paragraph.
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone", "a:buFont"):
        for existing in pPr.findall(qn(tag)):
            pPr.remove(existing)

    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", "Arial")
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", "•")


def _apply_run_formatting(
    run,
    *,
    font_size: int,
    font_color: RGBColor,
    font_name: str,
    bold: bool,
    italic: bool,
) -> None:
    """Apply font formatting to a single run."""
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = font_color


def add_paragraph(
    text_frame,
    text: str,
    *,
    font_size: int,
    font_color: RGBColor = BODY,
    font_name: str = BODY_FONT,
    bold: bool = False,
    italic: bool = False,
    align: str = "left",
    bullet: bool = False,
    space_before_pt: int = 0,
    space_after_pt: int = 6,
) -> None:
    """Append a paragraph to an existing text frame with explicit formatting.

    The paragraph is added via ``text_frame.add_paragraph()`` and populated
    with a single run carrying ``text``. Font name, size, color, weight, and
    italic styling are applied to that run; alignment, bullet flag, and
    before/after spacing are applied to the paragraph itself.
    """
    paragraph = text_frame.add_paragraph()
    paragraph.alignment = _ALIGN_MAP[align]
    paragraph.space_before = Pt(space_before_pt)
    paragraph.space_after = Pt(space_after_pt)

    run = paragraph.add_run()
    run.text = text
    _apply_run_formatting(
        run,
        font_size=font_size,
        font_color=font_color,
        font_name=font_name,
        bold=bold,
        italic=italic,
    )

    if bullet:
        _set_bullet(paragraph)


def add_text_box(
    slide,
    text: str | list[str],
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    font_size: int,
    font_color: RGBColor = BODY,
    font_name: str = BODY_FONT,
    bold: bool = False,
    italic: bool = False,
    align: str = "left",
    anchor: str = "top",
    line_spacing: float = 1.15,
    bullets: bool = False,
) -> None:
    """Create a text frame at the given geometry and fill it with text.

    If ``text`` is a ``str``, it is rendered as a single paragraph. If it is
    a ``list[str]``, the first entry populates the text frame's default first
    paragraph and the remaining entries are appended via ``add_paragraph``.
    All paragraphs share the same formatting; when ``bullets`` is ``True``
    each paragraph is rendered with a ``•`` bullet.
    """
    textbox = slide.shapes.add_textbox(
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = _ANCHOR_MAP[anchor]

    if isinstance(text, str):
        items = [text]
    else:
        items = list(text)

    if not items:
        return

    # The text frame ships with a single empty paragraph. Reuse it for the
    # first item rather than appending a new one so we don't leave a blank
    # leading line.
    first_paragraph = text_frame.paragraphs[0]
    # Clear any default runs to populate the paragraph cleanly.
    for existing_run in list(first_paragraph.runs):
        existing_run._r.getparent().remove(existing_run._r)

    first_paragraph.alignment = _ALIGN_MAP[align]
    first_paragraph.line_spacing = line_spacing

    first_run = first_paragraph.add_run()
    first_run.text = items[0]
    _apply_run_formatting(
        first_run,
        font_size=font_size,
        font_color=font_color,
        font_name=font_name,
        bold=bold,
        italic=italic,
    )

    if bullets:
        _set_bullet(first_paragraph)

    for entry in items[1:]:
        add_paragraph(
            text_frame,
            entry,
            font_size=font_size,
            font_color=font_color,
            font_name=font_name,
            bold=bold,
            italic=italic,
            align=align,
            bullet=bullets,
        )
        # Match the configured line spacing on appended paragraphs as well.
        text_frame.paragraphs[-1].line_spacing = line_spacing


# ---------------------------------------------------------------------------
# Shape / icon helpers
# ---------------------------------------------------------------------------

# Fallback mapping used when ``USE_GLYPH_EMOJI`` is False and the requested
# glyph is non-ASCII or multi-character. Keys cover every glyph the deck
# intentionally renders (Slide 4 activity row icons and Slide 6 differentiator
# icons) plus a generic bullet fallback.
_GLYPH_ASCII_FALLBACK: dict[str, str] = {
    "♪": "A",
    "🎧": "B",
    "🌍": "C",
    "🃏": "D",
    "✎": "E",
    "AI": "AI",
    "♫": "M",
    "♥": "H",
    "✓": "V",
}


def _resolve_glyph(glyph: str) -> str:
    """Return the glyph to render inside an icon circle.

    If ``USE_GLYPH_EMOJI`` is True, the glyph is used as-is. Otherwise this
    picks a short ASCII-safe substitute: the explicit mapping entry if any,
    else the first ASCII character of the glyph, else ``"•"`` as a last
    resort.
    """
    if USE_GLYPH_EMOJI:
        return glyph
    if glyph in _GLYPH_ASCII_FALLBACK:
        return _GLYPH_ASCII_FALLBACK[glyph]
    for ch in glyph:
        if ch.isascii() and ch.isprintable() and not ch.isspace():
            return ch
    return "•"


def add_icon_circle(
    slide,
    *,
    left_in: float,
    top_in: float,
    diameter_in: float = 0.6,
    fill: RGBColor = ACCENT,
    glyph: str | None = None,
    glyph_color: RGBColor = ON_ACCENT,
    glyph_size: int = 20,
) -> None:
    """Draw a circular accent shape with an optional centered glyph.

    The oval is placed at ``(left_in, top_in)`` with side length
    ``diameter_in`` (inches). It is filled with ``fill`` (defaulting to the
    module-level ``ACCENT``) and renders without a stroke so the circle reads
    as a clean color disc. When ``glyph`` is provided, the oval's own text
    frame is used to center the glyph both horizontally (paragraph alignment)
    and vertically (``MSO_ANCHOR.MIDDLE``). Glyphs honor the module-level
    ``USE_GLYPH_EMOJI`` flag via ``_resolve_glyph``: with emoji disabled, a
    single ASCII letter is substituted so icons still render on hosts without
    emoji-capable fonts.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left_in),
        Inches(top_in),
        Inches(diameter_in),
        Inches(diameter_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()

    if glyph is None:
        return

    text_frame = shape.text_frame
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Tight insets keep the glyph visually centered inside a small oval.
    text_frame.margin_left = Inches(0)
    text_frame.margin_right = Inches(0)
    text_frame.margin_top = Inches(0)
    text_frame.margin_bottom = Inches(0)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    # Clear any default run so the glyph is the only content.
    for existing_run in list(paragraph.runs):
        existing_run._r.getparent().remove(existing_run._r)

    run = paragraph.add_run()
    run.text = _resolve_glyph(glyph)
    _apply_run_formatting(
        run,
        font_size=glyph_size,
        font_color=glyph_color,
        font_name=BODY_FONT,
        bold=True,
        italic=False,
    )


def add_lyric_motif(
    slide,
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
) -> None:
    """Compose the Slide 1 "listener with lyrics" motif from pure shapes.

    The motif fits inside the rectangle defined by ``(left_in, top_in)`` and
    ``(width_in, height_in)`` and uses only palette colors (``ACCENT`` and
    ``MUTED``). It is deliberately stylized rather than realistic: a head
    oval with two small headphone ovals flanking it, a rounded-rectangle
    torso below, and three floating lyric-line rectangles of varying widths
    drifting near the figure. No pictures, no extra accent colors.
    """
    # ------------------------------------------------------------------
    # Figure geometry (all values in inches, relative to the motif box)
    # ------------------------------------------------------------------
    head_diameter = min(width_in * 0.38, height_in * 0.28)
    head_left = left_in + (width_in - head_diameter) / 2.0
    head_top = top_in + height_in * 0.08

    head = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(head_left),
        Inches(head_top),
        Inches(head_diameter),
        Inches(head_diameter),
    )
    head.fill.solid()
    head.fill.fore_color.rgb = ACCENT
    head.line.fill.background()

    # Headphone cups: small ovals on either side of the head, vertically
    # centered on the head.
    phone_diameter = head_diameter * 0.35
    phone_top = head_top + (head_diameter - phone_diameter) / 2.0
    phone_left_x = head_left - phone_diameter * 0.55
    phone_right_x = head_left + head_diameter - phone_diameter * 0.45

    for phone_left in (phone_left_x, phone_right_x):
        phone = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(phone_left),
            Inches(phone_top),
            Inches(phone_diameter),
            Inches(phone_diameter),
        )
        phone.fill.solid()
        phone.fill.fore_color.rgb = ACCENT
        phone.line.fill.background()

    # Torso: a rounded rectangle centered horizontally below the head.
    torso_width = width_in * 0.55
    torso_height = height_in * 0.32
    torso_left = left_in + (width_in - torso_width) / 2.0
    torso_top = head_top + head_diameter + height_in * 0.04

    torso = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(torso_left),
        Inches(torso_top),
        Inches(torso_width),
        Inches(torso_height),
    )
    torso.fill.solid()
    torso.fill.fore_color.rgb = ACCENT
    torso.line.fill.background()

    # ------------------------------------------------------------------
    # Floating lyric-line rectangles
    # ------------------------------------------------------------------
    # Each tuple is (x_offset_frac, y_offset_frac, width_frac, color).
    # x_offset_frac/y_offset_frac are anchored to the top-left of the motif
    # box; width_frac is a fraction of the motif width. Line thickness stays
    # constant so the rectangles read as "lyric lines" rather than blocks.
    line_height = 0.08  # inches, thin horizontal bar
    lyric_lines = [
        (0.05, 0.18, 0.28, ACCENT),
        (0.62, 0.32, 0.32, MUTED),
        (0.08, 0.72, 0.42, ACCENT),
        (0.55, 0.86, 0.30, MUTED),
    ]

    for x_frac, y_frac, w_frac, color in lyric_lines:
        line_left = left_in + width_in * x_frac
        line_top = top_in + height_in * y_frac
        line_width = width_in * w_frac
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(line_left),
            Inches(line_top),
            Inches(line_width),
            Inches(line_height),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = color
        line_shape.line.fill.background()


def add_qr_placeholder(
    slide,
    *,
    left_in: float,
    top_in: float,
    size_in: float = 1.8,
    stroke: RGBColor = ACCENT,
) -> None:
    """Draw a QR-code placeholder: a bordered square with a "QR" label and footnote.

    No real QR encoding is performed. The placeholder is a square rectangle at
    ``(left_in, top_in)`` with side length ``size_in`` (inches), filled with the
    background color and stroked with ``stroke`` at roughly 3 pt so the border
    reads clearly. The word "QR" is centered inside the square (both
    horizontally and vertically) in the stroke color at ~40 pt bold. A small
    muted italic footnote, "Mã QR sẽ được cập nhật", is placed directly below
    the square so reviewers understand the final QR will ship later.
    """
    # ------------------------------------------------------------------
    # Bordered square
    # ------------------------------------------------------------------
    square = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(size_in),
        Inches(size_in),
    )
    square.fill.solid()
    square.fill.fore_color.rgb = BG
    square.line.color.rgb = stroke
    square.line.width = Pt(3)

    # Centered "QR" label inside the square.
    text_frame = square.text_frame
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = Inches(0)
    text_frame.margin_right = Inches(0)
    text_frame.margin_top = Inches(0)
    text_frame.margin_bottom = Inches(0)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    for existing_run in list(paragraph.runs):
        existing_run._r.getparent().remove(existing_run._r)

    run = paragraph.add_run()
    run.text = "QR"
    _apply_run_formatting(
        run,
        font_size=40,
        font_color=stroke,
        font_name=BODY_FONT,
        bold=True,
        italic=False,
    )

    # ------------------------------------------------------------------
    # Muted footnote directly below the square
    # ------------------------------------------------------------------
    footnote_top = top_in + size_in + 0.1
    footnote_height = 0.35
    add_text_box(
        slide,
        "Mã QR sẽ được cập nhật",
        left_in=left_in,
        top_in=footnote_top,
        width_in=size_in,
        height_in=footnote_height,
        font_size=12,
        font_color=MUTED,
        font_name=BODY_FONT,
        italic=True,
        align="center",
        anchor="top",
    )


# ---------------------------------------------------------------------------
# Speaker notes helper
# ---------------------------------------------------------------------------
def add_speaker_notes(slide, notes_text: str) -> None:
    """Attach Vietnamese speaker notes to the slide, replacing any existing text.

    Uses the slide's ``notes_slide.notes_text_frame`` so the notes are
    accessible in PowerPoint's Notes pane and via ``python-pptx`` readers.
    Assigning ``notes_frame.text`` replaces any prior notes rather than
    appending.
    """
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes_text


# ---------------------------------------------------------------------------
# Image helper
# ---------------------------------------------------------------------------
def add_picture(
    slide,
    image_bytes: bytes,
    *,
    left_in: float,
    top_in: float,
    width_in: float | None = None,
    height_in: float | None = None,
) -> None:
    """Embed a raster image (PNG/JPEG bytes) into the slide at the given geometry.

    Wraps ``slide.shapes.add_picture`` so builders can pass raw bytes (as
    returned by the ``assets`` module) without managing a ``BytesIO`` buffer
    themselves. Either ``width_in`` or ``height_in`` may be left ``None``;
    ``python-pptx`` preserves aspect ratio when only one dimension is given.
    """
    buf = io.BytesIO(image_bytes)
    kwargs = {}
    if width_in is not None:
        kwargs["width"] = Inches(width_in)
    if height_in is not None:
        kwargs["height"] = Inches(height_in)
    slide.shapes.add_picture(buf, Inches(left_in), Inches(top_in), **kwargs)


# ---------------------------------------------------------------------------
# Slide builders
#
# One function per slide. Each builder reads its Vietnamese content from the
# ``SLIDES`` module constant (by index) and composes the helpers above into
# the per-slide layout defined in ``design.md`` (§Per-slide layout specs).
# Builders mutate the passed-in ``Presentation`` and return ``None``.
# ---------------------------------------------------------------------------
def build_slide_1_hook(prs) -> None:
    """Build Slide 1 — the emotional hook.

    Layout per ``design.md`` §Per-slide layout specs — Slide 1:
      - full-bleed warm-neutral background
      - thin vertical accent bar on the left
      - large headline text box (54 pt, body color, left-aligned, top-anchored)
      - optional supporting italic line in muted color below the headline
      - the Slide 1 lyric motif on the right-hand region
      - 2–3 sentence Vietnamese speaker notes
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    content = SLIDES[0]  # index 1 → array index 0

    add_bg(slide, BG)
    add_accent_bar(slide, left_in=0.6, top_in=1.2, width_in=0.12, height_in=2.6)

    # Headline (from content["title"]).
    add_text_box(
        slide,
        content["title"],
        left_in=0.9,
        top_in=1.1,
        width_in=8.2,
        height_in=3.0,
        font_size=SIZE_TITLE_XL,
        font_color=BODY,
        bold=True,
        align="left",
        anchor="top",
    )

    # Supporting italic line (single string from content["body"][0]).
    if content["body"]:
        add_text_box(
            slide,
            content["body"][0],
            left_in=0.9,
            top_in=4.3,
            width_in=8.2,
            height_in=0.8,
            font_size=SIZE_BODY_M,
            font_color=MUTED,
            italic=True,
            align="left",
        )

    # Real illustration occupies the right-hand region. Generated via
    # Vertex Gemini on first run, cached thereafter.
    from assets import slide_1_illustration

    add_picture(
        slide,
        slide_1_illustration(),
        left_in=9.2,
        top_in=1.3,
        width_in=3.5,
        height_in=4.9,
    )

    # Speaker notes.
    add_speaker_notes(slide, content["notes"])


def build_slide_2_gap(prs) -> None:
    """Build Slide 2 — why traditional methods fail.

    Layout per ``design.md`` §Per-slide layout specs — Slide 2:
      - full-bleed warm-neutral background
      - thin horizontal accent bar above the title
      - title text box (44 pt bold, accent color, left-aligned)
      - three bulleted body lines, each in its own text box at the
        per-bullet geometry (24 pt, body color, left-aligned, with
        a ``•`` bullet mark)
      - 2–3 sentence Vietnamese speaker notes
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    content = SLIDES[1]  # index 2 → array index 1

    add_bg(slide, BG)
    add_accent_bar(slide, left_in=0.6, top_in=0.7, width_in=1.2, height_in=0.12)

    # Title (from content["title"]).
    add_text_box(
        slide,
        content["title"],
        left_in=0.6,
        top_in=0.95,
        width_in=12.1,
        height_in=1.4,
        font_size=SIZE_TITLE_L,
        font_color=ACCENT,
        bold=True,
        align="left",
        anchor="top",
    )

    # Three body bullets, one per text box, at the geometries from the
    # Slide 2 layout spec.
    bullet_tops = (3.0, 4.2, 5.4)
    for bullet_text, top_in in zip(content["body"], bullet_tops):
        add_text_box(
            slide,
            bullet_text,
            left_in=0.6,
            top_in=top_in,
            width_in=12.1,
            height_in=1.0,
            font_size=SIZE_BODY_L,
            font_color=BODY,
            align="left",
            anchor="top",
            bullets=True,
        )

    # Speaker notes.
    add_speaker_notes(slide, content["notes"])


def build_slide_3_intro(prs) -> None:
    """Build Slide 3 — introducing Step via a reframe and tagline.

    Layout per ``design.md`` §Per-slide layout specs — Slide 3:
      - full-bleed warm-neutral background
      - old reframe question in muted italic, centered (Req 6.1)
      - small accent down-arrow between the old and new questions
      - new reframe question in bold body color, centered (Req 6.1)
      - thin horizontal accent rule centered below the question pair
      - ``Core_Message`` tagline (``emphasis``) in accent color, bold, centered,
        visually distinct from body per Req 2.4 and Req 6.2
      - one-sentence explanation line in body color, centered (Req 6.3)
      - 2–3 sentence Vietnamese speaker notes
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    content = SLIDES[2]  # index 3 → array index 2

    add_bg(slide, BG)

    # Old reframe question — muted italic, centered.
    add_text_box(
        slide,
        content["body"][0],
        left_in=0.6,
        top_in=1.0,
        width_in=12.1,
        height_in=0.9,
        font_size=SIZE_BODY_L,
        font_color=MUTED,
        italic=True,
        align="center",
    )

    # Accent down-arrow between the two reframe questions.
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(6.5),
        Inches(2.0),
        Inches(0.3),
        Inches(0.4),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()

    # New reframe question — bold body color, centered.
    add_text_box(
        slide,
        content["body"][1],
        left_in=0.6,
        top_in=2.5,
        width_in=12.1,
        height_in=0.9,
        font_size=28,
        font_color=BODY,
        bold=True,
        align="center",
    )

    # Thin horizontal accent rule separating the reframe pair from the tagline.
    add_accent_bar(
        slide,
        left_in=3.5,
        top_in=3.7,
        width_in=6.3,
        height_in=0.04,
    )

    # Core_Message tagline (from content["emphasis"]) — accent, bold, centered,
    # larger than any non-title body run on the slide so it reads as emphasis.
    add_text_box(
        slide,
        content["emphasis"],
        left_in=0.6,
        top_in=4.0,
        width_in=12.1,
        height_in=1.4,
        font_size=40,
        font_color=ACCENT,
        bold=True,
        align="center",
    )

    # One-sentence explanation line — body color, centered.
    add_text_box(
        slide,
        content["body"][2],
        left_in=0.6,
        top_in=5.6,
        width_in=12.1,
        height_in=1.3,
        font_size=22,
        font_color=BODY,
        align="center",
    )

    # Speaker notes.
    add_speaker_notes(slide, content["notes"])


def build_slide_4_heal(prs) -> None:
    """Build Slide 4 — the "Heal the World" worked example.

    Layout per ``design.md`` §Per-slide layout specs — Slide 4:
      - full-bleed warm-neutral background
      - title text box (40 pt bold, accent color, left-aligned)
      - italic muted subtitle framing the song (20 pt)
      - five activity rows, each composed of an ``add_icon_circle(...)``
        glyph disc plus a matching body text box at 22 pt; glyphs in
        order ``♪``, ``🎧``, ``🌍``, ``🃏``, ``✎`` (honoring
        ``USE_GLYPH_EMOJI`` via the icon helper)
      - a closing line in accent color, bold italic, 24 pt, left-aligned
      - 2–3 sentence Vietnamese speaker notes

    This is the one slide that intentionally exceeds the 4-bullet default
    (Req 3.5) because Req 7.2 mandates five concrete activity types. The
    trade-off is documented in ``design.md`` (§Design decision: Slide 4
    bullet count exception); row copy is kept short so total body word
    count stays within the 60-word budget.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    content = SLIDES[3]  # index 4 → array index 3

    add_bg(slide, BG)

    # Title (from content["title"]).
    add_text_box(
        slide,
        content["title"],
        left_in=0.6,
        top_in=0.6,
        width_in=12.1,
        height_in=1.0,
        font_size=SIZE_TITLE_M,
        font_color=ACCENT,
        bold=True,
        align="left",
    )

    # Italic muted subtitle framing the song (optional).
    if content.get("subtitle"):
        add_text_box(
            slide,
            content["subtitle"],
            left_in=0.6,
            top_in=1.5,
            width_in=12.1,
            height_in=0.6,
            font_size=SIZE_BODY_M,
            font_color=MUTED,
            italic=True,
            align="left",
        )

    # Five activity rows: icon circle + body text, at the per-row
    # geometries from the Slide 4 layout spec. Glyphs map 1:1 to the
    # activity order documented in design.md. Row widths are shortened
    # so a right-side mood illustration can sit next to them.
    glyphs = ["♪", "🎧", "🌍", "🃏", "✎"]
    row_tops = [2.3, 3.0, 3.7, 4.4, 5.1]
    for glyph, row_text, row_top in zip(glyphs, content["body"], row_tops):
        add_icon_circle(
            slide,
            left_in=0.8,
            top_in=row_top,
            diameter_in=0.6,
            fill=ACCENT,
            glyph=glyph,
            glyph_color=ON_ACCENT,
            glyph_size=20,
        )
        add_text_box(
            slide,
            row_text,
            left_in=1.6,
            top_in=row_top,
            width_in=7.6,
            height_in=0.6,
            font_size=22,
            font_color=BODY,
            align="left",
            anchor="middle",
        )

    # Right-side mood illustration evoking the 'Heal the World' song.
    # Generated via Vertex; cached thereafter.
    from assets import slide_4_illustration

    add_picture(
        slide,
        slide_4_illustration(),
        left_in=9.4,
        top_in=2.3,
        width_in=3.3,
        height_in=3.3,
    )

    # Closing line — accent color, bold italic, left-aligned.
    if content.get("closing"):
        add_text_box(
            slide,
            content["closing"],
            left_in=0.6,
            top_in=6.2,
            width_in=12.1,
            height_in=0.8,
            font_size=SIZE_BODY_L,
            font_color=ACCENT,
            bold=True,
            italic=True,
            align="left",
        )

    # Speaker notes.
    add_speaker_notes(slide, content["notes"])


def build_slide_5_why(prs) -> None:
    """Build Slide 5 — why learning through loved content works.

    Layout per ``design.md`` §Per-slide layout specs — Slide 5:
      - full-bleed warm-neutral background
      - thin horizontal accent bar above the title
      - title text box (44 pt bold, accent color, left-aligned)
      - insight paragraph 1 — the worldview statement (24 pt, body color)
      - insight paragraph 2 — the heal-vs-cure example (22 pt, body color)
      - thin horizontal accent divider rule centered below the insights
      - ``emphasis`` one-liner rendered in accent color, 32 pt bold italic,
        centered, making it visually distinct from body per Req 8.3
      - 2–3 sentence Vietnamese speaker notes
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    content = SLIDES[4]  # index 5 → array index 4

    add_bg(slide, BG)
    add_accent_bar(slide, left_in=0.6, top_in=0.7, width_in=1.2, height_in=0.12)

    # Title (from content["title"]).
    add_text_box(
        slide,
        content["title"],
        left_in=0.6,
        top_in=0.95,
        width_in=12.1,
        height_in=1.2,
        font_size=SIZE_TITLE_L,
        font_color=ACCENT,
        bold=True,
        align="left",
    )

    # Insight paragraph 1 — worldview (24 pt, body color, left-aligned).
    add_text_box(
        slide,
        content["body"][0],
        left_in=0.6,
        top_in=2.5,
        width_in=12.1,
        height_in=1.0,
        font_size=SIZE_BODY_L,
        font_color=BODY,
        align="left",
    )

    # Insight paragraph 2 — heal vs cure (22 pt, body color, left-aligned).
    add_text_box(
        slide,
        content["body"][1],
        left_in=0.6,
        top_in=3.7,
        width_in=12.1,
        height_in=1.5,
        font_size=22,
        font_color=BODY,
        align="left",
    )

    # Thin horizontal accent divider rule between insights and the one-liner.
    add_accent_bar(
        slide,
        left_in=3.5,
        top_in=5.4,
        width_in=6.3,
        height_in=0.04,
    )

    # Emphasis one-liner — accent color, bold italic, centered, at tagline size
    # so it reads as visually distinct from the body paragraphs above.
    add_text_box(
        slide,
        content["emphasis"],
        left_in=0.6,
        top_in=5.7,
        width_in=12.1,
        height_in=1.2,
        font_size=SIZE_TAGLINE,
        font_color=ACCENT,
        bold=True,
        italic=True,
        align="center",
    )

    # Speaker notes.
    add_speaker_notes(slide, content["notes"])


def build_slide_6_diff(prs) -> None:
    """Build Slide 6 — what makes Step different (2×2 differentiator grid).

    Layout per ``design.md`` §Per-slide layout specs — Slide 6:
      - full-bleed warm-neutral background
      - title text box (44 pt bold, accent color, left-aligned)
      - four differentiator cells arranged in a 2×2 grid; each cell is an
        ``add_icon_circle(...)`` glyph disc on the left, a 22 pt bold label
        in body color to its right, and an 18 pt muted body block beneath
        the label, all at the per-cell geometries from the Slide 6 layout
        spec
      - 2–3 sentence Vietnamese speaker notes

    Glyphs, in cell order 0..3 (top-left, top-right, bottom-left,
    bottom-right), are ``"AI"``, ``"♫"``, ``"♥"``, ``"✓"`` per design.md.
    The ``"AI"`` glyph is rendered a few points smaller than the others so
    both characters fit comfortably inside the 0.7 inch icon circle.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    content = SLIDES[5]  # index 6 → array index 5

    add_bg(slide, BG)

    # Title (from content["title"]).
    add_text_box(
        slide,
        content["title"],
        left_in=0.6,
        top_in=0.6,
        width_in=12.1,
        height_in=1.0,
        font_size=SIZE_TITLE_L,
        font_color=ACCENT,
        bold=True,
        align="left",
    )

    # 2×2 differentiator grid. ``cell_layouts`` matches the index order
    # of ``content["body"]`` and of the ``glyphs`` list.
    glyphs = ["AI", "♫", "♥", "✓"]
    cell_layouts = [
        {"icon": (0.8, 2.1), "label": (1.7, 2.0, 4.5, 0.6), "body": (1.7, 2.6, 4.5, 1.6)},
        {"icon": (7.0, 2.1), "label": (7.9, 2.0, 4.6, 0.6), "body": (7.9, 2.6, 4.6, 1.6)},
        {"icon": (0.8, 4.6), "label": (1.7, 4.5, 4.5, 0.6), "body": (1.7, 5.1, 4.5, 1.6)},
        {"icon": (7.0, 4.6), "label": (7.9, 4.5, 4.6, 0.6), "body": (7.9, 5.1, 4.6, 1.6)},
    ]

    for cell, glyph, layout in zip(content["body"], glyphs, cell_layouts):
        icon_left, icon_top = layout["icon"]
        # ``AI`` is two characters; drop the point size slightly so both
        # letters fit within the 0.7 inch disc. Single-character glyphs use
        # the standard 22 pt.
        glyph_size = 18 if glyph == "AI" else 22
        add_icon_circle(
            slide,
            left_in=icon_left,
            top_in=icon_top,
            diameter_in=0.7,
            fill=ACCENT,
            glyph=glyph,
            glyph_color=ON_ACCENT,
            glyph_size=glyph_size,
        )

        # Cell label — 22 pt bold, body color, left-aligned, vertically
        # centered inside its 0.6 inch band so it sits next to the icon.
        label_left, label_top, label_width, label_height = layout["label"]
        add_text_box(
            slide,
            cell["label"],
            left_in=label_left,
            top_in=label_top,
            width_in=label_width,
            height_in=label_height,
            font_size=22,
            font_color=BODY,
            bold=True,
            align="left",
            anchor="middle",
        )

        # Cell body — 18 pt muted, left-aligned, top-anchored beneath the
        # label. Copy is kept to ≤ 16 words per Req 9.3.
        body_left, body_top, body_width, body_height = layout["body"]
        add_text_box(
            slide,
            cell["text"],
            left_in=body_left,
            top_in=body_top,
            width_in=body_width,
            height_in=body_height,
            font_size=SIZE_BODY_S,
            font_color=MUTED,
            align="left",
            anchor="top",
        )

    # Speaker notes.
    add_speaker_notes(slide, content["notes"])


def build_slide_7_cta(prs) -> None:
    """Build Slide 7 — the closing call to action.

    Layout per ``design.md`` §Per-slide layout specs — Slide 7 (adapted):
      - full-bleed warm-neutral background
      - centered invitation headline (46 pt bold, accent color) occupying the
        upper band of the slide (Req 10.1)
      - centered italic sub-text (22 pt, body color) directly below the
        headline (Req 10.2)
      - the real Step logo rasterized to the lower-left, paired with
        the brand mark ``Step`` and a next-step line that includes the URL
        ``https://step.is/vi`` (Req 10.3)
      - the real App Store and Google Play download badges stacked on
        the lower-right (replaces the earlier QR placeholder since the
        badges + URL carry the CTA directly)
      - 2–3 sentence Vietnamese speaker notes

    The headline is phrased as an invitation rather than a demand per
    Req 10.4; copy is sourced verbatim from ``SLIDES[6]``.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    content = SLIDES[6]  # index 7 → array index 6

    add_bg(slide, BG)

    # Invitation headline — 46 pt bold accent, centered horizontally and
    # vertically within its band so the phrase reads as an open invitation.
    add_text_box(
        slide,
        content["title"],
        left_in=0.6,
        top_in=1.0,
        width_in=12.1,
        height_in=2.0,
        font_size=46,
        font_color=ACCENT,
        bold=True,
        align="center",
        anchor="middle",
    )

    # Sub-text — 22 pt italic body color, centered (Req 10.2).
    add_text_box(
        slide,
        content["body"][0],
        left_in=0.6,
        top_in=3.2,
        width_in=12.1,
        height_in=1.6,
        font_size=22,
        font_color=BODY,
        italic=True,
        align="center",
    )

    # Real Step logo rasterized from step-logo.svg, placed just above the
    # brand mark on the lower-left. Height ≈ 1 inch.
    from assets import render_logo_png

    add_picture(
        slide,
        render_logo_png(height_px=256),
        left_in=0.6,
        top_in=5.05,
        height_in=1.0,
    )

    # Brand mark ``Step`` — rendered at ``SIZE_TITLE_XL`` bold in accent
    # color, placed to the right of the logo so the two read as one
    # signature mark.
    add_text_box(
        slide,
        content["emphasis"],
        left_in=1.8,
        top_in=5.2,
        width_in=3.0,
        height_in=1.0,
        font_size=SIZE_TITLE_XL,
        font_color=ACCENT,
        bold=True,
        align="left",
        anchor="middle",
    )

    # Next-step line — 20 pt body color, left-aligned directly below the
    # brand mark. Mentions the public URL ``https://step.is/vi`` so
    # visitors can reach the product without a QR scan.
    add_text_box(
        slide,
        content["body"][1],
        left_in=0.6,
        top_in=6.25,
        width_in=9.0,
        height_in=0.6,
        font_size=SIZE_BODY_M,
        font_color=BODY,
        align="left",
    )

    # App-store badges on the lower-right — real Apple App Store and
    # Google Play badges vendored from https://step.is/assets/. Stacked
    # vertically at ~2.1 inches wide (badges are ~3:1 aspect ratio, so
    # each lands at roughly 2.1 × 0.7 inches). Replaces the earlier QR
    # placeholder and the Vertex closing image since the badges + URL
    # already carry the CTA directly.
    from assets import appstore_badge_png, playstore_badge_png

    add_picture(
        slide,
        appstore_badge_png(),
        left_in=10.6,
        top_in=5.2,
        width_in=2.1,
    )
    add_picture(
        slide,
        playstore_badge_png(),
        left_in=10.6,
        top_in=6.05,
        width_in=2.1,
    )

    # Speaker notes.
    add_speaker_notes(slide, content["notes"])


def build(prs) -> None:
    """Build all 7 slides in order on the given ``Presentation``.

    Invokes each per-slide builder against ``prs`` in narrative order
    (Slide 1 hook → Slide 7 CTA). Each builder appends exactly one slide
    and mutates ``prs`` in place.
    """
    build_slide_1_hook(prs)
    build_slide_2_gap(prs)
    build_slide_3_intro(prs)
    build_slide_4_heal(prs)
    build_slide_5_why(prs)
    build_slide_6_diff(prs)
    build_slide_7_cta(prs)


def main() -> None:
    """Build the deck in memory and upload the .pptx to Cloudflare R2.

    Constructs a new ``Presentation``, configures the 16:9 widescreen slide
    size (13.333 × 7.5 inches), composes all 7 slides via ``build(prs)``,
    serializes the result to bytes, and uploads the bytes to the
    ``step-pitch-decks`` R2 bucket under two keys: a timestamped snapshot
    and a ``.../latest`` alias. No ``.pptx`` is written to local disk.
    The public URL (served via ``slides.nspace.is``) is printed on stdout.
    """
    from assets import upload_pptx

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    build(prs)

    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()

    stamp = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")
    snapshot_key = _R2_KEY_TEMPLATE.format(stamp=stamp)

    snapshot_url = upload_pptx(data, key=snapshot_key)
    latest_url = upload_pptx(data, key=_R2_KEY_LATEST)

    print(f"Uploaded snapshot: {snapshot_url}")
    print(f"Latest:            {latest_url}")


if __name__ == "__main__":
    main()
