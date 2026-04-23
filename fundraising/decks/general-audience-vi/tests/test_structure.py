"""Structural tests for the general-audience-vi-pitch-deck.

The final deck is assembled by ``main()`` as:

  slide 1        — the shared Vietnamese title slide prepended by
                   ``assets.open_title_slide_presentation(language='vi')``
  slides 2..8    — the 7 minimalist content slides built by
                   ``generate_deck.build(prs)`` in order (hook, step,
                   founder, insight, why, promise, closing)

Validates:
  - The deck contains exactly 8 slides (1 title + 7 content).
  - Slide size is 16:9 widescreen (13.333 in × 7.5 in) on the merged
    deck.
  - The shared title slide comes first and carries the Step wordmark.
  - The 7 content slides follow in narrative order; each carries its
    ``hero`` phrase as a visible text run.
"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Inches


_DIM_TOLERANCE_IN = 0.01


def _slide_text(slide) -> list[str]:
    """Collect visible text runs from every text frame on ``slide``."""
    texts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    texts.append(run.text)
    return texts


def _merged_deck(generator):
    """Build the merged deck (title + 7 content slides) in memory."""
    import assets

    prs = assets.open_title_slide_presentation(language=generator.DECK_LANGUAGE)
    generator.build(prs)
    return prs


# -----------------------------------------------------------------------------
# Slide count — title + 7 content
# -----------------------------------------------------------------------------
def test_slide_count(generator):
    """The merged deck contains exactly 8 slides (1 title + 7 content)."""
    prs = _merged_deck(generator)
    assert len(prs.slides) == 8


# -----------------------------------------------------------------------------
# Dimensions — 16:9 widescreen preserved end-to-end
# -----------------------------------------------------------------------------
def test_slide_dimensions(generator):
    """Slide width and height match 13.333 × 7.5 inches within EMU tolerance."""
    prs = _merged_deck(generator)
    assert abs(prs.slide_width.inches - 13.333) < _DIM_TOLERANCE_IN
    assert abs(prs.slide_height.inches - 7.5) < _DIM_TOLERANCE_IN


# -----------------------------------------------------------------------------
# Title slide — first, carries STEP wordmark
# -----------------------------------------------------------------------------
def test_title_slide_is_first(generator):
    """The shared Vietnamese title slide is the first slide in the deck."""
    prs = _merged_deck(generator)
    first_slide_text = _slide_text(prs.slides[0])
    assert any("STEP" in run.upper() for run in first_slide_text), (
        f"Slide 1 should be the title slide carrying the 'STEP' wordmark; "
        f"got runs: {first_slide_text!r}"
    )


# -----------------------------------------------------------------------------
# Narrative order — content slides 2..8 carry their hero phrases
# -----------------------------------------------------------------------------
def test_content_slide_order(generator):
    """Each content slide carries its expected ``hero`` phrase in order."""
    prs = _merged_deck(generator)
    slides_after_title = list(prs.slides)[1:]
    assert len(slides_after_title) == 7

    for i, (slide, content) in enumerate(zip(slides_after_title, generator.SLIDES)):
        runs = _slide_text(slide)
        joined = "\n".join(runs)
        hero = content["hero"]
        assert hero in joined, (
            f"Content slide #{i + 1} (deck position {i + 2}) should carry "
            f"the hero phrase {hero!r}; got runs: {runs!r}"
        )
