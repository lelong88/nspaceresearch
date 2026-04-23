"""Integration tests for the general-audience-vi-pitch-deck generator.

These tests cover the end-to-end path from ``build(prs)`` through
serialization to upload, with R2 network calls stubbed so the suite runs
offline. They verify the documented behavior:

  - Req 13.3: ``main()`` runs cleanly and uploads the deck to R2 under the
    ``step-pitch-decks`` bucket, printing the public ``slides.nspace.is``
    URL on stdout.
  - Req 13.4: A new ``main()`` invocation overwrites the ``.../latest``
    alias in R2 (the snapshot key is timestamped so prior versions are
    preserved under their own keys).
  - Req 13.6: ``build(prs)`` + serialization completes in under 30 seconds
    (measured in-process to avoid an unnecessary subprocess fork).
  - No-local-file rule: ``main()`` MUST NOT write a ``.pptx`` to local
    disk. The integration suite asserts the deck directory still contains
    no ``.pptx`` after ``main()`` returns.

The R2 client is stubbed at the ``assets._r2_client`` seam so no real
network traffic leaves the test machine. Vertex image generation is
stubbed by monkey-patching ``assets.generate_illustration`` (and the
logo rasterizer is content-addressable via the local SVG, so no stub is
needed there).
"""
from __future__ import annotations

import io
import time
from pathlib import Path

import pytest
from pptx import Presentation


_DECK_DIR = Path(__file__).resolve().parent.parent


# -----------------------------------------------------------------------------
# Stubs
# -----------------------------------------------------------------------------
class _FakeR2:
    """In-memory stub for the R2 boto3 client.

    Captures every ``put_object`` call so tests can inspect the bucket,
    key, body bytes, and content type without touching the network.
    """

    def __init__(self) -> None:
        self.calls: list[dict] = []

    def put_object(self, *, Bucket: str, Key: str, Body: bytes, ContentType: str):
        self.calls.append(
            {
                "bucket": Bucket,
                "key": Key,
                "body": Body,
                "content_type": ContentType,
            }
        )
        return {"ETag": '"stub"'}


def _stub_illustration(prompt: str, aspect_ratio: str = "4:3") -> bytes:
    """Return a valid 1×1 PNG so ``add_picture`` accepts the bytes.

    The exact pixel content doesn't matter — the integration tests only
    care about presentation structure and upload behavior, not the
    generated artwork.
    """
    # Smallest valid PNG: a 1×1 fully transparent pixel.
    return bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
        "0000000d49444154789c6300010000000500010d0a2db40000000049454e44ae426082"
    )


@pytest.fixture
def stubbed_assets(monkeypatch):
    """Stub R2 and Vertex in ``assets`` for offline integration tests.

    Yields the ``_FakeR2`` instance so tests can inspect captured
    ``put_object`` calls.
    """
    import assets

    fake = _FakeR2()
    monkeypatch.setattr(assets, "_r2_client", lambda: fake)

    # All illustration and vendored-asset helpers share the same
    # underlying bytes return. Stub every helper the deck builders might
    # pull from so no Vertex calls are made and no vendor_assets/ file
    # needs to exist on CI.
    monkeypatch.setattr(assets, "slide_1_illustration", lambda: _stub_illustration(""))
    monkeypatch.setattr(assets, "slide_4_illustration", lambda: _stub_illustration(""))
    monkeypatch.setattr(assets, "hero_image_png", lambda: _stub_illustration(""))
    monkeypatch.setattr(assets, "appstore_badge_png", lambda: _stub_illustration(""))
    monkeypatch.setattr(assets, "playstore_badge_png", lambda: _stub_illustration(""))

    return fake


# -----------------------------------------------------------------------------
# Tests
# -----------------------------------------------------------------------------
def test_main_uploads_to_r2_and_prints_url(generator, stubbed_assets, capsys):
    """``main()`` builds the deck, uploads to R2, and prints the public URL (Req 13.3)."""
    generator.main()

    # Exactly two uploads: timestamped snapshot + ``.../latest`` alias.
    assert len(stubbed_assets.calls) == 2, (
        f"Expected 2 R2 uploads (snapshot + latest); got "
        f"{len(stubbed_assets.calls)}: {[c['key'] for c in stubbed_assets.calls]}"
    )

    buckets = {call["bucket"] for call in stubbed_assets.calls}
    assert buckets == {"step-pitch-decks"}, (
        f"All uploads must target the ``step-pitch-decks`` bucket; "
        f"got {buckets!r}."
    )

    content_types = {call["content_type"] for call in stubbed_assets.calls}
    assert content_types == {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    }, f"Unexpected content types: {content_types!r}"

    keys = [call["key"] for call in stubbed_assets.calls]
    assert any(key.endswith("step-pitch-vi.pptx") for key in keys), (
        f"Expected a ``.../latest`` alias key; got {keys!r}"
    )
    assert any("step-pitch-vi-" in key and key.endswith(".pptx") for key in keys), (
        f"Expected a timestamped snapshot key; got {keys!r}"
    )

    # Each upload body must be a valid, 7-slide .pptx.
    for call in stubbed_assets.calls:
        prs = Presentation(io.BytesIO(call["body"]))
        assert len(prs.slides) == 7

    # The public URL printed on stdout must point at the CDN domain.
    out = capsys.readouterr().out
    assert "https://slides.nspace.is/" in out, (
        f"Expected public URL on stdout; got:\n{out}"
    )


def test_main_overwrites_latest_alias(generator, stubbed_assets):
    """Repeated ``main()`` runs overwrite the same ``.../latest`` key (Req 13.4)."""
    generator.main()
    generator.main()

    latest_calls = [
        call
        for call in stubbed_assets.calls
        if call["key"] == "step-pitch-vi.pptx"
    ]
    # Two invocations of main ⇒ two uploads to the ``.../latest`` alias.
    assert len(latest_calls) == 2, (
        f"Expected exactly 2 writes to the .../latest alias across 2 "
        f"main() invocations; got {len(latest_calls)}."
    )


def test_main_does_not_write_local_pptx(generator, stubbed_assets):
    """``main()`` MUST NOT write a ``.pptx`` to local disk.

    Snapshots the deck directory before ``main()`` runs and asserts no
    ``.pptx`` file appears afterwards (beyond any that existed before,
    which per the project rules should be zero).
    """
    before = {p for p in _DECK_DIR.rglob("*.pptx")}
    generator.main()
    after = {p for p in _DECK_DIR.rglob("*.pptx")}
    new_pptx = after - before
    assert not new_pptx, (
        f"main() wrote .pptx files to local disk (should upload to R2 "
        f"only): {sorted(new_pptx)!r}"
    )


def test_build_completes_within_30s(generator, stubbed_assets):
    """``build(prs)`` + serialization finishes in under 30 seconds (Req 13.6)."""
    from pptx.util import Inches

    start = time.monotonic()
    prs = Presentation()
    prs.slide_width = Inches(generator.SLIDE_WIDTH_IN)
    prs.slide_height = Inches(generator.SLIDE_HEIGHT_IN)
    generator.build(prs)

    buf = io.BytesIO()
    prs.save(buf)

    duration_s = time.monotonic() - start
    assert duration_s < 30.0, (
        f"Building + serializing the deck took {duration_s:.2f}s, "
        f"exceeding the 30s budget (Req 13.6)."
    )
    assert buf.getvalue(), "Expected non-empty .pptx bytes."
